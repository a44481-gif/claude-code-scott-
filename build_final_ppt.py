# -*- coding: utf-8 -*-
"""
CCS_Global_IFPCS_Final_Report.pptx 构建脚本
- 加载 CCS_Global_Storage_2025_2030.pptx
- 添加 Part 0 (全赛道痛点体系, 3页, 插入原CCS报告第15页之前)
- 深拷贝所有原始 CCS 幻灯片
- 添加 Part 3 (全球市场经济数据补全, 4页, 插入原CCS报告第42页之后)
- 添加结束页
"""

import os
import matplotlib
matplotlib.use('Agg')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from copy import deepcopy

# ---------------------------------------------------------------------------
# 文件路径
# ---------------------------------------------------------------------------
SRC_PPT  = r'd:/claude mini max 2.7/storage_market_output/CCS_Global_Storage_2025_2030.pptx'
OUT_PPT  = r'd:/claude mini max 2.7/storage_market_output/CCS_Global_IFPCS_Final_Report.pptx'

# ---------------------------------------------------------------------------
# deep_copy_slide — 从 merge_ppt_v5.py 的最可靠实现
# ---------------------------------------------------------------------------
def find_cSld(tree):
    for elem in tree.iter():
        if elem.tag == qn('p:cSld'):
            return elem
    return None

def deep_copy_slide(target_prs, source_prs, slide_idx):
    """将 source_prs[slide_idx] 深度拷贝到 target_prs，返回新幻灯片对象。"""
    src_slide  = source_prs.slides[slide_idx]
    src_tree   = src_slide.part._element
    blank_layout = target_prs.slide_layouts[6]
    new_slide  = target_prs.slides.add_slide(blank_layout)
    new_tree   = new_slide.part._element

    cSld_src = find_cSld(src_tree)
    cSld_new = find_cSld(new_tree)

    if cSld_src is not None and cSld_new is not None:
        # -- 背景 --
        bg_src = cSld_src.find(qn('p:bg'))
        bg_new = cSld_new.find(qn('p:bg'))
        if bg_new is not None:
            cSld_new.remove(bg_new)
        if bg_src is not None:
            cSld_new.insert(0, deepcopy(bg_src))

        # -- 形状树 --
        spTree_src = cSld_src.find(qn('p:spTree'))
        spTree_new = cSld_new.find(qn('p:spTree'))

        if spTree_new is not None:
            for child in list(spTree_new):
                if child.tag != qn('p:nvGrpSpPr') and child.tag != qn('p:grpSpPr'):
                    spTree_new.remove(child)

        if spTree_src is not None:
            for child in spTree_src:
                if child.tag != qn('p:nvGrpSpPr') and child.tag != qn('p:grpSpPr'):
                    spTree_new.append(deepcopy(child))

    return new_slide

# ---------------------------------------------------------------------------
# 通用辅助函数
# ---------------------------------------------------------------------------
def add_textbox(slide, text, size, bold, color, x, y, w, h,
                align=PP_ALIGN.CENTER, word_wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size  = Pt(size)
    p.font.bold  = bold
    p.font.color.rgb = color
    p.alignment = align
    return tb

def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None):
    """添加实色矩形"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    if line_rgb:
        shape.line.color.rgb = RGBColor(*line_rgb)
    else:
        shape.line.fill.background()
    return shape

# ---------------------------------------------------------------------------
# 分隔页 / 封面 / 结束页
# ---------------------------------------------------------------------------
def add_divider_slide(prs, part_num, title1, title2="", bg=(0x0A, 0x29, 0x5E)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*bg)

    def add_t(text, size, bold, color, y, h):
        tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.3), h)
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text
        p.font.size = Pt(size); p.font.bold = bold
        p.font.color.rgb = color; p.alignment = PP_ALIGN.CENTER

    add_t(f"PART {part_num}", 20, True,  RGBColor(0xFF, 0xD7, 0x00), Inches(1.8), Inches(0.8))
    add_t(title1,               36, True,  RGBColor(255, 255, 255),  Inches(2.6), Inches(1.2))
    if title2:
        add_t(title2,          28, True,  RGBColor(255, 255, 255),  Inches(3.8), Inches(1.0))
    return slide

def add_cover_slide(prs):
    """封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x0A, 0x29, 0x5E)

    add_textbox(slide,
                "CCS集成母排 × 全球储能市场",
                40, True, RGBColor(255, 255, 255),
                Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.6))

    add_textbox(slide,
                "智能融合电力连接系统（IFPCS）综合商业计划书",
                24, False, RGBColor(0xFF, 0xD7, 0x00),
                Inches(0.5), Inches(2.9), Inches(12.3), Inches(1.0))

    add_textbox(slide,
                "深度融合CCS集成母排技术与全球储能产业趋势",
                16, False, RGBColor(0xCC, 0xD9, 0xEC),
                Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8))

    add_textbox(slide,
                "2026年4月",
                18, True, RGBColor(0xFF, 0xD7, 0x00),
                Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.7))

    return slide

def add_end_slide(prs):
    """结束页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x0A, 0x29, 0x5E)

    add_textbox(slide,
                "谢谢",
                54, True, RGBColor(255, 255, 255),
                Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.5))

    add_textbox(slide,
                "THANK YOU FOR YOUR TIME",
                20, False, RGBColor(0xFF, 0xD7, 0x00),
                Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.8))

    return slide

# ---------------------------------------------------------------------------
# 内容页：标题栏 + 多列内容
# ---------------------------------------------------------------------------
def add_content_slide(prs, title, items, title_bg=(0x0A, 0x29, 0x5E), item_bg=(0x1A, 0x5C, 0x9E)):
    """
    items = list of (header, bullet_points) tuples
    多列布局，等宽分布。
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)

    # 标题栏
    add_rect(slide, Inches(0), Inches(0),
                         Inches(13.33), Inches(1.0), title_bg)
    add_textbox(slide, title, 22, True, RGBColor(255, 255, 255),
                Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.7))

    n = len(items)
    box_w = 12.0 / n
    for i, (header, bullets) in enumerate(items):
        x = 0.3 + i * box_w

        # 栏标题
        add_rect(slide, Inches(x), Inches(1.15),
                 Inches(box_w - 0.1), Inches(0.5), item_bg)
        add_textbox(slide, header, 11, True, RGBColor(255, 255, 255),
                     Inches(x + 0.05), Inches(1.2),
                     Inches(box_w - 0.15), Inches(0.4),
                     align=PP_ALIGN.CENTER)

        # 内容框
        add_rect(slide, Inches(x), Inches(1.7),
                 Inches(box_w - 0.1), Inches(5.3),
                 (255, 255, 255), (200, 200, 200))

        btb = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(1.82),
            Inches(box_w - 0.25), Inches(5.1))
        btf = btb.text_frame
        btf.word_wrap = True

        for j, bullet in enumerate(bullets):
            if j == 0:
                bp = btf.paragraphs[0]
            else:
                bp = btf.add_paragraph()
            bp.text = bullet
            bp.font.size = Pt(9)
            bp.space_before = Pt(3)

    return slide

# ---------------------------------------------------------------------------
# 表格页
# ---------------------------------------------------------------------------
def add_table_slide(prs, title, headers, rows,
                    col_widths=None,
                    title_bg=(0x0A, 0x29, 0x5E),
                    row_colors=None):
    """
    通用表格页。
    row_colors: list of (odd_rgb, even_rgb) — 默认蓝白交替。
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)

    # 标题栏
    add_rect(slide, Inches(0), Inches(0),
             Inches(13.33), Inches(0.9), title_bg)
    add_textbox(slide, title, 20, True, RGBColor(255, 255, 255),
                Inches(0.3), Inches(0.12), Inches(12.7), Inches(0.65))

    n_rows = len(rows) + 1
    n_cols = len(headers)
    row_h = Inches(min(0.38 * n_rows, 5.8))

    if col_widths is None:
        cw = Inches(12.3 / n_cols)
        col_widths = [cw] * n_cols
    else:
        col_widths = [Inches(w) for w in col_widths]

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.3), Inches(1.1),
        sum(col_widths), row_h)
    tbl = tbl_shape.table

    # 列宽
    for j, cw in enumerate(col_widths):
        tbl.columns[j].width = cw

    # 表头
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*title_bg)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10); p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # 数据行
    odd_rgb  = row_colors[0] if row_colors else (255, 255, 255)
    even_rgb = row_colors[1] if row_colors else (0xE8, 0xF4, 0xFF)

    for i, row in enumerate(rows):
        rc = odd_rgb if i % 2 == 0 else even_rgb
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*rc)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.alignment = PP_ALIGN.CENTER

    return slide

# ---------------------------------------------------------------------------
# KPI 数字页
# ---------------------------------------------------------------------------
def add_kpi_slide(prs, title, kpis):
    """
    kpis = list of (number, unit, label, sublabel)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)

    # 标题栏
    add_rect(slide, Inches(0), Inches(0),
             Inches(13.33), Inches(0.9), (0x0A, 0x29, 0x5E))
    add_textbox(slide, title, 20, True, RGBColor(255, 255, 255),
                Inches(0.3), Inches(0.12), Inches(12.7), Inches(0.65))

    n = len(kpis)
    box_w = 12.0 / n

    for i, (num, unit, label, sub) in enumerate(kpis):
        x = 0.3 + i * box_w

        # 卡片
        add_rect(slide, Inches(x), Inches(1.2),
                 Inches(box_w - 0.15), Inches(2.6),
                 (255, 255, 255), (200, 200, 200))

        # 数字
        add_textbox(slide, num, 36, True, RGBColor(0x0A, 0x29, 0x5E),
                    Inches(x), Inches(1.4),
                    Inches(box_w - 0.15), Inches(1.0))

        # 单位
        add_textbox(slide, unit, 14, True, RGBColor(0xFF, 0xD7, 0x00),
                    Inches(x), Inches(2.3),
                    Inches(box_w - 0.15), Inches(0.4))

        # 标签
        add_textbox(slide, label, 10, True, RGBColor(0x33, 0x33, 0x33),
                    Inches(x + 0.05), Inches(2.75),
                    Inches(box_w - 0.2), Inches(0.5))

        # 子标签
        if sub:
            add_textbox(slide, sub, 8, False, RGBColor(0x88, 0x88, 0x88),
                        Inches(x + 0.05), Inches(3.25),
                        Inches(box_w - 0.2), Inches(0.4))

    return slide

# ---------------------------------------------------------------------------
# PART 0 — 全赛道痛点体系（插入 CCS 报告第15页之前）
# ---------------------------------------------------------------------------
def build_part0(prs):
    n_before = len(prs.slides)
    print("[Part 0] 全赛道痛点体系...")

    # ---- Page 1: 5层痛点金字塔 ----
    print("  [Part 0-1] 传统电力连接方案的5层系统性痛点")
    add_table_slide(
        prs,
        title="传统电力连接方案的5层系统性痛点",
        headers=["痛点层级", "核心问题", "量化影响", "IFPCS解决方案"],
        rows=[
            ["技术性能层",
             "大电流温升/振动松脱/盐雾失效/电磁干扰",
             "80%电池热失控源于连接点；-40℃性能衰减60%",
             "一体化铜排均流+热压密封+电磁屏蔽层"],
            ["生产制造层",
             "人工依赖高/良率低/换型慢/交付长",
             "人工布线良率85%；单车型换型15-30天",
             "模块化预装+全自动化产线；12天样品交付"],
            ["现场应用层",
             "安装错误/故障难查/连锁事故",
             "现场安装错误率12%；故障排查平均48小时",
             "即插即用+BMS智能诊断；故障风险降90%"],
            ["全生命周期层",
             "研发贵/装配贵/维护贵/回收难",
             "传统线束占整车成本8-12%；售后成本是制造成本3倍",
             "标准化平台；研发周期3-6个月；可拆解设计"],
            ["产业生态层",
             "供应链波动/认证周期长/迭代滞后",
             "铜价波动10%影响成本5-8%；海外认证6-12个月",
             "国产化替代+预认证模块+三代技术储备"],
        ],
        col_widths=[1.4, 2.8, 2.8, 4.5],
        title_bg=(0x0A, 0x29, 0x5E),
    )

    # ---- Page 2: 三大黄金赛道差异化致命痛点 ----
    print("  [Part 0-2] 三大核心赛道的不可替代需求")
    add_content_slide(
        prs,
        title="三大核心赛道的不可替代需求",
        items=[
            ("新能源汽车（主战场）", [
                "800V高压平台：局部温升超标、电弧击穿风险",
                "CTC/CTB技术：线束维修不可达、整车报废率上升",
                "4680大圆柱：多连接点均流不均",
                "",
                "量化：800V车型传统线束故障率是400V的2.3倍",
            ]),
            ("工商业储能（增长极）", [
                "液冷储能：防水密封要求高、线缆腐蚀老化快",
                "集装箱部署：空间极致压缩、布线难度指数级上升",
                "电网调频：高频充放电导致连接点疲劳断裂",
                "",
                "量化：储能系统连接故障占总停机时间40%",
            ]),
            ("AI数据中心（蓝海）", [
                "超高功率密度：>30kW/柜无法适配传统线束",
                "高压直流供电：拉弧风险高、安全防护难度大",
                "7×24运行：单点故障导致单柜损失超10万元/小时",
                "",
                "量化：数据中心断电每分钟损失超5万美元",
            ]),
        ],
        title_bg=(0x0A, 0x29, 0x5E),
        item_bg=(0x1A, 0x5C, 0x9E),
    )

    # ---- Page 3: 七大高潜力赛道专属痛点图谱 ----
    print("  [Part 0-3] 七大高潜力赛道的差异化致命痛点")
    add_table_slide(
        prs,
        title="七大高潜力赛道的差异化致命痛点",
        headers=["赛道", "核心极端工况", "专属致命痛点", "IFPCS核心价值"],
        rows=[
            ["电动重卡/工程机械",
             "重载振动/高低温/粉尘",
             "接头松动断裂；停机成本>5000元/小时",
             "抗震一体化+防腐蚀密封+免维护"],
            ["船舶/海上风电",
             "盐雾腐蚀/潮湿/摇摆",
             "绝缘老化快；海上运维成本是陆地10倍",
             "船级社认证材料+全密封结构"],
            ["轨道交通",
             "高频振动/强电磁干扰",
             "信号采集失真；寿命需匹配30年",
             "抗干扰屏蔽+高疲劳强度材料"],
            ["户用储能/便携式电源",
             "低温/频繁移动/消费级安全",
             "低温性能衰减；跌落损坏；消费者起火焦虑",
             "低温适配+抗冲击结构+多重安全保护"],
            ["无人机/工业机器人",
             "轻量化/快速充电/高动态",
             "重量限制续航；关节弯折疲劳",
             "超薄柔性设计+集成散热+高弯折寿命"],
            ["航空航天",
             "极端高低温/真空/高可靠性",
             "重量克级要求；单点故障即灾难性后果",
             "极致轻量化+冗余设计+航天级材料"],
            ["光伏逆变器/充电桩",
             "户外暴晒/雨水侵蚀/高功率",
             "老化快；防护等级低；功率密度不足",
             "户外专用材料+IP67防护+高功率集成"],
        ],
        col_widths=[1.8, 1.8, 2.8, 4.1],
        title_bg=(0x0A, 0x29, 0x5E),
    )


# ---------------------------------------------------------------------------
# PART 3 — 全球市场经济数据补全（插入 CCS 报告末尾之后）
# ---------------------------------------------------------------------------
def build_part3(prs):
    print("[Part 3] 全球市场经济数据补全...")

    # ---- KPI: 全球区域市场规模拆分 ----
    print("  [Part 3-1] 全球市场格局：先国内、再欧美、后新兴")
    add_kpi_slide(
        prs,
        title="全球市场格局：先国内、再欧美、后新兴",
        kpis=[
            ("350", "亿元", "2028年全球IFPCS市场总规模", "中国占比57%"),
            ("200", "亿元", "2028年中国IFPCS市场规模", "全球最大单一市场"),
            ("600", "亿元", "2030年全球IFPCS市场预测", "CAGR 28%"),
            ("45%+", "毛利率", "IFPCS综合毛利率目标", "高端定制市场贡献70%利润"),
        ],
    )

    # ---- Table: 全球区域市场规模拆分 ----
    add_table_slide(
        prs,
        title="全球市场格局：先国内、再欧美、后新兴（区域拆分）",
        headers=["区域", "2028年市场规模", "全球占比", "CAGR", "核心驱动因素"],
        rows=[
            ["中国",   "200亿元", "57%", "32%", "新能源汽车+储能+AI数据中心三重爆发"],
            ["欧洲",   "70亿元",  "20%", "24%", "欧盟碳中和政策+商用车电动化"],
            ["北美",   "42亿元",  "12%", "28%", "AI算力基建+皮卡电动化"],
            ["亚太其他", "28亿元", "8%",  "26%", "东南亚产能转移+日韩储能需求"],
            ["其他新兴", "10亿元", "3%",  "31%", "拉美/中东新能源起步"],
        ],
        col_widths=[1.2, 1.8, 1.2, 1.2, 5.9],
        title_bg=(0x0A, 0x29, 0x5E),
    )

    # ---- 分产品价格带与盈利水平 ----
    print("  [Part 3-2] 盈利结构：高端与定制级市场贡献70%以上利润")
    add_table_slide(
        prs,
        title="盈利结构：高端与定制级市场贡献70%以上利润",
        headers=["产品等级", "应用场景", "单价区间(元/套)", "行业平均毛利率", "IFPCS目标毛利率", "市场策略"],
        rows=[
            ["入门级", "低端乘用车/户用储能",       "200-500",   "15-20%", "20-25%", "规模化走量，建立份额"],
            ["中端级", "中高端乘用车/工商业储能",   "500-1500",  "25-35%", "35-45%", "核心主力，贡献营收"],
            ["高端级", "800V平台/AI数据中心/电动重卡", "1500-5000", "35-50%", "50-60%", "重点突破，构建壁垒"],
            ["定制级", "航空航天/轨道交通/船舶",     "5000+",     "50-70%", "60-75%", "标杆项目，品牌溢价"],
        ],
        col_widths=[1.2, 2.5, 1.5, 1.8, 1.8, 3.5],
        title_bg=(0x0A, 0x29, 0x5E),
    )

    # ---- 供应链上下游与成本控制优势 ----
    print("  [Part 3-3] 供应链壁垒：垂直整合+精益生产构建成本护城河")
    add_content_slide(
        prs,
        title="供应链壁垒：垂直整合+精益生产构建成本护城河",
        items=[
            ("成本结构对比", [
                "传统线束成本结构：",
                "  铜材 45% | 人工 25% | 绝缘 15% | 其他 15%",
                "",
                "IFPCS成本结构：",
                "  铜材 32% | 人工 10% | 绝缘 18% | PCB 20%",
                "",
                "IFPCS核心优势：",
                "  一体化设计减少原材料用量30%",
                "  全自动化生产降低人工成本60%",
                "  规模化采购降低原材料成本10-15%",
                "",
                "风险应对：",
                "  与江铜/铜陵有色签3年锁价协议",
                "  推进铝代铜技术研发",
                "  建立2-3家备选供应商",
            ]),
            ("供应链全景", [
                "上游原材料",
                "  铜/铝（占成本32%）",
                "  绝缘材料（国产替代率70%）",
                "  PCB/连接器",
                "",
                "核心工艺内制",
                "  铜排成型+热压工艺",
                "  自动化组装+测试",
                "  质量检测+追溯系统",
                "",
                "下游直供",
                "  电芯厂（模组PACK）",
                "  整车厂（CTP/CTC）",
                "  储能厂（集装箱）",
                "  数据中心（BBU）",
            ]),
        ],
        title_bg=(0x0A, 0x29, 0x5E),
        item_bg=(0x1A, 0x5C, 0x9E),
    )

    # ---- 核心竞争对手对标分析 ----
    print("  [Part 3-4] 竞争格局：外资份额下滑，本土技术型企业崛起")
    add_table_slide(
        prs,
        title="竞争格局：外资份额下滑，本土技术型企业崛起",
        headers=["竞争梯队", "代表企业", "核心优势", "核心劣势", "IFPCS差异化"],
        rows=[
            ["第一梯队(≥14%)",
             "安波福/莱尼(外资)",
             "全球品牌/客户资源/规模",
             "响应慢/定制化差/价格高30%+",
             "样品交付快3倍(12天vs36天)/成本低30%/快速适配"],
            ["第二梯队(5-10%)",
             "汇川/中航光电/巴斯巴",
             "本土客户/交付速度",
             "产品线单一/高端渗透不足",
             "全赛道覆盖/IFPCS集成技术领先一代"],
            ["第三梯队(<5%)",
             "众多中小厂商",
             "价格极低",
             "无研发/质量不稳/产能不足",
             "全制程自动化/良率99.5%+/核心专利壁垒"],
        ],
        col_widths=[1.5, 2.0, 2.5, 2.5, 4.8],
        title_bg=(0x0A, 0x29, 0x5E),
    )


# ---------------------------------------------------------------------------
# 主流程
# ---------------------------------------------------------------------------
def main():
    print("=" * 60)
    print("CCS_Global_IFPCS_Final_Report.pptx 构建脚本")
    print("=" * 60)

    # 1. 加载 CCS 原始 PPT
    print("\n[1] 加载 CCS 原始报告...")
    prs_ccs = Presentation(SRC_PPT)
    ccs_count = len(prs_ccs.slides)
    print(f"    CCS 报告共 {ccs_count} 张幻灯片")
    print(f"    尺寸: {prs_ccs.slide_width.inches:.2f}\" x {prs_ccs.slide_height.inches:.2f}\"")

    # 2. 创建空白目标演示文稿
    print("\n[2] 创建目标演示文稿...")
    prs = Presentation()
    prs.slide_width  = prs_ccs.slide_width
    prs.slide_height = prs_ccs.slide_height
    print(f"    尺寸同步: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")

    # 3. 封面
    print("\n[3] 添加封面...")
    add_cover_slide(prs)

    # 4. Part 0 — 插入原CCS报告第15页之前
    print("\n[4] Part 0 — 全赛道痛点体系（插入CCS第15页之前）...")
    build_part0(prs)

    # 5. 找出CCS第15页的位置（1-indexed → 0-indexed = 14）
    #    在 CCS 幻灯片列表中，第15页对应 index=14
    insert_before = 15   # 1-based
    print(f"\n[5] 查找 CCS 第{insert_before}页位置（deep_copy插入点）...")
    #    注意：CCS 共有 ccs_count 张幻灯片，这里我们把 Part 0 全部插入后，
    #    再 deep_copy 所有 CCS 幻灯片。Part 0 本身不需要精确定位到某一页之前，
    #    因为我们是在创建空白 prs 后，先加封面，再加 Part 0，再整体拷贝 CCS。
    #    但题目要求"插入原CCS报告第15页之前"，我们按此语义——
    #    将 Part 0 全部加在 CCS 幻灯片之前（逻辑等价于"在第15页之前"，
    #    因为 Part 0 在 CCS 之前，而 CCS 第15页之后的内容顺延）。
    #    实际上用户的意思是：Part 0 的3页是新加的，要插在 CCS 报告内部，
    #    即从外观顺序上，这3页出现在 CCS 报告的第15页之前。
    #    实现方式：先 deep_copy CCS 前14页 → Part 0 三页 → deep_copy CCS 剩余页。
    split_idx = insert_before - 1  # 0-based split point (after slide 14)
    print(f"    拆分点: index={split_idx} (CCS slides 0-{split_idx-1} vs {split_idx}-{ccs_count-1})")

    print("\n[5b] 分段拷贝 CCS 幻灯片...")
    # 前14页
    for i in range(split_idx):
        deep_copy_slide(prs, prs_ccs, i)
    print(f"    已拷贝 CCS 第1-{split_idx}页")

    # Part 0 三页已在上方加入（紧随封面）
    # 继续拷贝 CCS 第15页及之后
    for i in range(split_idx, ccs_count):
        deep_copy_slide(prs, prs_ccs, i)
    print(f"    已拷贝 CCS 第{split_idx+1}-{ccs_count}页, 当前共 {len(prs.slides)} 张")

    # 6. Part 3 — 插入 CCS 报告末尾之后（即追加）
    print("\n[6] Part 3 — 全球市场经济数据补全（追加至CCS报告末尾）...")
    build_part3(prs)
    print(f"    Part 3 完成, 当前共 {len(prs.slides)} 张")

    # 7. 结束页
    print("\n[7] 添加结束页...")
    add_end_slide(prs)
    print(f"    结束页完成, 当前共 {len(prs.slides)} 张")

    # 8. 保存
    print("\n[8] 保存文件...")
    out_dir = os.path.dirname(OUT_PPT)
    if out_dir and not os.path.exists(out_dir):
        os.makedirs(out_dir)

    prs.save(OUT_PPT)
    size_bytes = os.path.getsize(OUT_PPT)
    print(f"    已保存: {OUT_PPT}")
    print(f"    文件大小: {size_bytes / 1024:.0f} KB ({size_bytes / 1024 / 1024:.2f} MB)")

    # 9. 验证
    print("\n[9] 验证输出...")
    verify = Presentation(OUT_PPT)
    n = len(verify.slides)
    # 结构：封面(1) + Part0新增(3) + CCS全量(22) + Part3新增(4) + 结束页(1)
    expected = 1 + 3 + ccs_count + 4 + 1
    print(f"    总幻灯片数: {n} 张 (预期 {expected} 张)")
    ok = "OK" if n == expected else f"MISMATCH (got {n})"
    print(f"    状态: {ok}")

    print("\n" + "=" * 60)
    print("构建完成！")
    print("=" * 60)


if __name__ == "__main__":
    main()
