# -*- coding: utf-8 -*-
"""
合并两个PPT文件
Global_Storage_Market_2025_2030.pptx + CCS市场分析报告.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

SRC1 = r"d:/claude mini max 2.7/storage_market_output/Global_Storage_Market_2025_2030.pptx"
SRC2 = r"d:/李總項目/CCS集成母排与线束技术整合市场分析报告.pptx"
OUTPUT = r"d:/claude mini max 2.7/storage_market_output/CCS_全球储能市场综合报告_2025-2030.pptx"

# ── 加载两个演示文稿 ──
prs1 = Presentation(SRC1)
prs2 = Presentation(SRC2)

print(f"源文件1 (全球储能): {len(prs1.slides)} 张幻灯片")
print(f"源文件2 (CCS报告): {len(prs2.slides)} 张幻灯片")

# ── 创建新的合并演示文稿 ──
merged = Presentation()
merged.slide_width = Inches(13.33)
merged.slide_height = Inches(7.5)

def clone_slide(merged_prs, source_prs, slide_idx):
    """将 source_prs 的第 slide_idx 张幻灯片克隆到 merged_prs"""
    src_slide = source_prs.slides[slide_idx]
    # 使用空白布局
    blank_layout = merged_prs.slide_layouts[6]
    new_slide = merged_prs.slides.add_slide(blank_layout)

    # 复制背景
    try:
        if src_slide.background.fill.type is not None:
            new_slide.background.fill.solid()
            try:
                new_slide.background.fill.fore_color.rgb = src_slide.background.fill.fore_color.rgb
            except Exception:
                pass
    except Exception:
        pass

    # 复制所有形状
    for shape in src_slide.shapes:
        if shape.has_text_frame:
            # 文本框
            el = shape.element
            new_el = type(el)(el)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        elif shape.shape_type == 13:  # 图片
            # 图片：复制元素
            el = shape.element
            new_el = type(el)(el)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        else:
            # 其他形状（表格、线条等）
            el = shape.element
            new_el = type(el)(el)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

def add_divider_slide(merged_prs, text):
    """插入分隔页"""
    slide = merged_prs.slides.add_slide(merged_prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x0A, 0x29, 0x5E)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    return slide

# ── 第一步：插入封面（综合标题页） ──
cover = merged.slides.add_slide(merged.slide_layouts[6])
cover.background.fill.solid()
cover.background.fill.fore_color.rgb = RGBColor(0x0A, 0x29, 0x5E)

txTitle = cover.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.3), Inches(2.0))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "CCS集成母排 × 全球储能市场\n综合分析报告"
p.font.size = Pt(42)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

txSub = cover.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.3), Inches(1.0))
tf2 = txSub.text_frame
p2 = tf2.paragraphs[0]
p2.text = "CCS集成母排与线束技术整合市场分析  |  2025-2030全球储能市场经济规模"
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(0xCC, 0xD9, 0xEC)
p2.alignment = PP_ALIGN.CENTER

txDate = cover.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5))
tf3 = txDate.text_frame
p3 = tf3.paragraphs[0]
p3.text = "2026年4月 | 内部研究资料"
p3.font.size = Pt(14)
p3.font.color.rgb = RGBColor(0x88, 0x99, 0xBB)
p3.alignment = PP_ALIGN.CENTER

print("封面已创建")

# ── 第二步：复制 CCS 报告全部幻灯片 ──
print("\n正在复制 CCS 报告幻灯片...")
for i in range(len(prs2.slides)):
    clone_slide(merged, prs2, i)
    print(f"  CCS 第 {i+1}/{len(prs2.slides)} 张完成")

# ── 第三步：插入储能市场分析分隔页 ──
add_divider_slide(merged, "第二部分\n全球储能及新能源关联市场分析 2025-2030")
print("\n分隔页已插入")

# ── 第四步：复制全球储能分析PPT全部幻灯片 ──
print("\n正在复制全球储能分析幻灯片...")
for i in range(len(prs1.slides)):
    clone_slide(merged, prs1, i)
    print(f"  储能 第 {i+1}/{len(prs1.slides)} 张完成")

# ── 保存 ──
merged.save(OUTPUT)
print(f"\n合并完成！")
print(f"输出文件: {OUTPUT}")
print(f"总计: {len(merged.slides)} 张幻灯片")
print(f"  - CCS报告: {len(prs2.slides)} 张")
print(f"  - 全球储能分析: {len(prs1.slides)} 张")
