# -*- coding: utf-8 -*-
"""
PPT合併腳本 v4 - 完整保留所有媒體資源
使用 _element 深層複製 + 完整關聯修復
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree
import os, shutil, re, zipfile

SRC1 = r'd:/claude mini max 2.7/storage_market_output/Global_Storage_Market_2025_2030.pptx'
SRC2 = r'd:/claude mini max 2.7/storage_market_output/CCS_Global_Storage_2025_2030.pptx'
OUTPUT = r'd:/claude mini max 2.7/storage_market_output/CCS_Global_Report_MERGED.pptx'

# 先複製 SRC1 作為基礎
shutil.copy2(SRC1, OUTPUT)
prs = Presentation(OUTPUT)
print(f"Loaded base PPT: {len(prs.slides)} slides")

# 讀取 CCS PPT
prs2 = Presentation(SRC2)
print(f"CCS PPT: {len(prs2.slides)} slides")

NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_CT = 'http://schemas.openxmlformats.org/package/2006/content-types'

def get_rels_map(pres):
    """取得 presentation.xml.rels 中的所有Relationship"""
    tree = pres.part._element
    rels_tree = pres.part.rels_container.rels_element
    rels = {}
    for rel in rels_tree:
        rel_id = rel.get('Id')
        rel_target = rel.get('Target')
        rel_type = rel.get('Type', '').split('/')[-1]
        rels[rel_id] = {'target': rel_target, 'type': rel_type}
    return rels

def copy_slide_to_presentation(target_prs, source_prs, slide_idx):
    """使用 _element 深層複製方式將幻燈片加入目標簡報"""
    src_slide = source_prs.slides[slide_idx]
    src_tree = src_slide.part._element

    # 使用空白佈局新增一張新幻燈片
    blank_layout = target_prs.slide_layouts[6]  # 通常是 blank layout
    new_slide = target_prs.slides.add_slide(blank_layout)
    new_tree = new_slide.part._element

    # 清空新slide內容（保留sldId元素外殼）
    # 找到 cSld 元素
    cSld_new = new_tree.find(qn('p:cSld'))
    cSld_src = src_tree.find(qn('p:cSld'))

    if cSld_new is not None and cSld_src is not None:
        # 深層複製 spTree（形狀樹）
        spTree_new = cSld_new.find(qn('p:spTree'))
        spTree_src = cSld_src.find(qn('p:spTree'))

        # 清除新slide的spTree內容
        if spTree_new is not None:
            for child in list(spTree_new):
                spTree_new.remove(child)

        # 深層複製所有形狀
        if spTree_src is not None:
            for child in list(spTree_src):
                spTree_new.append(deepcopy(child))

        # 處理背景
        bg_new = cSld_new.find(qn('p:bg'))
        bg_src = cSld_src.find(qn('p:bg'))
        if bg_new is not None:
            cSld_new.remove(bg_new)
        if bg_src is not None:
            cSld_new.insert(0, deepcopy(bg_src))

    return new_slide

# 找出現有最大 rId
def get_max_rid(pres):
    rels_tree = pres.part.rels_container.rels_element
    max_rid = 0
    for rel in rels_tree:
        try:
            rid = int(rel.get('Id').replace('rId', ''))
            max_rid = max(max_rid, rid)
        except:
            pass
    return max_rid

# 找出現有最大 slide id
def get_max_sldid(pres):
    pres_tree = pres.part._element
    sldIdLst = pres_tree.find(qn('p:sldIdLst'))
    max_id = 0
    if sldIdLst is not None:
        for sldId in sldIdLst:
            try:
                eid = int(sldId.get('id'))
                max_id = max(max_id, eid)
            except:
                pass
    return max_id

# ============================================================
# 第一步：新增 CCS PPT 的空白幻燈片結構（收集所有關聯）
# ============================================================
print("\n第一步：建立 CCS 報告結構...")

# 收集CCS PPT的所有媒體關聯
ccs_media_rels = {}
for i, slide in enumerate(prs2.slides):
    for rel_id, rel in slide.part.rels.items():
        if rel.reltype.endswith('image'):
            ccs_media_rels[rel_id] = rel.target_ref

print(f"CCS PPT 圖片關聯數: {len(ccs_media_rels)}")

# ============================================================
# 第二步：使用 ZIP 直接合併（最完整方式）
# ============================================================
print("\n第二步：使用ZIP合併...")

shutil.copy2(SRC1, OUTPUT + '.tmp')

with zipfile.ZipFile(SRC1, 'r') as z1, \
     zipfile.ZipFile(SRC2, 'r') as z2, \
     zipfile.ZipFile(OUTPUT + '.tmp', 'r') as ztmp:

    files1 = set(z1.namelist())
    files2 = set(z2.namelist())

    # 讀取 CCS PPT 的所有媒體文件
    ccs_media = {}
    for f in files2:
        if f.startswith('ppt/media/'):
            ccs_media[f] = z2.read(f)

    # 讀取 CCS PPT 的 slides
    ccs_slides = {}
    for f in sorted(files2):
        if re.match(r'ppt/slides/slide\d+\.xml', f):
            ccs_slides[f] = z2.read(f)

    ccs_slide_rels = {}
    for f in sorted(files2):
        if re.match(r'ppt/slides/_rels/slide\d+\.xml\.rels', f):
            ccs_slide_rels[f] = z2.read(f)

# ============================================================
# 第三步：重建合併後的簡報
# ============================================================
print("\n第三步：重建簡報結構...")

# 以 CCS PPT 為基礎，加入全球儲能 PPT 的內容
# 策略：將兩個 PPT 的所有媒體都複製到合併文件中

TEMP_OUTPUT = OUTPUT + '.build'
shutil.copy2(SRC2, TEMP_OUTPUT)  # 以 CCS PPT 為基礎

# 讀取 CCS PPT 的結構
with zipfile.ZipFile(TEMP_OUTPUT, 'r') as z:
    base_files = set(z.namelist())
    base_ct = z.read('[Content_Types].xml')
    base_pres = z.read('ppt/presentation.xml')
    base_pres_rels = z.read('ppt/_rels/presentation.xml.rels')

# 解析現有結構
pres_tree = etree.fromstring(base_pres)
pres_rels_tree = etree.fromstring(base_pres_rels)
ct_tree = etree.fromstring(base_ct)

# 計算新的 slide ID 和 rId
existing_ids = [int(el.get('id')) for el in pres_tree.iter(qn('p:sldId'))]
next_sld_id = max(existing_ids) + 1 if existing_ids else 256

existing_rids = []
for rel in pres_rels_tree:
    try:
        existing_rids.append(int(rel.get('Id').replace('rId', '')))
    except:
        pass
next_rid = max(existing_rids) + 1 if existing_rids else 1

print(f"next_sld_id={next_sld_id}, next_rid={next_rid}")

# 讀取全球儲能 PPT 的所有 slides 和媒體
with zipfile.ZipFile(SRC1, 'r') as z1:
    storage_slides = {}
    for f in sorted(z1.namelist()):
        if re.match(r'ppt/slides/slide\d+\.xml', f):
            storage_slides[f] = z1.read(f)

    storage_slide_rels = {}
    for f in sorted(z1.namelist()):
        if re.match(r'ppt/slides/_rels/slide\d+\.xml\.rels', f):
            storage_slide_rels[f] = z1.read(f)

    storage_media = {}
    for f in sorted(z1.namelist()):
        if f.startswith('ppt/media/'):
            storage_media[f] = z1.read(f)

    # 讀取 global ct 和 pres
    g_ct = z1.read('[Content_Types].xml')
    g_pres = z1.read('ppt/presentation.xml')
    g_pres_rels = z1.read('ppt/_rels/presentation.xml.rels')
    g_rels_tree = etree.fromstring(g_pres_rels)

    # 全球儲能的 slide 數量
    g_slide_files = sorted([f for f in z1.namelist() if re.match(r'ppt/slides/slide\d+\.xml', f)])
    g_slide_count = len(g_slide_files)
    print(f"Storage slides to add: {g_slide_count}")

# 在 presentation.xml 的 sldIdLst 中加入全球儲能 PPT 的 slides
g_pres_tree = etree.fromstring(g_pres)
g_sldIdLst = g_pres_tree.find(qn('p:sldIdLst'))
sldIdLst = pres_tree.find(qn('p:sldIdLst'))

# 建立新 slide 的映射表：舊 slide id -> 新 rId + target
# 同時更新 g_slide_files 的 rId 映射

# 讀取全球儲能 PPT 的 presentation.xml.rels 建立舊->新 rId 映射
g_rels_map = {}
for rel in g_rels_tree:
    rid = rel.get('Id')
    target = rel.get('Target')
    g_rels_map[rid] = target

# 為全球儲能 PPT 的 slides 建立新 rId
storage_rid_map = {}  # old_rId -> new_rId
for rel in g_rels_tree:
    rid = rel.get('Id')
    if rid and rid.startswith('rId'):
        storage_rid_map[rid] = f'rId{next_rid}'
        next_rid += 1

# 讀取全球儲能 PPT 的 slides 數量（presentation.xml 中的 sldIdLst）
g_sld_ids = []
for sldId in g_sldIdLst:
    g_sld_ids.append({
        'id': sldId.get('id'),
        'rId': sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    })

print(f"Storage PPT has {len(g_sld_ids)} slides")

# 建立 storage slide index: slide index -> slide file
g_slide_file_map = {}
for f in g_slide_files:
    m = re.search(r'slide(\d+)\.xml', f)
    if m:
        g_slide_file_map[int(m.group(1))] = f

# 讀取 CCS PPT 的最大 slide number
with zipfile.ZipFile(TEMP_OUTPUT, 'r') as z:
    ccs_slide_nums = []
    for f in sorted(z.namelist()):
        m = re.match(r'ppt/slides/slide(\d+)\.xml', f)
        if m:
            ccs_slide_nums.append(int(m.group(1)))
    ccs_max_slide_num = max(ccs_slide_nums) if ccs_slide_nums else 0

print(f"CCS max slide number: {ccs_max_slide_num}")

# 建立 storage slide 的新名稱（避免與 CCS slides 衝突）
# 新的 storage slides: slide_ccs_max+1, slide_ccs_max+2, ...

# 讀取 storage slide 對應的 rId
storage_slide_rid = {}  # new slide file name -> rId
for sld_info in g_sld_ids:
    old_rid = sld_info['rId']
    new_rid = storage_rid_map.get(old_rid, f'rId{next_rid}')
    old_target = g_rels_map.get(old_rid, '')

    # 找到對應的 slide file
    m = re.search(r'slide(\d+)\.xml', old_target)
    if m:
        slide_num = int(m.group(1))
        if slide_num in g_slide_file_map:
            orig_file = g_slide_file_map[slide_num]
            storage_slide_rid[orig_file] = new_rid

# 將全球儲能 PPT 的 slides 加入到 sldIdLst
for sld_info in g_sld_ids:
    old_rid = sld_info['rId']
    new_rid = storage_rid_map.get(old_rid)
    new_sld_id = etree.SubElement(sldIdLst, qn('p:sldId'))
    new_sld_id.set('id', str(next_sld_id))
    new_sld_id.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', new_rid)
    next_sld_id += 1

# 更新 presentation.xml.rels 加入 storage slides 的關聯
for sld_info in g_sld_ids:
    old_rid = sld_info['rId']
    new_rid = storage_rid_map.get(old_rid)
    old_target = g_rels_map.get(old_rid, '')

    # 找到對應的 slide file
    m = re.search(r'slide(\d+)\.xml', old_target)
    if m:
        slide_num = int(m.group(1))
        if slide_num in g_slide_file_map:
            orig_file = g_slide_file_map[slide_num]
            # 新的 slide file 名稱
            ccs_max_num = ccs_max_slide_num
            new_slide_num = ccs_max_num + slide_num
            new_slide_file = f'ppt/slides/slide{new_slide_num}.xml'

            rel_elem = etree.SubElement(pres_rels_tree, 'Relationship')
            rel_elem.set('Id', new_rid)
            rel_elem.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
            rel_elem.set('Target', new_slide_file)

# 更新 [Content_Types].xml
def ensure_ct(ov_el, part_name, content_type):
    existing = ov_el.find(f'{{{NS_CT}}}Override[@PartName="{part_name}"]')
    if existing is None:
        new_ov = etree.SubElement(ov_el, f'{{{NS_CT}}}Override')
        new_ov.set('PartName', part_name)
        new_ov.set('ContentType', content_type)

# 加入 storage slides 的 content type
for i in range(1, g_slide_count + 1):
    new_slide_num = ccs_max_slide_num + i
    ensure_ct(ct_tree, f'/ppt/slides/slide{new_slide_num}.xml',
              'application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
    ensure_ct(ct_tree, f'/ppt/slides/_rels/slide{new_slide_num}.xml.rels',
              'application/vnd.openxmlformats-officedocument.presentationml.slide.related+xml')

# ============================================================
# 第四步：寫入新的 ZIP 文件
# ============================================================
print("\n第四步：寫入合併後的 PPTX 文件...")

FINAL_OUTPUT = OUTPUT
with zipfile.ZipFile(TEMP_OUTPUT, 'r') as z_orig, \
     zipfile.ZipFile(FINAL_OUTPUT, 'w', zipfile.ZIP_DEFLATED) as zout:

    # 寫入更新後的 presentation.xml
    zout.writestr('ppt/presentation.xml',
                  etree.tostring(pres_tree, xml_declaration=True, encoding='UTF-8', standalone=True))

    # 寫入更新後的 presentation.xml.rels
    zout.writestr('ppt/_rels/presentation.xml.rels',
                  etree.tostring(pres_rels_tree, xml_declaration=True, encoding='UTF-8', standalone=True))

    # 寫入更新後的 [Content_Types].xml
    zout.writestr('[Content_Types].xml',
                  etree.tostring(ct_tree, xml_declaration=True, encoding='UTF-8', standalone=True))

    # 複製 CCS PPT 的所有內容（以它為基礎）
    with zipfile.ZipFile(SRC2, 'r') as z2:
        for f in sorted(z2.namelist()):
            if f not in ['ppt/presentation.xml', 'ppt/_rels/presentation.xml.rels', '[Content_Types].xml']:
                zout.writestr(f, z2.read(f))

    # 加入全球儲能 PPT 的 slides（重新命名）
    for i in range(1, g_slide_count + 1):
        orig_file = g_slide_file_map.get(i)
        if orig_file:
            new_slide_num = ccs_max_slide_num + i
            new_slide_file = f'ppt/slides/slide{new_slide_num}.xml'
            # 加入 slide 內容
            slide_xml = storage_slides[orig_file]

            # 修正 slide rels 中的 rId 引用
            slide_tree = etree.fromstring(slide_xml)

            # 找到並修正 slide 的 rId
            # 需要更新 slide rels
            new_slide_rels_file = f'ppt/slides/_rels/slide{new_slide_num}.xml.rels'
            if orig_file in storage_slide_rels:
                orig_rels = storage_slide_rels[orig_file]
                rels_tree = etree.fromstring(orig_rels)

                # 更新 rels 中的所有 target（圖片路徑）
                # 圖片在 ppt/media/ 下，直接引用即可
                # 但需要修正舊 rId -> 新 rId

                # 建立舊->新 rId 映射（只用於 slide rels）
                slide_old_rid_map = {}
                slide_new_rid_counter = 1
                for rel in rels_tree:
                    old_id = rel.get('Id')
                    new_id = f'rId{1000 + slide_new_rid_counter}'  # 使用較大的起始值避免衝突
                    slide_old_rid_map[old_id] = new_id
                    rel.set('Id', new_id)
                    slide_new_rid_counter += 1

                # 寫入修正後的 rels
                zout.writestr(new_slide_rels_file,
                              etree.tostring(rels_tree, xml_declaration=True, encoding='UTF-8', standalone=True))

                # 更新 slide XML 中的 rId 引用
                for elem in slide_tree.iter():
                    for attr in elem.attrib:
                        if attr.endswith('}id') or attr == 'id':
                            val = elem.attrib[attr]
                            if val in slide_old_rid_map:
                                elem.attrib[attr] = slide_old_rid_map[val]

            zout.writestr(new_slide_file,
                          etree.tostring(slide_tree, xml_declaration=True, encoding='UTF-8', standalone=True))

    # 加入全球儲能 PPT 的所有媒體文件
    media_count = 0
    for f, data in sorted(storage_media.items()):
        zout.writestr(f, data)
        media_count += 1

    # 複製 global theme 等其他必要文件（如果CCS PPT中沒有）
    with zipfile.ZipFile(SRC1, 'r') as z1:
        for f in sorted(z1.namelist()):
            if f not in zout.namelist() and f not in storage_slides and f not in storage_slide_rels:
                try:
                    zout.writestr(f, z1.read(f))
                except:
                    pass

    print(f"Media files added: {media_count}")

# 清理
os.remove(TEMP_OUTPUT)
if os.path.exists(OUTPUT + '.tmp'):
    os.remove(OUTPUT + '.tmp')

size = os.path.getsize(FINAL_OUTPUT)
print(f"\nFile size: {size/1024:.0f} KB ({size/1024/1024:.2f} MB)")

# 驗證
try:
    verify_prs = Presentation(FINAL_OUTPUT)
    print(f"Verified: {len(verify_prs.slides)} slides")
except Exception as e:
    print(f"Verification ERROR: {e}")

print(f"\nOutput: {FINAL_OUTPUT}")
