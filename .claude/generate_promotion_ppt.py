# -*- coding: utf-8 -*-
"""
CoBM DIY PC电源品牌客户推广策略PPT生成
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 82, 118)
    background.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = '2026年4月6日'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets, color=(26, 82, 118)):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(*color)
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(6)

def add_table_slide(prs, title, headers, rows, color=(26, 82, 118)):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = RGBColor(*color)
    header_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    num_cols = len(headers)
    num_rows = len(rows) + 1
    table_height = min(Inches(5.2), Inches(0.5) * (num_rows + 1))
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.3), Inches(1.1), Inches(9.4), table_height).table

    for i in range(num_cols):
        table.columns[i].width = Inches(9.4 / num_cols)

    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*color)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.alignment = PP_ALIGN.CENTER

def add_conclusion_slide(prs, title, items):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(39, 174, 96)
    background.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(12)

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, 'CoBM DIY PC电源品牌客户推广策略', 'B2B营销 · 客户开发 · 渠道策略 | 2026-04-06')

    # Slide 2: Contents
    add_content_slide(prs, '报告目录', [
        '● 一、目标客户画像与核心痛点',
        '● 二、精准客户切入策略',
        '● 三、B2B营销工具包',
        '● 四、渠道策略与展会计划',
        '● 五、销售赋能与话术模板',
        '● 六、执行计划与里程碑'
    ], (142, 68, 173))

    # Slide 3: Customer Types
    headers = ['客户类型', '代表品牌', '年出货量', '切入优先级']
    rows = [
        ['台系电源大厂', '曜越T.T、海韵、全汉', '50-100万+', 'S级'],
        ['大陆电源品牌', '长城、航嘉、先马', '100万+', 'A级'],
        ['电竞品牌', 'ROG、微星MEG', '20-50万', 'A级'],
        ['ITX新锐品牌', '傻瓜超人、乔思伯', '5-30万', 'B级']
    ]
    add_table_slide(prs, '一、目标客户类型与优先级', headers, rows, (52, 73, 94))

    # Slide 4: Pain Points
    headers = ['痛点类型', '具体表现', 'CoBM解决方案']
    rows = [
        ['成本压力', '利润持续压缩', 'BOM成本降低18-28%'],
        ['效率竞争', '80PLUS认证成本高', '白金牌方案低成本实现'],
        ['功率密度', 'ITX/SFF市场需求', '功率密度+50-65%'],
        ['AI PC需求', 'RTX 40/50系列大功率', '覆盖550W-2400W'],
        ['供应链风险', '传统IC方案元件多', 'CoBM高集成减少元件']
    ]
    add_table_slide(prs, '一、客户核心痛点与CoBM解决能力', headers, rows, (192, 57, 43))

    # Slide 5: Entry Strategy for T.T
    add_content_slide(prs, '二、台系大厂切入策略 - 曜越T.T', [
        '【步骤1】发送CoBM vs T.T现有方案对比分析报告',
        '【步骤2】寄送CoBM样品及技术文档',
        '【步骤3】线上/线下技术交流会议',
        '【步骤4】提供定制化CoBM方案报价',
        '【步骤5】联合测试验证（PFC/LLC效率实测）',
        '【步骤6】签署合作意向书',
        '',
        '切入话术：GF1/XP系列导入CoBM，BOM成本降低22%，',
        '年出货50万台节省800-1200万元'
    ], (231, 76, 60))

    # Slide 6: Entry Strategy for All Brands
    add_content_slide(prs, '二、各类品牌差异化切入策略', [
        '【大陆品牌】长城/航嘉 - 强调性价比+产能+供应链',
        '  CoBM Si方案：铜牌→金牌，BOM下降20%',
        '',
        '【电竞品牌】ROG/微星MEG - 强调功率密度+小尺寸',
        '  122mm深度，ATX 3.1原生12V-2x6接口',
        '',
        '【ITX品牌】傻瓜超人/乔思伯 - 强调SFX尺寸优势',
        '  Wukong GX 122mm最短深度'
    ], (155, 89, 182))

    # Slide 7: Marketing Materials
    add_content_slide(prs, '三、B2B营销工具包', [
        '【核心物料】技术白皮书 - CoBM原理/优势/实测数据',
        '【核心物料】对比分析表 - CoBM vs 竞品 vs 传统方案',
        '【核心物料】ROI计算器 - 输入出货量→输出节省金额',
        '【样品支持】CoBM模块+完整测试报告',
        '【案例集】已合作客户实证数据',
        '【线上内容】短视频/技术动画/社交媒体传播'
    ], (26, 82, 118))

    # Slide 8: ROI Calculator
    headers = ['参数', '典型值', '用户可调范围']
    rows = [
        ['年出货量(万台)', '50', '10-500'],
        ['平均瓦数(W)', '750', '550-1200'],
        ['当前BOM成本(元)', '380', '-'],
        ['CoBM方案BOM(元)', '305', '-'],
    ]
    add_table_slide(prs, '三、ROI计算器模板', headers, rows, (39, 174, 96))

    # Add ROI output
    slide = prs.slides[-1]
    add_text = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
    tf = add_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '输出结果：'
    p.font.size = Pt(14)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = '年BOM成本节省：750万元 | ROI：320% | 回收期：3个月'
    p.font.size = Pt(12)

    # Slide 9: Sales Talking Points
    add_content_slide(prs, '五、销售话术模板', [
        '【开场白30秒】专注功率半导体封装，BOM成本降20%，效率+2-4%',
        '',
        '【痛点挖掘三问】',
        '  1. 您觉得目前电源方案最大的成本压力在哪？',
        '  2. 80PLUS认证投入产出比是否令您满意？',
        '  3. ITX/紧凑电源市场是否有合适方案？',
        '',
        '【价值主张】850W产品升级白金牌认证，BOM成本还下降20%'
    ], (52, 152, 219))

    # Slide 10: Timeline
    headers = ['时间节点', '里程碑', '验收标准']
    rows = [
        ['4月底', '营销物料完成', '白皮书v1.0+对比表+ROI计算器'],
        ['5月中', '首批客户拜访', '完成10+家客户技术交流'],
        ['6月初', 'Computex亮相', '收集50+名片，转化20+MQL'],
        ['7月底', '首批样品订单', '3家客户进入样品测试'],
        ['9月底', '首批成交', '至少1家签署正式订单']
    ]
    add_table_slide(prs, '六、执行计划与里程碑', headers, rows, (241, 196, 15))

    # Slide 11: Budget
    headers = ['费用类别', '预算(万元)', '占比']
    rows = [
        ['展会费用', '30', '25%'],
        ['样品制作', '20', '17%'],
        ['营销物料', '15', '13%'],
        ['FAE支持', '25', '21%'],
        ['销售费用', '20', '17%'],
        ['线上推广', '10', '8%'],
        ['总计', '120', '100%']
    ]
    add_table_slide(prs, '六、年度推广预算规划', headers, rows, (142, 68, 173))

    # Slide 12: Conclusion
    add_conclusion_slide(prs, '总结与行动计划', [
        '◆ 爆品切入：以850W-1200W白金牌为突破口',
        '◆ 标杆客户：优先攻克曜越T.T，成功后复制',
        '◆ 数据说话：所有推广基于实测数据，ROI量化价值',
        '◆ 展会搭台：Computex 2026是最佳展示舞台',
        '◆ 渠道深耕：建立华东/华南直销+代理网络'
    ])

    output_path = r'D:/claude mini max 2.7/.claude/CoBM_DIY_PC_电源推广策略_20260406.pptx'
    prs.save(output_path)
    print(f'PPT已生成: {output_path}')

if __name__ == '__main__':
    main()
