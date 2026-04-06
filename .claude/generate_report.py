import pandas as pd
from pptx import Presentation
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from datetime import datetime

# ========== 讀取數據 ==========
# CoBM Excel
cobm_path = r'D:\深圳\台灣瑞聲源\副本CoBM产品资料.xlsx'
cobm_xl = pd.ExcelFile(cobm_path)
cobm_data = {}
for sheet in cobm_xl.sheet_names:
    cobm_data[sheet] = pd.read_excel(cobm_path, sheet_name=sheet)

# BQT PPT
bqt_path = r'D:\深圳\台灣瑞聲源\BQT高瓦数电源双技术深度分析报告（专家版）.pptx'
prs = Presentation(bqt_path)
bqt_text = []
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            bqt_text.append(shape.text)

# ========== 競品數據（已研究） ==========
competitors = [
    {"brand": "Corsair", "wattage": "550-1600W", "ic": "GaN+DSP", "cert": "Titanium", "tech": "iCUE軟件生態"},
    {"brand": "Seasonic", "wattage": "650-2200W", "ic": "雙變壓器", "cert": "Titanium", "tech": "12年質保"},
    {"brand": "Super Flower", "wattage": "650-2800W", "ic": "專利連接器", "cert": "Platinum", "tech": "全球最高2800W"},
    {"brand": "be quiet!", "wattage": "550-1600W", "ic": "Silent Wings", "cert": "Titanium", "tech": "極致靜音"},
    {"brand": "ASUS ROG", "wattage": "1000-1600W", "ic": "GaN+OLED", "cert": "Titanium", "tech": "GPU-First專利"},
    {"brand": "MSI MEG", "wattage": "550-1300W", "ic": "AI整合", "cert": "Platinum", "tech": "MSI Center"},
    {"brand": "Gigabyte AORUS", "wattage": "850-1200W", "ic": "LCD顯示", "cert": "Platinum", "tech": "AORUS生態"},
    {"brand": "Thermaltake", "wattage": "550-1500W", "ic": "ATX 3.1", "cert": "Platinum", "tech": "RGB燈效"},
    {"brand": "EVGA", "wattage": "550-1600W", "ic": "雙滾珠軸承", "cert": "Titanium", "tech": "10年質保"},
    {"brand": "Cooler Master", "wattage": "550-1300W", "ic": "TRM溫控", "cert": "Platinum", "tech": "性價比高"},
]

# ========== 生成HTML報告 ==========
html = f"""
<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8">
<title>CoBM-BQT 電源產品線互補分析報告</title>
<style>
body {{ font-family: Arial, sans-serif; margin: 40px; background: #f5f5f5; }}
.container {{ max-width: 1200px; margin: auto; background: white; padding: 40px; box-shadow: 0 2px 10px rgba(0,0,0,0.1); }}
h1 {{ color: #1a1a1a; border-bottom: 3px solid #0066cc; padding-bottom: 15px; }}
h2 {{ color: #0066cc; margin-top: 30px; }}
h3 {{ color: #333; }}
table {{ width: 100%; border-collapse: collapse; margin: 20px 0; }}
th, td {{ border: 1px solid #ddd; padding: 12px; text-align: left; }}
th {{ background: #0066cc; color: white; }}
tr:nth-child(even) {{ background: #f9f9f9; }}
.highlight {{ background: #e6f3ff; padding: 20px; border-radius: 8px; margin: 20px 0; }}
.footer {{ margin-top: 40px; padding-top: 20px; border-top: 1px solid #ddd; color: #666; }}
</style>
</head>
<body>
<div class="container">
<h1>CoBM-BQT 電源產品線互補分析報告</h1>
<p><strong>生成日期：</strong>{datetime.now().strftime('%Y-%m-%d %H:%M')}</p>

<h2>一、CoBM 產品 IC 方案分析</h2>
<div class="highlight">
<p>以下為 CoBM 產品資料的核心 IC 方案摘要：</p>
<pre>{str(cobm_data)}</pre>
</div>

<h2>二、BQT 電源產品線缺口分析</h2>
<div class="highlight">
<p><strong>BQT 技術特點：</strong></p>
<pre>{chr(10).join(bqt_text[:20])}</pre>
</div>

<h2>三、競品對標分析</h2>
<table>
<tr><th>品牌</th><th>功率段</th><th>IC方案</th><th>效率認證</th><th>技術特點</th></tr>
"""

for c in competitors:
    html += f"<tr><td>{c['brand']}</td><td>{c['wattage']}</td><td>{c['ic']}</td><td>{c['cert']}</td><td>{c['tech']}</td></tr>"

html += """
</table>

<h2>四、互補策略建議</h2>
<div class="highlight">
<h3>短期機會（填補BQT缺口）：</h3>
<ol>
<li>開發800W-1200W主流段產品</li>
<li>推出80+ Platinum認證系列</li>
<li>支援ATX 3.1標準</li>
</ol>

<h3>中期策略：</h3>
<ol>
<li>引進GaN技術，提升效率</li>
<li>開發1500W+旗艦產品</li>
<li>建立軟件生態（監控/調控）</li>
</ol>

<h3>長期規劃：</h3>
<ol>
<li>2000W+超旗艦產品線</li>
<li>數字化/智能電源</li>
<li>全模組化/定制化服務</li>
</ol>
</div>

<h2>五、市場機會分析</h2>
<table>
<tr><th>功率段</th><th>市場機會</th><th>建議品牌</th></tr>
<tr><td>2000W+</td><td>僅Seasonic/Super Flower，競爭少</td><td>優先引進</td></tr>
<tr><td>1500W</td><td>Corsair/ASUS/ROG覆蓋</td><td>中期目標</td></tr>
<tr><td>1000-1200W</td><td>紅海市場，所有品牌佈局</td><td>性價比突圍</td></tr>
<tr><td>800W以下</td><td>利基市場</td><td>定制化服務</td></tr>
</table>

<div class="footer">
<p><strong>免責聲明：</strong>本報告僅供參考，不構成投資建議。</p>
</div>
</div>
</body>
</html>
"""

# 保存HTML報告
report_path = r'D:\claude mini max 2.7\.claude\CoBM_BQT_Analysis_Report.html'
with open(report_path, 'w', encoding='utf-8') as f:
    f.write(html)
print(f"報告已生成: {report_path}")

# ========== 發送郵件 ==========
def send_email(html_content, to_email):
    from_email = "your_163_email@163.com"  # 替換為您的163郵箱
    password = "your_auth_code"  # 替換為您的授權碼
    
    msg = MIMEMultipart('alternative')
    msg['Subject'] = 'CoBM-BQT 電源產品線互補分析報告'
    msg['From'] = from_email
    msg['To'] = to_email
    
    part = MIMEText(html_content, 'html', 'utf-8')
    msg.attach(part)
    
    try:
        with smtplib.SMTP_SSL('smtp.163.com', 465) as server:
            server.login(from_email, password)
            server.sendmail(from_email, to_email, msg.as_string())
        print(f"郵件已發送至: {to_email}")
    except Exception as e:
        print(f"郵件發送失敗: {e}")

# 發送郵件
send_email(html, 'h13751019800@163.com')
print("完成!")
