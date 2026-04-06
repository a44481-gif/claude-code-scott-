# -*- coding: utf-8 -*-
"""
PPT 生成器 v3 - 增強版
增加細分市場分析頁面 + PDF 導出
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
import json
import os
from datetime import datetime


# ============================================================================
# 配色方案
# ============================================================================
class Colors:
    PRIMARY = RGBColor(0, 51, 102)
    SECONDARY = RGBColor(0, 102, 153)
    ACCENT = RGBColor(51, 153, 255)
    SUCCESS = RGBColor(46, 204, 113)
    WARNING = RGBColor(241, 196, 15)
    TEXT_DARK = RGBColor(51, 51, 51)
    TEXT_LIGHT = RGBColor(102, 102, 102)
    WHITE = RGBColor(255, 255, 255)
    BG = RGBColor(245, 247, 250)


# ============================================================================
# 工具函數
# ============================================================================
def set_cell_bg(cell, rgb_color):
    """設置單元格背景色"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    color_str = f"{rgb_color[0]:02X}{rgb_color[1]:02X}{rgb_color[2]:02X}"
    solidFill = parse_xml(
        f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:srgbClr val="{color_str}"/></a:solidFill>'
    )
    tcPr.append(solidFill)


def add_text(slide, left, top, width, height, text, size=14, bold=False,
             color=None, align=PP_ALIGN.LEFT, bg_color=None):
    """添加文字框"""
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or Colors.TEXT_DARK
    p.alignment = align

    if bg_color:
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color

    return box


def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    """添加形狀"""
    shape = slide.shapes.add_shape(shape_type, Cm(left), Cm(top), Cm(width), Cm(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_page_number(slide, num, total):
    """添加頁碼"""
    add_text(slide, 30, 17.5, 3, 0.8, f"{num} / {total}", 10, False, Colors.TEXT_LIGHT, PP_ALIGN.RIGHT)


def add_header(slide, title, page_num, total):
    """添加頁面標題頭"""
    add_text(slide, 1, 0.3, 31, 1.2, title, 26, True, Colors.PRIMARY)
    add_shape(slide, 1, 1.5, 31.8, 0.05, Colors.ACCENT)
    add_page_number(slide, page_num, total)


# ============================================================================
# 頁面模板
# ============================================================================
def create_prs():
    """創建簡報"""
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    return prs


def slide_cover(prs, title, subtitle):
    """封面頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_shape(slide, 0, 0, 33.867, 0.5, Colors.PRIMARY)
    add_shape(slide, 0, 18.55, 33.867, 0.5, Colors.PRIMARY)

    add_text(slide, 1, 5, 31, 3, title, 38, True, Colors.PRIMARY, PP_ALIGN.CENTER)
    add_text(slide, 1, 9, 31, 2, subtitle, 18, False, Colors.TEXT_LIGHT, PP_ALIGN.CENTER)
    add_text(slide, 1, 13, 31, 1, datetime.now().strftime('%Y-%m-%d'), 14, False, Colors.TEXT_LIGHT, PP_ALIGN.CENTER)


def slide_summary(prs, data, page_num, total):
    """執行摘要頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Executive Summary | 執行摘要", page_num, total)

    segments = data.get('segments', [])
    total_market = data.get('total_market_5yr', 0)

    content = f"""
Analysis Scope: Global Energy Storage & New Energy Market (2025-2030)
Coverage: {len(segments)} Market Segments

KEY FINDINGS:

Market Scale:
  - 5-Year Total: {total_market:,.0f} Billion RMB
  - 2030 Projected: {sum(s['scale_2030'] for s in segments if segments):,.0f} Billion RMB

Fastest Growing (CAGR):
  - AI Storage: 63.9%
  - Robot Storage: 59.1%
  - Drone Storage: 56.3%

Largest Market:
  - EV: {segments[0]['cumulative_5yr'] if segments else 0:,.0f} Billion RMB ({segments[0]['share_2025_2030'] if segments else 0}%)

Data Sources: BloombergNEF / GGII / SNE Research
"""
    add_text(slide, 1, 2, 31, 15, content.strip(), 14, False, Colors.TEXT_DARK)


def slide_overview_table(prs, data, page_num, total):
    """市場總覽表格頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Market Overview | 市場總覽 (Billion RMB)", page_num, total)

    segments = data.get('segments', [])

    rows = len(segments) + 1
    cols = 7
    table = slide.shapes.add_table(rows, cols, Cm(0.5), Cm(2), Cm(33), Cm(15)).table

    widths = [2.5, 5, 5, 4.5, 5, 5, 4.5]
    for i, w in enumerate(widths):
        table.columns[i].width = Cm(w)

    headers = ["Rank", "Segment", "5Y Total", "CAGR%", "2030 Scale", "Share%", "Unit"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = Colors.WHITE
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_cell_bg(cell, Colors.PRIMARY)

    units = ["", "萬輛", "GWh", "GW", "GWh", "GWh", "GWh", "", "", "", "", ""]

    for row_idx, seg in enumerate(segments, 1):
        table.cell(row_idx, 0).text = str(row_idx)
        table.cell(row_idx, 1).text = seg['name']
        table.cell(row_idx, 2).text = f"{seg['cumulative_5yr']:,.0f}"
        table.cell(row_idx, 3).text = f"{seg['cagr_2025_2030']:.1f}%"
        table.cell(row_idx, 4).text = f"{seg['scale_2030']:,.0f}"
        table.cell(row_idx, 5).text = f"{seg['share_2025_2030']:.2f}%"
        table.cell(row_idx, 6).text = units[row_idx] if row_idx < len(units) else ""

        for col in range(7):
            cell = table.cell(row_idx, col)
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if col != 1 else PP_ALIGN.LEFT

            if row_idx <= 3:
                cell.text_frame.paragraphs[0].font.bold = True


def slide_cagr_ranking(prs, data, page_num, total):
    """CAGR 增速排名頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Growth Ranking | 增速排名 (CAGR 2025-2030)", page_num, total)

    fastest = data.get('fastest_growing', [])[:8]

    # 左側：增速排名
    add_text(slide, 1, 2, 15, 1, "Fastest Growing", 16, True, Colors.SECONDARY)
    y = 3.5
    for i, seg in enumerate(fastest[:5], 1):
        bar_width = seg['cagr_2025_2030'] / 70 * 12
        add_shape(slide, 1, y, bar_width, 1.3, Colors.ACCENT)
        add_text(slide, 1.2, y + 0.3, 12, 0.8, f"{i}. {seg['name']}", 12, True, Colors.WHITE)
        add_text(slide, 13.5, y + 0.3, 3, 0.8, f"{seg['cagr_2025_2030']}%", 12, True, Colors.TEXT_DARK)
        y += 1.6

    # 右側：市場份額
    segments = data.get('segments', [])[:5]
    add_text(slide, 17, 2, 15, 1, "Market Share", 16, True, Colors.SECONDARY)
    y = 3.5
    for i, seg in enumerate(segments, 1):
        bar_width = seg['share_2025_2030'] / 70 * 12
        add_shape(slide, 17, y, bar_width, 1.3, Colors.PRIMARY)
        add_text(slide, 17.2, y + 0.3, 12, 0.8, f"{i}. {seg['name']}", 12, True, Colors.WHITE)
        add_text(slide, 29.5, y + 0.3, 3, 0.8, f"{seg['share_2025_2030']}%", 12, True, Colors.TEXT_DARK)
        y += 1.6


def slide_trend(prs, data, page_num, total):
    """年度趨勢頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Annual Trend | 年度趨勢 (Billion RMB)", page_num, total)

    segments = data.get('segments', [])[:3]
    years = ['2025', '2026', '2027', '2028', '2029', '2030']
    colors_list = [Colors.PRIMARY, Colors.ACCENT, Colors.SUCCESS]

    # 圖例
    y_legend = 2.3
    for i, seg in enumerate(segments[:2]):
        add_text(slide, 1 + i*10, y_legend, 8, 0.6, f"■ {seg['name']}", 11, False, colors_list[i])
    add_text(slide, 1, y_legend + 0.8, 30, 0.5, "Market Scale Trends", 12, True, Colors.TEXT_DARK)

    # 柱狀圖
    x_start = 2
    bar_w = 1.2
    gap = 0.3
    y_base = 14
    scale = 0.003

    for seg_idx, seg in enumerate(segments[:2]):
        max_val = seg['scale_2030']
        for i, year in enumerate(years):
            val = seg['yearly_data'].get(year, 0)
            height = val * scale
            x = x_start + i * (bar_w * 2 + gap) + seg_idx * bar_w
            add_shape(slide, x, y_base - height, bar_w, height, colors_list[seg_idx])

    # X 軸標籤
    for i, year in enumerate(years):
        x = x_start + i * (bar_w * 2 + gap) + bar_w / 2
        add_text(slide, x, 14.5, 3, 0.6, year, 10, False, Colors.TEXT_DARK, PP_ALIGN.CENTER)


def slide_segment_detail(prs, data, page_num, total):
    """細分市場詳情頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Segment Analysis | 細分市場分析", page_num, total)

    segments = data.get('segments', [])

    # EV Market
    if len(segments) > 0:
        ev = segments[0]
        add_text(slide, 1, 2, 10, 1, "EV | 電動汽車", 14, True, Colors.PRIMARY)
        add_shape(slide, 1, 3, 10, 0.05, Colors.ACCENT)
        ev_data = f"5Y Total: {ev['cumulative_5yr']:,.0f}B | CAGR: {ev['cagr_2025_2030']}%\n2030: {ev['scale_2030']:,.0f}B (from {ev['scale_2025']:,.0f}B)"
        add_text(slide, 1, 3.3, 10, 3, ev_data, 11, False, Colors.TEXT_DARK)

    # Truck Market
    if len(segments) > 1:
        truck = segments[1]
        add_text(slide, 12, 2, 10, 1, "Truck | 電動貨卡車", 14, True, Colors.PRIMARY)
        add_shape(slide, 12, 3, 10, 0.05, Colors.ACCENT)
        truck_data = f"5Y Total: {truck['cumulative_5yr']:,.0f}B | CAGR: {truck['cagr_2025_2030']}%\n2030: {truck['scale_2030']:,.0f}B (from {truck['scale_2025']:,.0f}B)"
        add_text(slide, 12, 3.3, 10, 3, truck_data, 11, False, Colors.TEXT_DARK)

    # Emerging Markets
    if len(segments) > 5:
        add_text(slide, 23, 2, 10, 1, "Emerging | 新興市場", 14, True, Colors.SECONDARY)
        add_shape(slide, 23, 3, 10, 0.05, Colors.SUCCESS)
        emerging = segments[5:8]
        em_text = "\n".join([f"• {s['name']}: {s['cagr_2025_2030']}%" for s in emerging])
        add_text(slide, 23, 3.3, 10, 5, em_text, 11, False, Colors.TEXT_DARK)

    # Market Comparison
    add_text(slide, 1, 8, 31, 1, "Market Comparison | 市場對比", 14, True, Colors.PRIMARY)
    add_shape(slide, 1, 9, 31, 0.05, Colors.ACCENT)

    # 簡單對比表格
    rows = min(5, len(segments))
    table = slide.shapes.add_table(rows + 1, 3, Cm(1), Cm(9.3), Cm(31), Cm(5)).table
    table.columns[0].width = Cm(12)
    table.columns[1].width = Cm(10)
    table.columns[2].width = Cm(9)

    headers = ["Segment", "5Y Growth Rate", "Market Scale"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_cell_bg(cell, Colors.PRIMARY)
        cell.text_frame.paragraphs[0].font.color.rgb = Colors.WHITE

    for row_idx in range(rows):
        seg = segments[row_idx]
        table.cell(row_idx + 1, 0).text = seg['name']
        table.cell(row_idx + 1, 1).text = f"{seg['cagr_2025_2030']:.1f}%"
        table.cell(row_idx + 1, 2).text = f"{seg['cumulative_5yr']:,.0f}B"

        for col in range(3):
            cell = table.cell(row_idx + 1, col)
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT


def slide_conclusion(prs, data, page_num, total):
    """結論頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Conclusion | 結論與展望", page_num, total)

    segments = data.get('segments', [])
    fastest = data.get('fastest_growing', [])[:3]

    content = f"""
KEY INSIGHTS:

1. EV Market Dominates
   - 5-Year Total: {segments[0]['cumulative_5yr'] if segments else 0:,.0f} Billion RMB
   - Share: {segments[0]['share_2025_2030'] if segments else 0}%
   - CAGR: {segments[0]['cagr_2025_2030'] if segments else 0}%

2. Emerging Segments Show Explosive Growth
   - AI Storage: CAGR {fastest[0]['cagr_2025_2030'] if len(fastest) > 0 else 0}%
   - Robot Storage: CAGR {fastest[1]['cagr_2025_2030'] if len(fastest) > 1 else 0}%
   - Drone Storage: CAGR {fastest[2]['cagr_2025_2030'] if len(fastest) > 2 else 0}%

3. Market Trends
   - Accelerated electrification and intelligence
   - Explosive growth in storage demand
   - Continuous emergence of new technologies and applications

Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
"""
    add_text(slide, 1, 2, 31, 15, content.strip(), 14, False, Colors.TEXT_DARK)


def slide_thank_you(prs, page_num, total):
    """感謝頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_shape(slide, 0, 0, 33.867, 0.5, Colors.PRIMARY)
    add_shape(slide, 0, 18.55, 33.867, 0.5, Colors.PRIMARY)

    add_text(slide, 1, 6, 31, 3, "Thank You", 48, True, Colors.PRIMARY, PP_ALIGN.CENTER)
    add_text(slide, 1, 10, 31, 2, "感謝聆聽", 24, False, Colors.TEXT_LIGHT, PP_ALIGN.CENTER)

    add_page_number(slide, page_num, total)


# ============================================================================
# 主函數
# ============================================================================
def generate_enhanced_ppt(data, output_path):
    """生成增強版 PPT"""
    prs = create_prs()

    segments = data.get('segments', [])
    total_slides = 8

    # 頁面 1: 封面
    slide_cover(prs,
        "Global Energy Storage Market\nAnalysis Report 2025-2030",
        f"11 Segments | {data.get('total_market_5yr', 0):,.0f} Billion RMB")

    # 頁面 2: 執行摘要
    slide_summary(prs, data, 2, total_slides)

    # 頁面 3: 市場總覽
    slide_overview_table(prs, data, 3, total_slides)

    # 頁面 4: 增速排名
    slide_cagr_ranking(prs, data, 4, total_slides)

    # 頁面 5: 年度趨勢
    slide_trend(prs, data, 5, total_slides)

    # 頁面 6: 細分市場
    slide_segment_detail(prs, data, 6, total_slides)

    # 頁面 7: 結論
    slide_conclusion(prs, data, 7, total_slides)

    # 頁面 8: 感謝
    slide_thank_you(prs, 8, total_slides)

    # 保存
    os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
    prs.save(output_path)
    print(f"[OK] PPT generated: {output_path}")
    return output_path


def export_to_pdf(pptx_path):
    """導出為 PDF（需要手動在 PowerPoint 中完成）"""
    pdf_path = pptx_path.replace('.pptx', '.pdf')
    print(f"[INFO] To export PDF, open {pptx_path} in PowerPoint and Save As PDF")
    print(f"[INFO] PDF will be saved to: {pdf_path}")
    return pdf_path


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python ppt_generator_v3.py <market_json> [output_pptx]")
        sys.exit(1)

    json_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else "market_report_v3.pptx"

    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    generate_enhanced_ppt(data, output_path)

    # 提示 PDF 導出
    export_to_pdf(output_path)
