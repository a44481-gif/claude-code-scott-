# -*- coding: utf-8 -*-
"""
PPT 分析器 - 自訂 Skill 實現
分析 PowerPoint 檔案內容、提取數據、生成結構化報告
"""

import sys
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
from datetime import datetime


def analyze_presentation(file_path, output_path=None):
    """
    分析 PPT 完整內容

    Args:
        file_path: PPT 檔案路徑
        output_path: 輸出 JSON 檔案路徑（可選）

    Returns:
        dict: 分析結果
    """
    prs = Presentation(file_path)

    result = {
        "file": file_path,
        "file_name": os.path.basename(file_path),
        "analyzed_at": datetime.now().isoformat(),
        "total_slides": len(prs.slides),
        "statistics": {
            "total_shapes": 0,
            "total_tables": 0,
            "total_images": 0,
            "total_charts": 0,
            "total_text_boxes": 0
        },
        "slides": []
    }

    for idx, slide in enumerate(prs.slides, 1):
        slide_data = {
            "slide_number": idx,
            "title": "",
            "shapes_count": 0,
            "tables": [],
            "images": [],
            "charts": [],
            "text_content": [],
            "full_text": ""
        }

        for shape in slide.shapes:
            result["statistics"]["total_shapes"] += 1
            slide_data["shapes_count"] += 1

            # 提取標題和文字
            if shape.has_text_frame:
                text = shape.text_frame.text.strip() if shape.text_frame.text else ""
                if text:
                    if not slide_data["title"]:
                        slide_data["title"] = text[:100]
                    slide_data["text_content"].append(text)
                    result["statistics"]["total_text_boxes"] += 1

            # 提取表格
            if shape.has_table:
                result["statistics"]["total_tables"] += 1
                table_data = extract_table(shape.table)
                slide_data["tables"].append(table_data)

            # 提取圖表
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                result["statistics"]["total_charts"] += 1
                slide_data["charts"].append({"type": "chart", "name": shape.name})

            # 提取圖片
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                result["statistics"]["total_images"] += 1
                slide_data["images"].append({"type": "image", "name": shape.name})

        # 組合完整文字
        slide_data["full_text"] = "\n".join(slide_data["text_content"])
        result["slides"].append(slide_data)

    # 保存結果
    if output_path:
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(result, f, ensure_ascii=False, indent=2)
        print(f"[OK] 分析結果已保存: {output_path}")

    return result


def extract_table(table):
    """提取表格數據"""
    rows_data = []
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            row_data.append(cell.text.strip())
        rows_data.append(row_data)
    return rows_data


def print_summary(result):
    """打印摘要報告"""
    print("=" * 60)
    print(f"PPT 分析報告: {result['file_name']}")
    print("=" * 60)
    print(f"總頁數: {result['total_slides']}")
    print(f"總形狀: {result['statistics']['total_shapes']}")
    print(f"文字框: {result['statistics']['total_text_boxes']}")
    print(f"表格: {result['statistics']['total_tables']}")
    print(f"圖表: {result['statistics']['total_charts']}")
    print(f"圖片: {result['statistics']['total_images']}")
    print("-" * 60)

    for slide in result['slides']:
        title = slide['title'][:40] if slide['title'] else '(無標題)'
        print(f"\n[頁 {slide['slide_number']}] {title}")
        if slide['tables']:
            print(f"  - 包含 {len(slide['tables'])} 個表格")
        if slide['charts']:
            print(f"  - 包含 {len(slide['charts'])} 個圖表")
        if slide['images']:
            print(f"  - 包含 {len(slide['images'])} 張圖片")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python ppt_analyzer.py <pptx_file> [output_json]")
        print("範例: python ppt_analyzer.py report.pptx analysis.json")
        sys.exit(1)

    file_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None

    if not os.path.exists(file_path):
        print(f"[ERROR] 檔案不存在: {file_path}")
        sys.exit(1)

    print(f"正在分析: {file_path}")
    result = analyze_presentation(file_path, output_path)
    print_summary(result)
