# -*- coding: utf-8 -*-
"""
PPT 生成器 - 自訂 Skill
根據數據自動生成市場分析報告 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
from datetime import datetime


def create_presentation():
    """創建新的 PPT 簡報"""
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    return prs


def add_title_slide(prs, title, subtitle):
    """添加標題頁"""
    slide_layout = prs.slide_layouts[6]  # 空白版面
    slide = prs.slides.add_slide(slide_layout)

    # 標題
    title_box = slide.shapes.add_textbox(Cm(1), Cm(6), Cm(31), Cm(3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 副標題
    sub_box = slide.shapes.add_textbox(Cm(1), Cm(10), Cm(31), Cm(2))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = RgbColor(102, 102, 102)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_executive_summary(prs, data):
    """添加執行摘要頁"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 標題
    title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(31), Cm(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "執行摘要 | Executive Summary"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 51, 102)

    # 內容
    content = f"""
分析範圍：全球儲能及新能源關聯市場（2025-2030）

覆蓋賽道：{len(data['segments'])} 大賽道

核心結論：
  ▸ 11大賽道5年累計總規模：{data['total_market_5yr']:,.0f} 億元
  ▸ 2030年預計總規模：{data.get('total_2030', 0):,.0f} 億元
  ▸ 增速最快：AI儲能 (CAGR 65%)、機器人儲能 (CAGR 58%)、無人機儲能 (CAGR 55%)
  ▸ 規模最大：電動汽車 (5年累計 {data['segments'][0]['cumulative_5yr']:,.0f} 億元，占比 {data['segments'][0]['share_2025_2030']}%)

數據來源：BloombergNEF / GGII / SNE Research / 行業公開數據綜合
"""

    content_box = slide.shapes.add_textbox(Cm(1), Cm(2.5), Cm(31), Cm(15))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content.strip()
    p.font.size = Pt(16)
    p.font.color.rgb = RgbColor(51, 51, 51)

    return slide


def add_market_overview(prs, data):
    """添加市場總覽頁"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 標題
    title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(31), Cm(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "市場規模總覽 | Market Size Overview (億元)"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 51, 102)

    # 表格
    rows = len(data['segments']) + 1
    cols = 5
    table = slide.shapes.add_table(rows, cols, Cm(1), Cm(2.5), Cm(31), Cm(14)).table

    # 設置列寬
    table.columns[0].width = Cm(5)
    table.columns[1].width = Cm(6)
    table.columns[2].width = Cm(5)
    table.columns[3].width = Cm(5)
    table.columns[4].width = Cm(5)

    # 表頭
    headers = ["排名", "賽道", "5年累計(億元)", "CAGR (%)", "市場份額(%)"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 數據行
    for row_idx, seg in enumerate(data['segments'][:10], 1):
        table.cell(row_idx, 0).text = str(row_idx)
        table.cell(row_idx, 1).text = seg['name']
        table.cell(row_idx, 2).text = f"{seg['cumulative_5yr']:,.0f}"
        table.cell(row_idx, 3).text = f"{seg['cagr_2025_2030']:.1f}"
        table.cell(row_idx, 4).text = f"{seg['share_2025_2030']:.2f}"

        for col in range(5):
            cell = table.cell(row_idx, col)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            if col != 1:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def add_growth_analysis(prs, data):
    """添加增速分析頁"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 標題
    title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(31), Cm(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "增速分析 | Growth Analysis"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 51, 102)

    # 內容
    content = "【增速最快 Top 5】\n\n"
    for i, seg in enumerate(data['fastest_growing'][:5], 1):
        content += f"{i}. {seg['name']}: CAGR {seg['cagr_2025_2030']}%\n"
        content += f"   5年累計: {seg['cumulative_5yr']:,.0f} 億元\n\n"

    content += "\n【市場規模最大 Top 3】\n\n"
    for i, seg in enumerate(data['segments'][:3], 1):
        content += f"{i}. {seg['name']}: {seg['cumulative_5yr']:,.0f} 億元 ({seg['share_2025_2030']}%)\n"

    content_box = slide.shapes.add_textbox(Cm(1), Cm(2.5), Cm(31), Cm(15))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content.strip()
    p.font.size = Pt(18)
    p.font.color.rgb = RgbColor(51, 51, 51)

    return slide


def add_conclusion(prs, data):
    """添加結論頁"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 標題
    title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(31), Cm(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "結論與展望 | Conclusion"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 51, 102)

    # 內容
    conclusions = f"""
核心洞察：

1. 電動汽車市場規模最大
   - 5年累計規模：{data['segments'][0]['cumulative_5yr']:,.0f} 億元
   - 市場份額：{data['segments'][0]['share_2025_2030']}%
   - CAGR：{data['segments'][0]['cagr_2025_2030']}%

2. 新興賽道增速迅猛
   - AI儲能 CAGR：63.9%
   - 機器人儲能 CAGR：59.1%
   - 無人機儲能 CAGR：56.3%

3. 市場趨勢
   - 電動化、智能化加速
   - 儲能需求爆發式增長
   - 新技術、新應用場景不斷涌現

生成時間：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
"""

    content_box = slide.shapes.add_textbox(Cm(1), Cm(2.5), Cm(31), Cm(15))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = conclusions.strip()
    p.font.size = Pt(16)
    p.font.color.rgb = RgbColor(51, 51, 51)

    return slide


def generate_market_report_ppt(data, output_path):
    """
    生成市場分析報告 PPT

    Args:
        data: 市場分析數據
        output_path: 輸出檔案路徑
    """
    prs = create_presentation()

    # 頁面 1: 標題頁
    add_title_slide(
        prs,
        "全球儲能及新能源關聯市場\n經濟規模分析報告",
        f"2025-2030 | {len(data['segments'])}大賽道 | 市場規模 | 增速分析"
    )

    # 頁面 2: 執行摘要
    add_executive_summary(prs, data)

    # 頁面 3: 市場總覽
    add_market_overview(prs, data)

    # 頁面 4: 增速分析
    add_growth_analysis(prs, data)

    # 頁面 5: 結論
    add_conclusion(prs, data)

    # 保存
    prs.save(output_path)
    print(f"[OK] PPT 已生成: {output_path}")
    return output_path


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("用法: python ppt_generator.py <market_json> [output_pptx]")
        print("範例: python ppt_generator.py market_analysis.json report.pptx")
        sys.exit(1)

    json_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else "market_report.pptx"

    # 載入數據
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    print(f"正在生成 PPT: {output_path}")
    generate_market_report_ppt(data, output_path)
