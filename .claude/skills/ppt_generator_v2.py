# -*- coding: utf-8 -*-
"""
PPT 生成器 v2 - 優化版
根據市場數據生成專業的市場分析報告 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import json
from datetime import datetime
import os


# ============================================================================
# 配色方案
# ============================================================================
class Colors:
    PRIMARY = RGBColor(0, 51, 102)      # 深藍
    SECONDARY = RGBColor(0, 102, 153)   # 中藍
    ACCENT = RGBColor(51, 153, 255)     # 亮藍
    TEXT_DARK = RGBColor(51, 51, 51)    # 深灰文字
    TEXT_LIGHT = RGBColor(102, 102, 102) # 淺灰文字
    WHITE = RGBColor(255, 255, 255)
    BG_LIGHT = RGBColor(245, 247, 250)  # 淺灰背景


# ============================================================================
# 工具函數
# ============================================================================
def set_cell_style(cell, font_size=11, bold=False, align=PP_ALIGN.CENTER, bg_color=None):
    """設置表格單元格樣式"""
    cell.text_frame.paragraphs[0].font.size = Pt(font_size)
    cell.text_frame.paragraphs[0].font.bold = bold
    cell.text_frame.paragraphs[0].alignment = align
    cell.text_frame.paragraphs[0].font.color.rgb = Colors.TEXT_DARK

    if bg_color:
        set_cell_bg(cell, bg_color)


def set_cell_bg(cell, color):
    """設置單元格背景色"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = parse_xml(
        f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:srgbClr val="{color[0]:02X}{color[1]:02X}{color[2]:02X}"/></a:solidFill>'
    )
    tcPr.append(solidFill)


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    """添加文字框"""
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color or Colors.TEXT_DARK
    p.alignment = align
    return box


def add_page_number(slide, page_num, total):
    """添加頁碼"""
    add_text_box(slide, 30, 17.5, 3, 1, f"{page_num} / {total}", 10, False, Colors.TEXT_LIGHT, PP_ALIGN.RIGHT)


# ============================================================================
# 頁面生成函數
# ============================================================================
def create_presentation():
    """創建簡報"""
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    return prs


def add_cover_slide(prs, title, subtitle, date=None):
    """封面頁"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景裝飾線
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(33.867), Cm(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.PRIMARY
    shape.line.fill.background()

    shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(18.75), Cm(33.867), Cm(0.3))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = Colors.PRIMARY
    shape2.line.fill.background()

    # 主標題
    add_text_box(slide, 1, 5, 31, 4, title, 40, True, Colors.PRIMARY, PP_ALIGN.CENTER)

    # 副標題
    add_text_box(slide, 1, 9, 31, 2, subtitle, 18, False, Colors.TEXT_LIGHT, PP_ALIGN.CENTER)

    # 日期
    date_str = date or datetime.now().strftime('%Y-%m-%d')
    add_text_box(slide, 1, 13, 31, 1, date_str, 14, False, Colors.TEXT_LIGHT, PP_ALIGN.CENTER)

    return slide


def add_executive_summary(prs, data, page_num, total):
    """執行摘要頁"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 頁面標題
    add_text_box(slide, 1, 0.5, 31, 1.5, "執行摘要 | Executive Summary", 28, True, Colors.PRIMARY)
    add_page_number(slide, page_num, total)

    # 分隔線
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1), Cm(2), Cm(31.8), Cm(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.ACCENT
    shape.line.fill.background()

    # 內容
    total_2030 = sum(seg['scale_2030'] for seg in data['segments'])
    top3_cagr = [seg['cagr_2025_2030'] for seg in data['fastest_growing'][:3]]

    content = f"""
分析範圍：全球儲能及新能源關聯市場（2025-2030）
覆蓋賽道：{len(data['segments'])} 大賽道

【核心結論】

▸ 市場總規模
  • 11大賽道5年累計：{data['total_market_5yr']:,.0f} 億元
  • 2030年預計規模：{total_2030:,.0f} 億元

▸ 增速最快的賽道
  • AI儲能 (CAGR 63.9%)
  • 機器人儲能 (CAGR 59.1%)
  • 無人機儲能 (CAGR 56.3%)

▸ 規模最大的賽道
  • 電動汽車 (5年累計 {data['segments'][0]['cumulative_5yr']:,.0f} 億元，占比 {data['segments'][0]['share_2025_2030']}%)

▸ 數據來源
  BloombergNEF / GGII / SNE Research / 行業公開數據
"""
    add_text_box(slide, 1, 2.5, 31, 14, content.strip(), 14, False, Colors.TEXT_DARK)

    return slide


def add_market_overview_table(prs, data, page_num, total):
    """市場總覽表格頁"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 頁面標題
    add_text_box(slide, 1, 0.5, 31, 1.5, "市場規模總覽 | Market Size Overview (億元)", 26, True, Colors.PRIMARY)
    add_page_number(slide, page_num, total)

    # 分隔線
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1), Cm(2), Cm(31.8), Cm(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.ACCENT
    shape.line.fill.background()

    # 表格
    rows = len(data['segments']) + 1
    cols = 6
    table = slide.shapes.add_table(rows, cols, Cm(0.8), Cm(2.3), Cm(32), Cm(15)).table

    # 列寬
    widths = [3, 5.5, 6, 5, 5, 5]
    for i, w in enumerate(widths):
        table.columns[i].width = Cm(w)

    # 表頭
    headers = ["排名", "賽道", "5年累計(億元)", "CAGR (%)", "2030規模", "市場份額(%)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        set_cell_style(cell, 11, True, PP_ALIGN.CENTER, Colors.PRIMARY)
        cell.text_frame.paragraphs[0].font.color.rgb = Colors.WHITE

    # 數據
    for row_idx, seg in enumerate(data['segments'], 1):
        table.cell(row_idx, 0).text = str(row_idx)
        table.cell(row_idx, 1).text = seg['name']
        table.cell(row_idx, 2).text = f"{seg['cumulative_5yr']:,.0f}"
        table.cell(row_idx, 3).text = f"{seg['cagr_2025_2030']:.1f}"
        table.cell(row_idx, 4).text = f"{seg['scale_2030']:,.0f}"
        table.cell(row_idx, 5).text = f"{seg['share_2025_2030']:.2f}"

        for col in range(6):
            cell = table.cell(row_idx, col)
            font_size = 10
            bold = (row_idx <= 3)  # 前三名加粗
            set_cell_style(cell, font_size, bold, PP_ALIGN.CENTER)

    return slide


def add_growth_ranking(prs, data, page_num, total):
    """增速排名頁"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 頁面標題
    add_text_box(slide, 1, 0.5, 31, 1.5, "增速排名 | Growth Ranking (CAGR 2025-2030)", 26, True, Colors.PRIMARY)
    add_page_number(slide, page_num, total)

    # 分隔線
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1), Cm(2), Cm(31.8), Cm(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.ACCENT
    shape.line.fill.background()

    # 左側：增速最快
    add_text_box(slide, 1, 2.5, 15, 1, "增速最快 Top 5", 16, True, Colors.SECONDARY)
    y = 4
    for i, seg in enumerate(data['fastest_growing'][:5], 1):
        bar_width = seg['cagr_2025_2030'] / 70 * 12  # 按比例
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1), Cm(y), Cm(bar_width), Cm(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = Colors.ACCENT
        shape.line.fill.background()

        add_text_box(slide, 1.2, y + 0.3, 12, 0.8, f"{i}. {seg['name']}", 12, True, Colors.WHITE)
        add_text_box(slide, 13.5, y + 0.3, 3, 0.8, f"{seg['cagr_2025_2030']}%", 12, True, Colors.TEXT_DARK, PP_ALIGN.LEFT)
        y += 1.5

    # 右側：市場份額
    add_text_box(slide, 17, 2.5, 15, 1, "市場份額 | Market Share", 16, True, Colors.SECONDARY)
    y = 4
    for i, seg in enumerate(data['segments'][:5], 1):
        bar_width = seg['share_2025_2030'] / 70 * 12
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(17), Cm(y), Cm(bar_width), Cm(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = Colors.PRIMARY
        shape.line.fill.background()

        add_text_box(slide, 17.2, y + 0.3, 12, 0.8, f"{i}. {seg['name']}", 12, True, Colors.WHITE)
        add_text_box(slide, 29.5, y + 0.3, 3, 0.8, f"{seg['share_2025_2030']}%", 12, True, Colors.TEXT_DARK, PP_ALIGN.LEFT)
        y += 1.5

    return slide


def add_trend_chart(prs, data, page_num, total):
    """年度趨勢頁"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 頁面標題
    add_text_box(slide, 1, 0.5, 31, 1.5, "年度趨勢 | Annual Trend (億元)", 26, True, Colors.PRIMARY)
    add_page_number(slide, page_num, total)

    # 分隔線
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1), Cm(2), Cm(31.8), Cm(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.ACCENT
    shape.line.fill.background()

    # 圖例
    add_text_box(slide, 1, 2.3, 31, 0.8, "電動汽車 vs 電動貨卡車 市場規模趨勢", 14, True, Colors.TEXT_DARK)

    # 選擇兩個主要賽道顯示
    seg1 = data['segments'][0]  # 電動汽車
    seg2 = data['segments'][1]  # 電動貨卡車

    # 繪製簡化的柱狀圖
    years = ['2025', '2026', '2027', '2028', '2029', '2030']
    x_start = 3
    bar_width = 1.5
    gap = 0.5
    y_scale = 0.003  # 調整係數

    # 電動汽車
    max_val = seg1['scale_2030']
    y = 14
    for i, year in enumerate(years):
        val = seg1['yearly_data'][year]
        height = val * y_scale
        x = x_start + i * (bar_width * 2 + gap)

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y - height), Cm(bar_width), Cm(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = Colors.PRIMARY
        shape.line.fill.background()

    # 電動貨卡車
    max_val2 = seg2['scale_2030']
    for i, year in enumerate(years):
        val = seg2['yearly_data'][year]
        height = val * y_scale * 0.6
        x = x_start + i * (bar_width * 2 + gap) + bar_width

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y - height), Cm(bar_width), Cm(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = Colors.ACCENT
        shape.line.fill.background()

    # 年份標籤
    for i, year in enumerate(years):
        x = x_start + i * (bar_width * 2 + gap) + bar_width / 2
        add_text_box(slide, x, 14.5, 3, 0.8, year, 11, False, Colors.TEXT_DARK, PP_ALIGN.CENTER)

    # 圖例
    add_text_box(slide, 1, 16, 4, 0.8, "■ 電動汽車", 11, False, Colors.PRIMARY)
    add_text_box(slide, 5, 16, 5, 0.8, "■ 電動貨卡車", 11, False, Colors.ACCENT)

    return slide


def add_conclusion(prs, data, page_num, total):
    """結論頁"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 頁面標題
    add_text_box(slide, 1, 0.5, 31, 1.5, "結論與展望 | Conclusion", 28, True, Colors.PRIMARY)
    add_page_number(slide, page_num, total)

    # 分隔線
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1), Cm(2), Cm(31.8), Cm(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.ACCENT
    shape.line.fill.background()

    top1 = data['segments'][0]
    top_ai = data['fastest_growing'][0]
    top_robot = data['fastest_growing'][1]
    top_drone = data['fastest_growing'][2]

    conclusions = f"""
【市場洞察】

1. 電動汽車主導市場
   • 5年累計規模：{top1['cumulative_5yr']:,.0f} 億元，市場份額 {top1['share_2025_2030']}%
   • 從 {top1['scale_2025']:,.0f} 億元（2025）增長至 {top1['scale_2030']:,.0f} 億元（2030）
   • CAGR {top1['cagr_2025_2030']}%，穩健成長

2. 新興賽道爆發式增長
   • AI儲能：CAGR {top_ai['cagr_2025_2030']}%，從 {top_ai['scale_2025']} 增至 {top_ai['scale_2030']} 億元
   • 機器人儲能：CAGR {top_robot['cagr_2025_2030']}%，從 {top_robot['scale_2025']} 增至 {top_robot['scale_2030']} 億元
   • 無人機儲能：CAGR {top_drone['cagr_2025_2030']}%，從 {top_drone['scale_2025']} 增至 {top_drone['scale_2030']} 億元

3. 市場趨勢
   • 電動化、智能化加速推進
   • 儲能需求爆發式增長
   • 新技術、新應用場景不斷涌現
"""
    add_text_box(slide, 1, 2.5, 31, 14, conclusions.strip(), 14, False, Colors.TEXT_DARK)

    return slide


def add_thank_you(prs, page_num, total):
    """感謝頁"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景裝飾
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(33.867), Cm(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.PRIMARY
    shape.line.fill.background()

    shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(18.75), Cm(33.867), Cm(0.3))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = Colors.PRIMARY
    shape2.line.fill.background()

    # 感謝文字
    add_text_box(slide, 1, 7, 31, 3, "感謝聆聽", 48, True, Colors.PRIMARY, PP_ALIGN.CENTER)
    add_text_box(slide, 1, 11, 31, 2, "Thank You", 24, False, Colors.TEXT_LIGHT, PP_ALIGN.CENTER)

    add_page_number(slide, page_num, total)

    return slide


# ============================================================================
# 主函數
# ============================================================================
def generate_market_report_ppt(data, output_path):
    """生成市場分析報告 PPT"""
    prs = create_presentation()
    total_slides = 7

    # 頁面 1: 封面
    add_cover_slide(
        prs,
        "全球儲能及新能源關聯市場\n經濟規模分析報告",
        f"2025-2030 | {len(data['segments'])}大賽道 | 市場規模分析"
    )

    # 頁面 2: 執行摘要
    add_executive_summary(prs, data, 2, total_slides)

    # 頁面 3: 市場總覽表格
    add_market_overview_table(prs, data, 3, total_slides)

    # 頁面 4: 增速排名
    add_growth_ranking(prs, data, 4, total_slides)

    # 頁面 5: 年度趨勢
    add_trend_chart(prs, data, 5, total_slides)

    # 頁面 6: 結論
    add_conclusion(prs, data, 6, total_slides)

    # 頁面 7: 感謝頁
    add_thank_you(prs, 7, total_slides)

    # 保存
    os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
    prs.save(output_path)
    print(f"[OK] PPT generated: {output_path}")
    return output_path


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python ppt_generator_v2.py <market_json> [output_pptx]")
        sys.exit(1)

    json_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else "market_report_v2.pptx"

    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    generate_market_report_ppt(data, output_path)
