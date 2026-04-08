# -*- coding: utf-8 -*-
"""
Power IC + Big Data + PPT 綜合報告生成器
整合電源IC技術分析與大數據報告
"""

import sys
import os
from pathlib import Path
from datetime import datetime

# Add skills directory to path
sys.path.insert(0, str(Path(__file__).parent))

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib import colors
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False


class PowerICReportGenerator:
    """Power IC 綜合報告生成器"""

    def __init__(self):
        self.power_ic_data = self._load_power_ic_samples()
        self.market_data = self._load_market_samples()

    def _load_power_ic_samples(self):
        """載入Power IC樣本數據"""
        return {
            "topology": {
                "Buck": {"market_share": 45, "growth": 8.5, "applications": ["消費電子", "工業控制", "汽車電子"]},
                "Boost": {"market_share": 20, "growth": 6.2, "applications": ["穿戴裝置", "IoT", "醫療設備"]},
                "LDO": {"market_share": 18, "growth": 4.1, "applications": ["模擬IC", "射頻", "感測器供電"]},
                "Buck-Boost": {"market_share": 10, "growth": 12.3, "applications": ["TWS耳機", "手持設備", "USB PD"]},
                "Motor Driver": {"market_share": 7, "growth": 9.8, "applications": ["風扇控制", "馬達驅動", "機器人"]},
            },
            "vendors": {
                "TI": {"market_share": 28, "focus": "全系列", "strength": "品質穩定"},
                "MPS": {"market_share": 15, "focus": "高效能", "strength": "高效率"},
                "DIODES": {"market_share": 12, "focus": "成本優化", "strength": "性價比"},
                "RICHTEK": {"market_share": 18, "focus": "亞太市場", "strength": "本地支援"},
                "SILERGY": {"market_share": 10, "focus": "中小功率", "strength": "快速交付"},
                "Others": {"market_share": 17, "focus": "細分市場", "strength": "差異化"},
            },
            "technologies": {
                "GaN": {"adoption": 15, "efficiency": 98, "cost_factor": 2.5, "trend": "快速增長"},
                "SiC": {"adoption": 8, "efficiency": 97, "cost_factor": 3.0, "trend": "高端應用"},
                "Advanced Si": {"adoption": 77, "efficiency": 95, "cost_factor": 1.0, "trend": "穩定"},
            }
        }

    def _load_market_samples(self):
        """載入市場樣本數據"""
        years = ["2020", "2021", "2022", "2023", "2024", "2025E", "2026E", "2027E"]
        return {
            "market_size": {
                "years": years,
                "values": [28.5, 32.1, 36.8, 41.2, 46.8, 52.5, 58.2, 65.0]  # 億美元
            },
            "unit_shipments": {
                "years": years,
                "values": [85, 98, 115, 132, 150, 172, 195, 220]  # 億顆
            },
            "avg_price": {
                "years": years,
                "values": [0.335, 0.328, 0.320, 0.312, 0.312, 0.305, 0.298, 0.295]  # 美元
            }
        }

    def generate_market_analysis(self) -> str:
        """生成市場分析報告"""
        report = []
        report.append("=" * 70)
        report.append("       IT Power IC Market Analysis Report")
        report.append("=" * 70)
        report.append(f"Report Date: {datetime.now().strftime('%Y-%m-%d')}")
        report.append("")

        # Market Overview
        report.append("[1. Market Overview]")
        data = self.market_data
        latest = data["market_size"]["values"][-3]  # 2024 actual
        forecast = data["market_size"]["values"][-1]  # 2027E
        cagr = ((forecast / latest) ** (1/3) - 1) * 100
        report.append(f"  - 2024 Market Size: ${latest}B USD")
        report.append(f"  - 2027E Market Size: ${forecast}B USD")
        report.append(f"  - CAGR (2024-2027): {cagr:.1f}%")
        report.append(f"  - 2024 Unit Shipments: {data['unit_shipments']['values'][-3]}B units")

        # Topology Analysis
        report.append("\n[2. Topology Market Share]")
        for topo, info in sorted(self.power_ic_data["topology"].items(), key=lambda x: x[1]["market_share"], reverse=True):
            report.append(f"  {topo:15} {info['market_share']:5}%  Growth: {info['growth']}%")

        # Vendor Analysis
        report.append("\n[3. Key Vendors]")
        for vendor, info in sorted(self.power_ic_data["vendors"].items(), key=lambda x: x[1]["market_share"], reverse=True):
            report.append(f"  {vendor:10} Share: {info['market_share']:5}%  Focus: {info['focus']}")

        # Technology Trends
        report.append("\n[4. Technology Trends]")
        for tech, info in self.power_ic_data["technologies"].items():
            report.append(f"  {tech}: Adoption {info['adoption']}% | Efficiency {info['efficiency']}% | Cost: {info['cost_factor']}x | Trend: {info['trend']}")

        # Key Insights
        report.append("\n[5. Key Insights]")
        report.append("  1. GaN technology adoption accelerating in consumer electronics")
        report.append("  2. USB PD driving Buck-Boost demand growth")
        report.append("  3. Asian vendors gaining market share through cost optimization")
        report.append("  4. Automotive and Industrial segments show resilience")
        report.append("  5. Miniaturization trend continues with advanced packaging")

        report.append("\n" + "=" * 70)
        return "\n".join(report)

    def generate_ic_selection_report(self) -> str:
        """生成IC選型報告"""
        report = []
        report.append("=" * 70)
        report.append("       Power IC Selection Guide")
        report.append("=" * 70)
        report.append("")

        # Common Applications
        apps = {
            "TWS Earphones": {
                "requirements": "Vin: 3.0-5.0V, Vout: 1.8-3.3V, Iout: 0.5-2A",
                "topology": "Buck / LDO",
                "key_ics": ["TPS62840", "SY6288", "RT6150"]
            },
            "Smart Watch": {
                "requirements": "Vin: 3.7-4.2V, Vout: 1.2-3.3V, Iout: 0.2-1A",
                "topology": "Buck (High Efficiency)",
                "key_ics": ["TPS62840", "AXP2101", "SY6288"]
            },
            "Industrial IoT": {
                "requirements": "Vin: 5-24V, Vout: 3.3-12V, Iout: 1-5A",
                "topology": "Buck",
                "key_ics": ["MP1484", "RT8289", "LM2596"]
            },
            "Automotive": {
                "requirements": "Vin: 9-16V, Vout: 3.3-5V, Iout: 2-10A",
                "topology": "Buck / LDO",
                "key_ics": ["TPS54540", "LM2596", "MPQ4470"]
            }
        }

        for app, details in apps.items():
            report.append(f"\n[{app}]")
            report.append(f"  Requirements: {details['requirements']}")
            report.append(f"  Topology: {details['topology']}")
            report.append(f"  Recommended ICs: {', '.join(details['key_ics'])}")

        report.append("\n" + "=" * 70)
        return "\n".join(report)

    def create_ppt_report(self, output_path: str = "PowerIC_Report.pptx") -> bool:
        """生成PPT報告"""
        if not HAS_PPTX:
            print("[ERROR] python-pptx not installed")
            return False

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Slide 1: Title
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "IT Power IC Market Analysis"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)

        subtitle = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
        tf2 = subtitle.text_frame
        tf2.paragraphs[0].text = f"Report Date: {datetime.now().strftime('%Y-%m-%d')}"
        tf2.paragraphs[0].font.size = Pt(24)

        # Slide 2: Market Overview
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title.text_frame
        tf.paragraphs[0].text = "Market Overview"
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True

        content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
        tf = content.text_frame
        tf.word_wrap = True

        lines = [
            "2024 Market Size: $46.8B USD",
            "2027E Market Size: $65.0B USD",
            "CAGR (2024-2027): 11.5%",
            "",
            "Key Growth Drivers:",
            "  - GaN adoption in fast charging",
            "  - USB PD protocol expansion",
            "  - EV and automotive electrification",
            "  - IoT device proliferation",
            "  - AI server power requirements"
        ]

        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(20)

        # Slide 3: Topology Analysis
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title.text_frame
        tf.paragraphs[0].text = "Topology Market Share"
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True

        # Add table
        rows = [["Topology", "Market Share", "Growth", "Key Applications"]]
        for topo, info in sorted(self.power_ic_data["topology"].items(), key=lambda x: x[1]["market_share"], reverse=True):
            rows.append([topo, f"{info['market_share']}%", f"{info['growth']}%", ", ".join(info['applications'][:2])])

        table = slide.shapes.add_table(len(rows), 4, Inches(0.5), Inches(1.5), Inches(12), Inches(4)).table
        for i, row in enumerate(rows):
            for j, cell in enumerate(row):
                cell_obj = table.cell(i, j)
                cell_obj.text = cell
                cell_obj.text_frame.paragraphs[0].font.size = Pt(14)
                if i == 0:
                    cell_obj.text_frame.paragraphs[0].font.bold = True

        # Slide 4: Technology Trends
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title.text_frame
        tf.paragraphs[0].text = "Technology Trends"
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True

        content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        tf = content.text_frame
        tf.word_wrap = True

        lines = [
            "GaN (Gallium Nitride):",
            "  - Market adoption: 15% | Efficiency: 98%",
            "  - Trend: Rapid growth in consumer fast chargers",
            "",
            "SiC (Silicon Carbide):",
            "  - Market adoption: 8% | Efficiency: 97%",
            "  - Trend: High-power applications, automotive",
            "",
            "Advanced Si (Traditional Silicon):",
            "  - Market adoption: 77% | Efficiency: 95%",
            "  - Trend: Stable, cost-optimized solutions"
        ]

        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            if line.endswith(":") and not line.startswith(" "):
                p.font.bold = True

        # Save
        prs.save(output_path)
        print(f"[OK] PPT saved: {output_path}")
        return True

    def create_pdf_report(self, output_path: str = "PowerIC_Report.pdf") -> bool:
        """生成PDF報告"""
        if not HAS_REPORTLAB:
            print("[ERROR] reportlab not installed")
            return False

        doc = SimpleDocTemplate(output_path, pagesize=A4)
        story = []
        styles = getSampleStyleSheet()

        # Title
        story.append(Paragraph("IT Power IC Market Analysis Report", styles['Title']))
        story.append(Spacer(1, 20))
        story.append(Paragraph(f"Date: {datetime.now().strftime('%Y-%m-%d')}", styles['Normal']))
        story.append(Spacer(1, 30))

        # Market Overview
        story.append(Paragraph("1. Market Overview", styles['Heading1']))
        market_text = f"""
        2024 Market Size: $46.8B USD<br/>
        2027E Market Size: $65.0B USD<br/>
        CAGR (2024-2027): 11.5%<br/>
        2024 Unit Shipments: 150B units
        """
        story.append(Paragraph(market_text, styles['Normal']))
        story.append(Spacer(1, 20))

        # Topology Table
        story.append(Paragraph("2. Topology Market Share", styles['Heading1']))
        data = [["Topology", "Share", "Growth", "Key Apps"]]
        for topo, info in self.power_ic_data["topology"].items():
            data.append([topo, f"{info['market_share']}%", f"{info['growth']}%", ", ".join(info['applications'][:2])])

        t = Table(data)
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(t)
        story.append(Spacer(1, 20))

        # IC Selection Guide
        story.append(Paragraph("3. IC Selection Guide", styles['Heading1']))
        selection_text = """
        TWS Earphones: TPS62840, SY6288 (Buck/LDO)<br/>
        Smart Watch: TPS62840, AXP2101 (High Efficiency Buck)<br/>
        Industrial IoT: MP1484, RT8289 (General Buck)<br/>
        Automotive: TPS54540, MPQ4470 (Wide Vin Buck)
        """
        story.append(Paragraph(selection_text, styles['Normal']))

        doc.build(story)
        print(f"[OK] PDF saved: {output_path}")
        return True


def main():
    generator = PowerICReportGenerator()

    import sys

    if len(sys.argv) < 2:
        # Demo
        print("\n[Power IC + Big Data + PPT Report Generator]")
        print("\n--- Market Analysis ---")
        print(generator.generate_market_analysis())

        print("\n--- IC Selection Guide ---")
        print(generator.generate_ic_selection_report())

        # Try to generate files
        print("\n--- Generating Reports ---")
        generator.create_ppt_report("PowerIC_Report.pptx")
        generator.create_pdf_report("PowerIC_Report.pdf")

    else:
        cmd = sys.argv[1]

        if cmd == "market":
            print(generator.generate_market_analysis())

        elif cmd == "selection":
            print(generator.generate_ic_selection_report())

        elif cmd == "ppt":
            output = sys.argv[2] if len(sys.argv) > 2 else "PowerIC_Report.pptx"
            generator.create_ppt_report(output)

        elif cmd == "pdf":
            output = sys.argv[2] if len(sys.argv) > 2 else "PowerIC_Report.pdf"
            generator.create_pdf_report(output)

        elif cmd == "all":
            generator.create_ppt_report("PowerIC_Report.pptx")
            generator.create_pdf_report("PowerIC_Report.pdf")

        else:
            print("Commands: market | selection | ppt [output] | pdf [output] | all")


if __name__ == "__main__":
    main()
