# -*- coding: utf-8 -*-
"""
PPT 內容提取器 v2 - 通用版
提取任意 PPT 的內容生成摘要報告
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import os
from datetime import datetime


def extract_presentation_content(file_path):
    """提取 PPT 完整內容"""
    prs = Presentation(file_path)

    result = {
        "file": file_path,
        "file_name": os.path.basename(file_path),
        "analyzed_at": datetime.now().isoformat(),
        "total_slides": len(prs.slides),
        "statistics": {
            "total_shapes": 0,
            "total_tables": 0,
            "total_images": 0,
            "total_charts": 0,
            "total_text_boxes": 0
        },
        "slides": [],
        "summary": {
            "title": "",
            "table_of_contents": [],
            "key_points": [],
            "full_text": ""
        }
    }

    all_text = []

    for idx, slide in enumerate(prs.slides, 1):
        slide_data = {
            "slide_number": idx,
            "title": "",
            "content": [],
            "images_count": 0,
            "tables_count": 0
        }

        for shape in slide.shapes:
            result["statistics"]["total_shapes"] += 1

            # 提取文字
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    all_text.append(text)

                    if not slide_data["title"]:
                        slide_data["title"] = text[:100]

                    # 過濾有意義的文字
                    if len(text) > 10 and not text.startswith("Picture"):
                        slide_data["content"].append(text)

                    result["statistics"]["total_text_boxes"] += 1

                    # 第一頁設為標題
                    if idx == 1 and not result["summary"]["title"]:
                        result["summary"]["title"] = text

            # 提取表格
            if shape.has_table:
                result["statistics"]["total_tables"] += 1
                slide_data["tables_count"] += 1

            # 提取圖片
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                result["statistics"]["total_images"] += 1
                slide_data["images_count"] += 1

            # 提取圖表
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                result["statistics"]["total_charts"] += 1

        result["slides"].append(slide_data)

    # 生成摘要
    result["summary"]["full_text"] = "\n\n".join(all_text)

    # 提取目錄（從第二頁）
    for slide in result["slides"][1:5]:
        if slide["title"] and "目錄" in slide["title"]:
            result["summary"]["table_of_contents"] = slide["content"]
            break

    # 提取關鍵點
    keywords = ["核心", "優勢", "策略", "機遇", "目標", "價值"]
    for text in all_text:
        for kw in keywords:
            if kw in text and text not in result["summary"]["key_points"]:
                result["summary"]["key_points"].append(text[:200])
                break

    return result


def generate_content_summary(data):
    """生成內容摘要"""
    summary = []
    summary.append("=" * 60)
    summary.append("PPT Content Summary")
    summary.append("=" * 60)
    summary.append(f"\nTitle: {data['summary']['title'] or data['file_name']}")
    summary.append(f"\nSlides: {data['total_slides']}")
    summary.append(f"Images: {data['statistics']['total_images']}")
    summary.append(f"Tables: {data['statistics']['total_tables']}")
    summary.append(f"Charts: {data['statistics']['total_charts']}")

    if data['summary']['table_of_contents']:
        summary.append("\nTable of Contents:")
        for item in data['summary']['table_of_contents'][:6]:
            summary.append(f"  - {item[:60]}")

    if data['summary']['key_points']:
        summary.append("\nKey Points:")
        for point in data['summary']['key_points'][:5]:
            summary.append(f"  - {point[:80]}...")

    return "\n".join(summary)


def save_json(data, output_path):
    """保存為 JSON"""
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"[OK] JSON: {output_path}")


def save_text(data, output_path):
    """保存為文字報告"""
    lines = []
    lines.append("=" * 60)
    lines.append("PPT 內容報告")
    lines.append("=" * 60)
    lines.append(f"\n檔案: {data['file_name']}")
    lines.append(f"時間: {data['analyzed_at']}")
    lines.append(f"頁數: {data['total_slides']}")
    lines.append(f"圖片: {data['statistics']['total_images']}")
    lines.append(f"表格: {data['statistics']['total_tables']}")
    lines.append(f"圖表: {data['statistics']['total_charts']}")

    lines.append("\n" + "=" * 60)
    lines.append("各頁內容")
    lines.append("=" * 60)

    for slide in data['slides']:
        lines.append(f"\n[第 {slide['slide_number']} 頁] {slide['title'] or '(無標題)'}")
        if slide['images_count']:
            lines.append(f"  含 {slide['images_count']} 張圖片")
        if slide['tables_count']:
            lines.append(f"  含 {slide['tables_count']} 個表格")
        for content in slide['content'][:3]:
            lines.append(f"  • {content[:100]}")

    lines.append("\n" + "=" * 60)
    lines.append("完整文字內容")
    lines.append("=" * 60)
    lines.append(data['summary']['full_text'])

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write("\n".join(lines))
    print(f"[OK] TXT: {output_path}")


def print_summary(data):
    """打印摘要"""
    print(generate_content_summary(data))


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("用法: python ppt_extractor.py <pptx_file> [output_dir]")
        sys.exit(1)

    file_path = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "."

    if not os.path.exists(file_path):
        print(f"[錯誤] 檔案不存在: {file_path}")
        sys.exit(1)

    print(f"\n正在分析: {file_path}")
    data = extract_presentation_content(file_path)
    print_summary(data)

    # 保存
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    save_json(data, os.path.join(output_dir, f"{base_name}_content.json"))
    save_text(data, os.path.join(output_dir, f"{base_name}_content.txt"))

    print("\n[完成]")
