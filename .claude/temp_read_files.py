# -*- coding: utf-8 -*-
import pandas as pd
import json
import sys

# Read CoBM Excel
try:
    cobm_file = r'D:/深圳/台灣瑞聲源/副本CoBM产品资料.xlsx'
    xl = pd.ExcelFile(cobm_file)
    print('=== CoBM Sheets ===')
    print(xl.sheet_names)
    
    cobm_data = {}
    for sheet in xl.sheet_names:
        df = pd.read_excel(cobm_file, sheet_name=sheet)
        cobm_data[sheet] = df.to_dict()
        print(f'\n=== {sheet} ===')
        print(f'Shape: {df.shape}')
        print('Columns:', list(df.columns))
        print(df.to_string())
except Exception as e:
    print(f'CoBM Error: {e}', file=sys.stderr)
    import traceback
    traceback.print_exc()

print('\n\n=== BQT PPT Content ===')
try:
    from pptx import Presentation
    bqt_file = r'D:/深圳/台灣瑞聲源/BQT高瓦数电源双技术深度分析报告（专家版）.pptx'
    prs = Presentation(bqt_file)
    
    bqt_data = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
        if slide_texts:
            bqt_data.append({'slide': i+1, 'text': '\n'.join(slide_texts)})
            print(f'\n--- Slide {i+1} ---')
            print('\n'.join(slide_texts))
except Exception as e:
    print(f'BQT Error: {e}', file=sys.stderr)
    import traceback
    traceback.print_exc()

# Save extracted data
try:
    with open(r'D:/claude mini max 2.7/.claude/extracted_data.json', 'w', encoding='utf-8') as f:
        json.dump({
            'cobm_data': cobm_data if 'cobm_data' in dir() else {},
            'bqt_data': bqt_data if 'bqt_data' in dir() else []
        }, f, ensure_ascii=False, default=str)
    print('\n\nData saved to extracted_data.json')
except Exception as e:
    print(f'Save Error: {e}', file=sys.stderr)
