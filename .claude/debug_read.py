# -*- coding: utf-8 -*-
import pandas as pd
import json
import traceback

output = []

# Read CoBM Excel
try:
    cobm_file = r'D:/深圳/台灣瑞聲源/副本CoBM产品资料.xlsx'
    output.append('Reading CoBM file...')
    xl = pd.ExcelFile(cobm_file)
    output.append(f'Sheets: {xl.sheet_names}')

    cobm_data = {}
    for sheet in xl.sheet_names:
        df = pd.read_excel(cobm_file, sheet_name=sheet)
        cobm_data[sheet] = df.to_dict()
        output.append(f'\n=== {sheet} ===')
        output.append(f'Shape: {df.shape}')
        output.append(f'Columns: {list(df.columns)}')
        output.append(df.to_string())
except Exception as e:
    output.append(f'CoBM Error: {e}')
    output.append(traceback.format_exc())

output.append('\n\n=== BQT PPT Content ===')
try:
    from pptx import Presentation
    bqt_file = r'D:/深圳/台灣瑞聲源/BQT高瓦数电源双技术深度分析报告（专家版）.pptx'
    output.append('Reading BQT file...')
    prs = Presentation(bqt_file)
    output.append(f'Number of slides: {len(prs.slides)}')

    bqt_data = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
        if slide_texts:
            bqt_data.append({'slide': i+1, 'text': '\n'.join(slide_texts)})
            output.append(f'\n--- Slide {i+1} ---')
            output.append('\n'.join(slide_texts))
except Exception as e:
    output.append(f'BQT Error: {e}')
    output.append(traceback.format_exc())

# Save extracted data
try:
    with open(r'D:/claude mini max 2.7/.claude/extracted_data.json', 'w', encoding='utf-8') as f:
        json.dump({
            'cobm_data': cobm_data if 'cobm_data' in dir() else {},
            'bqt_data': bqt_data if 'bqt_data' in dir() else []
        }, f, ensure_ascii=False, default=str)
    output.append('\n\nData saved to extracted_data.json')
except Exception as e:
    output.append(f'Save Error: {e}')

# Write output to file
with open(r'D:/claude mini max 2.7/.claude/read_output.txt', 'w', encoding='utf-8') as f:
    f.write('\n'.join(output))
