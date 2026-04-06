# -*- coding: utf-8 -*-
"""
CoBM电源IC方案专业技术分析报告 PPT生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 82, 118)  # 深蓝色
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    # 日期
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026年4月5日"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets, highlight_color=None):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8)
    )
    header.fill.solid()
    if highlight_color:
        header.fill.fore_color.rgb = RGBColor(*highlight_color)
    else:
        header.fill.fore_color.rgb = RGBColor(26, 82, 118)
    header.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(12)
        p.space_after = Pt(6)
        # 添加项目符号
        p.text = "● " + bullet

    return slide

def add_table_slide(prs, title, headers, rows, highlight_color=None):
    """添加表格页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8)
    )
    header.fill.solid()
    if highlight_color:
        header.fill.fore_color.rgb = RGBColor(*highlight_color)
    else:
        header.fill.fore_color.rgb = RGBColor(26, 82, 118)
    header.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 表格
    num_cols = len(headers)
    num_rows = len(rows) + 1
    table = slide.shapes.add_table(
        num_rows, num_cols, Inches(0.5), Inches(1.2), Inches(9), Inches(4.5)
    ).table

    # 设置列宽
    col_width = Inches(9.0 / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_width

    # 表头
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(26, 82, 118)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # 数据行
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            # 交替颜色
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.alignment = PP_ALIGN.CENTER

def add_comparison_slide(prs, title, left_title, left_items, right_title, right_items):
    """添加对比页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(26, 82, 118)
    header.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 左侧标题
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4.3), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 82, 118)

    # 左侧内容
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4.3), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "✓ " + item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(8)

    # 中间分隔线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.9), Inches(1.1), Inches(0.02), Inches(5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(200, 200, 200)
    line.line.fill.background()

    # 右侧标题
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(39, 174, 96)

    # 右侧内容
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(4.3), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "✓ " + item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(8)

def add_conclusion_slide(prs, title, items, recommendations):
    """添加结论页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(39, 174, 96)
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 结论内容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(3))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "◆ " + item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(10)

    # 建议框
    rec_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.5), Inches(8.6), Inches(1.8)
    )
    rec_box.fill.solid()
    rec_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rec_box.line.fill.background()

    # 建议内容
    rec_text = slide.shapes.add_textbox(Inches(0.9), Inches(4.7), Inches(8.2), Inches(1.4))
    tf = rec_text.text_frame
    tf.word_wrap = True
    for i, rec in enumerate(recommendations):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = "投资建议"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(39, 174, 96)
            p = tf.add_paragraph()
        else:
            p = tf.add_paragraph()
        p.text = rec
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(4)

def create_presentation():
    """创建PPT演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 第1页：标题页
    add_title_slide(
        prs,
        "CoBM电源IC方案专业技术分析",
        "GaN · SiC · Si · 数字控制 · 封装技术"
    )

    # 第2页：目录
    add_content_slide(
        prs,
        "报告目录",
        [
            "一、电源IC技术分类体系",
            "二、GaN/SiC/Si IC方案详细分析",
            "三、数字控制IC方案分析",
            "四、PFC与LLC控制器IC",
            "五、CoBM封装技术分析",
            "六、IC方案选型指南",
            "七、综合对比与供应商生态",
            "八、技术结论与总结"
        ],
        highlight_color=(142, 68, 173)
    )

    # 第3页：半导体材料对比
    headers = ["材料", "带隙(eV)", "耐压", "效率", "成本($/pcs)", "市场地位"]
    rows = [
        ["Si (硅)", "1.1", "650V级", "92-95%", "$0.1-0.5", "主流 70%"],
        ["GaN (氮化镓)", "3.4", "650/800V", "96-99%", "$0.6-2.5", "快速增长"],
        ["SiC (碳化硅)", "3.3", "650-1700V", "97-99%", "$1.5-8.0", "高端市场"]
    ]
    add_table_slide(prs, "一、半导体材料技术对比", headers, rows, highlight_color=(26, 82, 118))

    # 第4页：GaN厂商对比
    headers = ["厂商", "产品系列", "耐压", "峰值效率", "主要应用", "价格($/pcs)"]
    rows = [
        ["Infineon", "CoolGaN G5", "600V", ">98%", "服务器/汽车", "$1.2-2.5"],
        ["Navitas", "GaNFast Gen2", "650V", "99%", "USB-PD/消费", "$0.8-1.5"],
        ["Power Integrations", "InnoSwitch5-Pro", "650V", "97%", "AC-DC/充电", "$0.6-1.2"],
        ["EPC", "eGaN FET", "100-350V", "98%", "激光雷达/48V", "$0.5-2.0"]
    ]
    add_table_slide(prs, "二、GaN厂商产品线对比", headers, rows, highlight_color=(39, 174, 96))

    # 第5页：SiC厂商对比
    headers = ["厂商", "产品系列", "耐压", "峰值效率", "主要应用", "价格($/pcs)"]
    rows = [
        ["Wolfspeed", "SiC FET", "650-1700V", "98%+", "数据中心/汽车", "$2.0-8.0"],
        ["Infineon", "CoolSiC G3", "650-1700V", "98%+", "工业/汽车", "$1.8-6.0"],
        ["onsemi", "EliteSiC", "650-1700V", "97%+", "服务器/汽车", "$1.5-5.0"],
        ["ROHM", "SiC MOSFET", "650-1700V", "98%+", "工业/汽车", "$1.5-5.5"]
    ]
    add_table_slide(prs, "二、SiC厂商产品线对比", headers, rows, highlight_color=(231, 76, 60))

    # 第6页：数字控制器对比
    headers = ["型号", "厂商", "主频", "PWM分辨率", "PMBus", "主要应用"]
    rows = [
        ["TMS320F28379D", "TI", "200MHz", "150ps", "v1.3", "旗舰服务器"],
        ["UCD3138", "TI", "125MHz", "250ps", "v1.3", "通信电源"],
        ["XMC4500", "Infineon", "144MHz", "100ps", "可选", "工业"],
        ["NCP4208", "ON Semi", "数字核", "数字", "v1.3", "服务器"],
        ["dsPIC33EP", "Microchip", "70MHz", "250ps", "可选", "消费级"]
    ]
    add_table_slide(prs, "三、数字控制器对比", headers, rows, highlight_color=(52, 152, 219))

    # 第7页：PFC控制器对比
    headers = ["厂商", "型号", "拓扑", "开关频率", "GaN兼容", "应用"]
    rows = [
        ["Infineon", "EiceDRIVER 2ED4820", "CCM/CRM", "最高500kHz", "是", "服务器"],
        ["TI", "UCC28180", "CCM", "最高250kHz", "否", "通用"],
        ["onsemi", "NCP1654", "CCM", "最高250kHz", "否", "消费级"],
        ["MPS", "HR1201", "CRM/CCM", "最高300kHz", "是", "紧凑型"],
        ["Infineon", "EiceDRIVER 1ED-TC", "Totem-Pole", "高频率", "是", "高端服务器"]
    ]
    add_table_slide(prs, "四、PFC控制器IC对比", headers, rows, highlight_color=(155, 89, 182))

    # 第8页：LLC控制器对比
    headers = ["厂商", "型号", "谐振类型", "开关频率", "同步整流", "应用"]
    rows = [
        ["TI", "UCC25640", "半桥LLC", "50-500kHz", "是", "服务器"],
        ["TI", "LMG3410", "GaN集成", "最高1MHz", "是", "高密度"],
        ["Infineon", "XMC1300", "半桥LLC", "50-500kHz", "是", "工业"],
        ["onsemi", "NCP1399", "混合LLC", "最高500kHz", "是", "消费级"],
        ["ST", "L6699", "半桥LLC", "50-400kHz", "是", "通用"]
    ]
    add_table_slide(prs, "四、LLC谐振控制器IC对比", headers, rows, highlight_color=(241, 196, 15))

    # 第9页：CoBM封装对比
    headers = ["封装类型", "特点", "功率范围", "主要优势", "主要劣势"]
    rows = [
        ["CoBM", "多芯片同封", "200W-3kW", "最高集成度", "定制成本高"],
        ["SiP", "系统级封装", "10-200W", "较高集成", "散热受限"],
        ["分立", "独立器件", "全范围", "灵活/成本低", "集成度低"],
        ["模组", "标准化模块", "500W-5kW", "标准化/可靠", "定制性差"]
    ]
    add_table_slide(prs, "五、CoBM封装技术对比", headers, rows, highlight_color=(26, 82, 118))

    # 第10页：功率段选型
    headers = ["功率段", "推荐拓扑", "GaN方案", "Si方案", "效率目标"]
    rows = [
        ["30-65W", "QR Flyback", "Navitas NV6172", "InnoSwitch4", "92-94%"],
        ["65-120W", "QR/LLC Flyback", "PI InnoSwitch5", "NCP1399", "94-96%"],
        ["250-500W", "PFC+LLC", "CoolGaN", "NCP1654+LLC", "96-97%"],
        ["500-1000W", "PFC+LLC", "CoolGaN", "Si双级", "97-98%"],
        ["1000-3000W", "Totem PFC+LLC", "CoolGaN", "Si双级+SiC", "98%+"]
    ]
    add_table_slide(prs, "六、IC方案按功率段选型", headers, rows, highlight_color=(39, 174, 96))

    # 第11页：技术维度对比
    add_comparison_slide(
        prs,
        "七、技术维度对比",
        "传统分立Si方案",
        [
            "BOM成本: $50-55",
            "80+认证: 金牌",
            "功率密度: 18-22 W/in3",
            "纹波(12V): 25-35mV",
            "保持时间: >16ms",
            "MTBF: 100k小时",
            "边际毛利率: 18-22%"
        ],
        "CoBM(Si)方案",
        [
            "BOM成本: $38-42 (降低20-25%)",
            "80+认证: 白金牌 (升级2-4%)",
            "功率密度: 28-32 W/in3 (提升50%)",
            "纹波(12V): <15mV (降低50-60%)",
            "保持时间: >22ms (提升35%)",
            "MTBF: 120k+小时 (提升20%)",
            "边际毛利率: 28-32% (提升8-10pp)"
        ]
    )

    # 第12页：供应商生态
    add_content_slide(
        prs,
        "七、全链路IC方案供应商生态",
        [
            "【PFC控制器】Infineon, TI, onsemi, ST, MPS, Richtek, NXP",
            "【LLC控制器】TI, Infineon, onsemi, ST, MPS, Microchip",
            "【GaN器件】Infineon(CoolGaN), Navitas(GaNFast), Power Integrations, EPC",
            "【SiC器件】Wolfspeed, Infineon(CoolSiC), onsemi(EliteSiC), ROHM",
            "【Si MOSFET】Infineon(OptiMOS), onsemi(NexperSi), ST, Vishay, Toshiba",
            "【数字控制器】TI(C2000), Infineon(XMC), NXP, Microchip(dsPIC), onsemi"
        ],
        highlight_color=(26, 82, 118)
    )

    # 第13页：投资建议
    headers = ["阶段", "产品", "IC方案", "预期收益", "风险等级"]
    rows = [
        ["Phase 1 (2026-2027)", "RM550x/650x/750x", "CoBM(Si)", "ROI 126%", "低"],
        ["Phase 2 (2027-2028)", "SF750 SFX", "CoBM(GaN)", "体积-30%", "中"],
        ["Phase 3 (2028-2029)", "HX1200旗舰", "CoBM(SiC)", "效率99%+", "中"]
    ]
    add_table_slide(prs, "八、投资建议路线图", headers, rows, highlight_color=(39, 174, 96))

    # 第14页：结论
    add_conclusion_slide(
        prs,
        "结论与建议",
        [
            "CoBM(Si)方案性价比最优：750W电竞电源成本-25%, 效率+2%, 功率密度+50%",
            "GaN是未来方向：高功率密度场景，GaN CoBM可实现体积缩小30-40%",
            "SiC适合旗舰定位：1200W+产品，SiC CoBM可实现98%+效率",
            "数字控制是高端标配：PMBus监控/软件调优/固件升级已成标准"
        ],
        [
            "Phase 1: RM550x/650x/750x -> CoBM(Si) -> ROI 126%",
            "Phase 2: SF750 SFX -> CoBM(GaN) -> 体积-30%",
            "Phase 3: HX1200旗舰 -> CoBM(SiC) -> 效率99%+"
        ]
    )

    # 保存
    output_path = r'D:/claude mini max 2.7/.claude/CoBM_IC方案专业技术分析报告_20260405.pptx'
    prs.save(output_path)
    print(f"PPT已生成: {output_path}")
    return output_path

if __name__ == '__main__':
    create_presentation()
