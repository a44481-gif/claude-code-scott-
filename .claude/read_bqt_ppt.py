from pptx import Presentation

ppt_path = r'D:\深圳\台灣瑞聲源\BQT高瓦数电源双技术深度分析报告（专家版）.pptx'
prs = Presentation(ppt_path)

for i, slide in enumerate(prs.slides):
    print(f'\n=== Slide {i+1} ===')
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            print(shape.text)
