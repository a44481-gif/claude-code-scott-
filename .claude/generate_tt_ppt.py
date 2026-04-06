# -*- coding: utf-8 -*-
"""
曜越T.T电竞电源CoBM方案分析报告 PPT生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 82, 118)
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    # 日期
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026年4月6日"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets, highlight_color=None):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8)
    )
    header.fill.solid()
    if highlight_color:
        header.fill.fore_color.rgb = RGBColor(*highlight_color)
    else:
        header.fill.fore_color.rgb = RGBColor(26, 82, 118)
    header.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "● " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(10)
        p.space_after = Pt(4)

    return slide

def add_table_slide(prs, title, headers, rows, highlight_color=None):
    """添加表格页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8)
    )
    header_shape.fill.solid()
    if highlight_color:
        header_shape.fill.fore_color.rgb = RGBColor(*highlight_color)
    else:
        header_shape.fill.fore_color.rgb = RGBColor(26, 82, 118)
    header_shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 表格
    num_cols = len(headers)
    num_rows = len(rows) + 1

    # 计算表格高度
    table_height = min(Inches(5.2), Inches(0.5) * (num_rows + 1))

    table = slide.shapes.add_table(
        num_rows, num_cols, Inches(0.3), Inches(1.1), Inches(9.4), table_height
    ).table

    # 设置列宽
    col_width = Inches(9.4 / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_width

    # 表头
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(26, 82, 118)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # 数据行
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.alignment = PP_ALIGN.CENTER

def add_comparison_slide(prs, title, left_title, left_items, right_title, right_items):
    """添加对比页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(26, 82, 118)
    header.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 左侧标题
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4.3), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(231, 76, 60)

    # 左侧内容
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4.3), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "○ " + item
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.space_before = Pt(6)

    # 中间分隔线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.9), Inches(1.1), Inches(0.02), Inches(5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(200, 200, 200)
    line.line.fill.background()

    # 右侧标题
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(39, 174, 96)

    # 右侧内容
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(4.3), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "● " + item
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(6)

def add_conclusion_slide(prs, title, items):
    """添加结论页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(39, 174, 96)
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 结论内容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "◆ " + item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(14)

def create_presentation():
    """创建PPT演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 第1页：标题页
    add_title_slide(
        prs,
        "CoBM产品IC方案优势分析",
        "曜越T.T电竞电源市场切入策略 · 竞品分析 · 销售增长预测"
    )

    # 第2页：目录
    add_content_slide(
        prs,
        "报告目录",
        [
            "一、CoBM产品IC方案优势分析",
            "二、曜越T.T电竞电源产品线调研",
            "三、CoBM方案切入策略",
            "四、竞品分析",
            "五、竞争优势对比",
            "六、销售增长预测",
            "七、策略建议"
        ],
        highlight_color=(142, 68, 173)
    )

    # 第3页：CoBM技术优势
    add_content_slide(
        prs,
        "一、CoBM产品IC方案技术优势",
        [
            "【性能提升】功率密度+50%~65%，纹波降低50%~60%",
            "【能效升级】转换效率提升2%~4%，金牌可升级至白金牌",
            "【成本优化】BOM+制造综合成本降低18%~28%",
            "【可靠性增强】焊点失效率降低80%+，MTBF达120,000+小时",
            "【散热改善】热点温度一致性优化，风扇转速可降低10%~20%",
            "【集成度高】元件数量减少40%+，PCB面积节省30%~40%"
        ],
        highlight_color=(26, 82, 118)
    )

    # 第4页：曜越T.T产品线 - GF系列
    headers = ["系列", "瓦数覆盖", "80+认证", "模组化", "质保", "参考价(USD)"]
    rows = [
        ["GF3 (旗舰)", "650W-1650W", "Gold", "全模组", "10年", "$110-$440"],
        ["GF2 (高端)", "650W-1200W", "Gold", "全模组", "10年", "$95-$240"],
        ["GF1 (主流)", "650W-1200W", "Gold", "全模组", "5年", "$80-$220"],
        ["XP (白金旗舰)", "850W-1200W", "Platinum", "全模组", "10年", "$170-$310"],
        ["PF1 (白金)", "650W-1200W", "Platinum", "全模组", "10年", "$110-$280"]
    ]
    add_table_slide(prs, "二、曜越T.T电竞电源产品线（旗舰/高端）", headers, rows, highlight_color=(26, 82, 118))

    # 第5页：曜越T.T产品线 - RGB/主流
    headers = ["系列", "瓦数覆盖", "80+认证", "模组化", "质保", "特色功能"]
    rows = [
        ["Grand RGB DPS G", "650W-1050W", "Gold", "全模组", "10年", "RGB风扇+DPS监控"],
        ["XT Snow", "750W-850W", "Gold", "全模组", "5年", "全白外观"],
        ["Smart BX3 Pro SE", "650W-850W", "Bronze", "半模组", "5年", "ATX 3.1"],
        ["Smart BX1", "430W-600W", "Bronze", "非模组", "3年", "入门主流"],
        ["Litepower", "350W-600W", "White", "非模组", "2年", "最低成本"]
    ]
    add_table_slide(prs, "二、曜越T.T电竞电源产品线（RGB/主流/入门）", headers, rows, highlight_color=(39, 174, 96))

    # 第6页：CoBM切入优先级
    headers = ["优先级", "产品线", "瓦数段", "切入理由", "建议模式"]
    rows = [
        ["P1 - 最高", "Smart BX1", "430-600W", "非模组，用户量大，升级需求极强", "基础模组线材包"],
        ["P1 - 最高", "GF1系列", "650-1200W", "最大众化高端线，CoBM价值最高", "高性价比模组升级套装"],
        ["P2 - 高", "Smart BX3 Pro SE", "650-850W", "半模组，ATX 3.1增强接口需求", "ATX 3.1增强接口线缆"],
        ["P2 - 高", "GF2系列", "650-1200W", "主流玩家，Premium线材需求", "镀银/编织网版本"],
        ["P3 - 中", "GF3系列", "650-1650W", "全模组，CoBM价值较低", "仅高端定制线材合作"]
    ]
    add_table_slide(prs, "三、CoBM方案切入优先级", headers, rows, highlight_color=(231, 76, 60))

    # 第7页：最佳切入型号
    add_content_slide(
        prs,
        "三、最佳CoBM切入型号TOP5",
        [
            "【#1】Smart BX1系列全系 → 430-600W → 基础模组线材包，电商走量",
            "【#2】GF1 650W/750W → 650-750W → 高性价比模组升级套装",
            "【#3】GF1 850W/1050W → 850-1050W → 主流玩家Premium套装",
            "【#4】GF1 1200W → 1200W → 旗舰级镀银编织线材",
            "【#5】Smart BX3 Pro SE → 650-850W → ATX 3.1增强接口线缆"
        ],
        highlight_color=(155, 89, 182)
    )

    # 第8页：竞品分析
    headers = ["品牌", "代表系列", "80+认证", "功率密度", "智能功能", "价格区间", "质保"]
    rows = [
        ["华硕 ROG", "Thor系列", "Platinum", "高", "Aura Sync+OLED", "$145-$350", "10年"],
        ["微星", "MEG/MPG系列", "Platinum", "高", "MSI Center AI", "$160-$280", "10年"],
        ["酷冷至尊", "V Gold/Platinum", "Gold/Plat", "中高", "无", "$105-$210", "10年"],
        ["安钛克", "HCG/RM系列", "Gold", "中", "无", "$90-$180", "10年"],
        ["海韵", "Focus/PRIME", "Gold/Plat", "中", "无", "$130-$250", "10年"],
        ["曜越T.T", "GF3/XP系列", "Gold/Plat", "高", "T.T AI+RGB", "$110-$440", "10年"]
    ]
    add_table_slide(prs, "四、竞品分析", headers, rows, highlight_color=(52, 152, 219))

    # 第9页：CoBM vs 竞品优势
    add_comparison_slide(
        prs,
        "五、CoBM方案竞争优势",
        "传统电源方案",
        [
            "效率: Gold为主 (89-92%)",
            "功率密度: 18-22 W/in3",
            "BOM成本: $50-55 (750W)",
            "焊点数量: ~380个",
            "纹波(12V): 25-35mV",
            "MTBF: 100k小时",
            "质保: 10年(主流)"
        ],
        "CoBM方案",
        [
            "效率: Platinum (92-94%) (+2~4%)",
            "功率密度: 28-35 W/in3 (+50~65%)",
            "BOM成本: $38-48 (-20~25%)",
            "焊点数量: ~200个 (-47%)",
            "纹波(12V): <15mV (-50~60%)",
            "MTBF: 120k+小时 (+20%)",
            "质保: 12年(目标) (+20%)"
        ]
    )

    # 第10页：量化对比
    headers = ["对比维度", "CoBM优势", "量化提升", "年化收益估算"]
    rows = [
        ["转换效率", "vs Gold竞品", "+3-5%", "年省电$8-12/台"],
        ["智能功能", "vs 酷冷/海韵/安钛克", "+40%", "功能溢价$20-30"],
        ["RGB生态", "vs 无RGB竞品", "+50%", "品牌联动价值"],
        ["质保年限", "vs 行业标准", "+20%", "信心保障溢价"],
        ["功率密度", "vs Gold竞品", "+15-20%", "体积缩小/高性能"],
        ["可靠性", "vs 传统方案", "+80%", "返修率降低50%+"]
    ]
    add_table_slide(prs, "五、竞争优势量化分析", headers, rows, highlight_color=(241, 196, 15))

    # 第11页：销售增长预测
    headers = ["瓦数段", "目标市场", "当前市占", "CoBM导入后预测", "预计销量增长"]
    rows = [
        ["550W-650W", "主流电竞(RTX 4060级别)", "~18%", "~28%", "+35%~50%"],
        ["750W-850W", "中高端电竞(RTX 4070级别)", "~22%", "~32%", "+28%~42%"],
        ["1000W-1200W", "高端电竞(RTX 4080级别)", "~25%", "~32%", "+20%~30%"],
        ["1200W+", "旗舰/工作站", "~20%", "~25%", "+15%~25%"],
        ["650W SFX", "ITX紧凑电竞", "~15%", "~25%", "+50%~70%"]
    ]
    add_table_slide(prs, "六、销售增长预测（基准情景）", headers, rows, highlight_color=(39, 174, 96))

    # 第12页：定价策略
    add_content_slide(
        prs,
        "六、产品定位与定价策略",
        [
            "【650W段】$95-$115 → CoBM版:$85-$100 → 降幅10-15%",
            "【750W段】$115-$135 → CoBM版:$100-$120 → 降幅12-15%",
            "【850W段】$135-$170 → CoBM版:$120-$150 → 降幅10-15%",
            "【1000W段】$180-$210 → CoBM版:$160-$185 → 降幅10-12%",
            "【1200W段】$220-$260 → CoBM版:$200-$230 → 降幅8-12%",
            "【策略核心】加量不加价: 同价位提供更高效率+更长质保"
        ],
        highlight_color=(26, 82, 118)
    )

    # 第13页：策略建议
    add_content_slide(
        prs,
        "七、市场策略建议",
        [
            "【产品策略】优先推出GF1(650-1200W)和Smart BX1(430-600W)专用CoBM线材包",
            "【定价策略】Smart BX1 CoBM套装:$15-$25(走量)；GF1 CoBM套装:$25-$45(利润)",
            "【渠道策略】与曜越联合品牌认证，官方商城/京东/天猫同步上架",
            "【差异化策略】GF3/GX3 SE专用12V-2x6原生线材升级",
            "【推广策略】KOL装机评测+电商直播+用户晒单奖励",
            "【服务策略】目标12年质保，超越行业标准建立信心"
        ],
        highlight_color=(231, 76, 60)
    )

    # 第14页：结论
    add_conclusion_slide(
        prs,
        "结论",
        [
            "CoBM方案性价比最优：750W电竞电源成本-25%, 效率+2%, 功率密度+50%",
            "最佳切入型号：GF1系列(650-1200W)和Smart BX1系列(430-600W)",
            "核心竞争优势：T.T AI智能控制+T.T RGB生态，差异化明显",
            "销售增长预测：主流瓦数段增长28%-50%，旗舰段增长15%-30%",
            "合作建议：与曜越联合推出品牌认证CoBM升级套装"
        ]
    )

    # 保存
    output_path = r'D:/claude mini max 2.7/.claude/CoBM_曜越T.T分析报告_20260406.pptx'
    prs.save(output_path)
    print(f"PPT已生成: {output_path}")
    return output_path

if __name__ == '__main__':
    create_presentation()
