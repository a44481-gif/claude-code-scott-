# -*- coding: utf-8 -*-
"""
PPT 讀取腳本 - 分析 PowerPoint 檔案內容
支援：文字、表格、圖表、圖片分析
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
from datetime import datetime


def analyze_presentation(file_path):
    """分析 PPT 完整內容"""
    prs = Presentation(file_path)

    result = {
        "file": file_path,
        "total_slides": len(prs.slides),
        "slides": [],
        "statistics": {
            "total_shapes": 0,
            "total_tables": 0,
            "total_images": 0,
            "total_charts": 0,
            "total_text_boxes": 0
        }
    }

    for idx, slide in enumerate(prs.slides, 1):
        slide_data = {
            "slide_number": idx,
            "title": "",
            "shapes": [],
            "tables": [],
            "images": [],
            "charts": [],
            "text_content": []
        }

        for shape in slide.shapes:
            result["statistics"]["total_shapes"] += 1

            # 取得標題
            if shape.has_text_frame and hasattr(shape, 'text') and shape.text:
                if not slide_data["title"]:
                    slide_data["title"] = shape.text[:100]
                slide_data["text_content"].append({
                    "type": "text",
                    "content": shape.text.strip(),
                    "length": len(shape.text)
                })
                result["statistics"]["total_text_boxes"] += 1

            # 表格
            if shape.has_table:
                result["statistics"]["total_tables"] += 1
                table_data = extract_table_data(shape.table)
                slide_data["tables"].append(table_data)

            # 圖表
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                result["statistics"]["total_charts"] += 1
                slide_data["charts"].append({
                    "type": "chart",
                    "has_chart": True
                })

            # 圖片
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                result["statistics"]["total_images"] += 1
                slide_data["images"].append({
                    "type": "image",
                    "name": shape.name
                })

        # 組合所有文字
        all_text = []
        for tc in slide_data["text_content"]:
            all_text.append(tc["content"])
        slide_data["full_text"] = "\n".join(all_text)

        result["slides"].append(slide_data)

    return result


def extract_table_data(table):
    """提取表格資料"""
    rows_data = []
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            row_data.append(cell.text.strip())
        rows_data.append(row_data)
    return rows_data


def save_analysis(result, output_path):
    """儲存分析結果"""
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    print(f"分析結果已儲存: {output_path}")


def print_summary(result):
    """列印摘要"""
    print("=" * 60)
    print(f"[PPT] Analysis Report: {result['file']}")
    print("=" * 60)
    print(f"Total Slides: {result['total_slides']}")
    print(f"Total Shapes: {result['statistics']['total_shapes']}")
    print(f"Text Boxes: {result['statistics']['total_text_boxes']}")
    print(f"Tables: {result['statistics']['total_tables']}")
    print(f"Charts: {result['statistics']['total_charts']}")
    print(f"Images: {result['statistics']['total_images']}")
    print("-" * 60)

    for slide in result['slides']:
        print(f"\n[Slide {slide['slide_number']}] {slide['title'] or '(No Title)'}")
        if slide['tables']:
            print(f"  - Contains {len(slide['tables'])} table(s)")
        if slide['charts']:
            print(f"  - Contains {len(slide['charts'])} chart(s)")


if __name__ == "__main__":
    import sys

    # 分析指定檔案
    file_path = sys.argv[1] if len(sys.argv) > 1 else "Global_Storage_Market_2025_2030.pptx"
    output_path = sys.argv[2] if len(sys.argv) > 2 else "ppt_analysis_result.json"

    print(f"正在分析: {file_path}")

    result = analyze_presentation(file_path)
    print_summary(result)
    save_analysis(result, output_path)
