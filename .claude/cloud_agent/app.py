# -*- coding: utf-8 -*-
"""
CoBM Cloud Agent - FastAPI 主程式
雲端 AI 自動化系統核心
"""

import os
import json
import uuid
import asyncio
import httpx
from datetime import datetime
from typing import Optional, Dict, Any
from pathlib import Path

from fastapi import FastAPI, HTTPException, BackgroundTasks, Depends, Security, APIKeyHeader
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from mangum import Mangum

# ============== 配置 ==============
CLAUDE_API_KEY = os.environ.get("CLAUDE_API_KEY", "")
SMTP_AUTH_CODE = os.environ.get("SMTP_AUTH_CODE", "")
API_SECRET_KEY = os.environ.get("API_SECRET_KEY", "cobm-secret-key-2026")
REPORTS_DIR = Path("/tmp/reports") if os.environ.get("AWS_LAMBDA_FUNCTION_NAME") else Path("./reports")
REPORTS_DIR.mkdir(exist_ok=True)

# ============== FastAPI 應用 ==============
app = FastAPI(
    title="CoBM Cloud Agent API",
    description="雲端 AI 自動化系統 - 報告生成、數據分析、郵件發送",
    version="1.0"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# API Key 認證
API_KEY_HEADER = APIKeyHeader(name="X-API-Key")

async def verify_api_key(api_key: str = Security(API_KEY_HEADER)):
    if api_key != API_SECRET_KEY:
        raise HTTPException(status_code=403, detail="無效的 API Key")
    return api_key

# ============== 數據模型 ==============

class TaskRequest(BaseModel):
    task_type: str  # "full_pipeline" | "report" | "crawl" | "analyze"
    params: dict = {}
    callback_url: Optional[str] = None
    email_to: Optional[str] = None

class TaskStatus(BaseModel):
    task_id: str
    status: str  # "pending" | "running" | "completed" | "failed"
    created_at: str
    completed_at: Optional[str] = None
    result: Optional[dict] = None
    error: Optional[str] = None

# ============== 任務存儲 (內存 + 文件持久化) ==============

TASKS_DB: Dict[str, Dict] = {}

def save_task(task_id: str, task_data: dict):
    """保存任務到文件"""
    task_file = REPORTS_DIR / f"task_{task_id}.json"
    with open(task_file, "w", encoding="utf-8") as f:
        json.dump(task_data, f, ensure_ascii=False, indent=2)

def load_tasks():
    """加載歷史任務"""
    global TASKS_DB
    for f in REPORTS_DIR.glob("task_*.json"):
        try:
            with open(f, "r", encoding="utf-8") as fp:
                task = json.load(fp)
                TASKS_DB[task["task_id"]] = task
        except:
            pass

load_tasks()

# ============== 核心邏輯 ==============

async def call_claude_api(prompt: str, system: str = "") -> str:
    """調用 Claude API"""
    if not CLAUDE_API_KEY:
        # 演示模式
        return f"[演示模式] 分析結果:\n\n基於您的請求，已完成數據分析。\n\n時間戳: {datetime.now().isoformat()}"

    headers = {
        "Authorization": f"Bearer {CLAUDE_API_KEY}",
        "Content-Type": "application/json",
        "anthropic-version": "2023-06-01"
    }

    payload = {
        "model": "claude-sonnet-4-20250514",
        "max_tokens": 4096,
        "messages": [{"role": "user", "content": prompt}]
    }

    if system:
        payload["system"] = system

    async with httpx.AsyncClient(timeout=60.0) as client:
        response = await client.post(
            "https://api.anthropic.com/v1/messages",
            headers=headers,
            json=payload
        )

        if response.status_code == 200:
            data = response.json()
            return data["content"][0]["text"]
        else:
            raise Exception(f"Claude API 錯誤: {response.status_code} - {response.text}")

async def crawl_market_data(params: dict) -> dict:
    """爬取市場數據（模擬）"""
    await asyncio.sleep(1)  # 模擬網絡延遲

    return {
        "timestamp": datetime.now().isoformat(),
        "gpu_market": {
            "rtx_5090": {"price": "1599", "tdp": "450W", "demand": "高"},
            "rtx_5080": {"price": "999", "tdp": "360W", "demand": "中"},
            "rx_9070": {"price": "549", "tdp": "304W", "demand": "高"}
        },
        "psu_market": {
            "trend": "大功率化",
            "atx3_1": "標配",
            "avg_wattage": "850W"
        }
    }

async def analyze_data_with_ai(data: dict, task_type: str) -> dict:
    """使用 AI 分析數據"""

    if task_type == "full_pipeline":
        prompt = f"""
        請分析以下市場數據，生成CoBM電源方案競爭力報告：

        {json.dumps(data, ensure_ascii=False, indent=2)}

        分析維度：
        1. AI PC 電源需求趨勢
        2. CoBM 方案市場機會
        3. 目標客戶推薦
        4. 推廣策略建議

        請用繁體中文輸出。
        """
    else:
        prompt = f"請分析以下數據：\n{json.dumps(data, ensure_ascii=False)}"

    analysis_text = await call_claude_api(prompt)

    return {
        "timestamp": datetime.now().isoformat(),
        "analysis": analysis_text,
        "data_source": data
    }

async def generate_ppt_report(analysis: dict, report_name: str) -> str:
    """生成 PPT 報告"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # 標題頁
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(26, 82, 118)
        bg.line.fill.background()

        title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = report_name
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        date = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.8))
        tf = date.text_frame
        p = tf.paragraphs[0]
        p.text = datetime.now().strftime("%Y年%m月%d日 %H:%M")
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

        # 分析頁
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(39, 174, 96)
        header.line.fill.background()

        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "AI 分析結果"
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        content = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.5))
        tf = content.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = analysis.get("analysis", "分析結果")[0:2000]
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(50, 50, 50)

        # 保存
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"CoBM_Report_{timestamp}.pptx"
        filepath = REPORTS_DIR / filename
        prs.save(str(filepath))

        return str(filepath)

    except ImportError:
        # python-pptx 未安裝
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"CoBM_Report_{timestamp}.txt"
        filepath = REPORTS_DIR / filename
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(f"{report_name}\n\n")
            f.write(f"分析時間: {datetime.now().isoformat()}\n\n")
            f.write(f"分析結果:\n{analysis.get('analysis', '')}\n\n")
            f.write(f"數據來源:\n{json.dumps(analysis.get('data_source', {}), ensure_ascii=False, indent=2)}")
        return str(filepath)

async def send_email_report(report_path: str, email_to: str) -> dict:
    """發送郵件報告"""
    import smtplib
    from email.mime.multipart import MIMEMultipart
    from email.mime.text import MIMEText
    from email.mime.application import MIMEApplication

    if not SMTP_AUTH_CODE:
        return {"status": "demo", "message": "SMTP未配置，演示模式"}

    msg = MIMEMultipart('mixed')
    msg['Subject'] = f'CoBM 雲端分析報告 {datetime.now().strftime("%Y-%m-%d %H:%M")}'
    msg['From'] = 'scott365888@gmail.com'
    msg['To'] = email_to

    # HTML 正文
    html = f"""
    <html><body style="font-family: Arial; padding: 20px;">
    <h2 style="color: #1a5276;">CoBM 雲端分析報告</h2>
    <p>報告生成時間：{datetime.now().strftime("%Y-%m-%d %H:%M:%S")}</p>
    <p>任務類型：全流程自動化分析</p>
    <hr>
    <p>完整報告請查看附件。</p>
    <p style="color: #666; font-size: 12px;">本報告由 CoBM Cloud Agent 自動生成</p>
    </body></html>
    """
    msg.attach(MIMEText(html, 'html', 'utf-8'))

    # 附件
    if os.path.exists(report_path):
        with open(report_path, 'rb') as f:
            part = MIMEApplication(f.read(), Name=os.path.basename(report_path))
            part['Content-Disposition'] = f'attachment; filename="{os.path.basename(report_path)}"'
            msg.attach(part)

    try:
        # Gmail SMTP
        server = smtplib.SMTP_SSL('smtp.gmail.com', 465)
        server.login('scott365888@gmail.com', SMTP_AUTH_CODE)
        server.sendmail('scott365888@gmail.com', email_to, msg.as_string())
        server.quit()
        return {"status": "sent", "to": email_to}
    except Exception as e:
        return {"status": "error", "message": str(e)}

async def execute_task(task_id: str, req: TaskRequest):
    """執行任務"""
    try:
        TASKS_DB[task_id]["status"] = "running"
        save_task(task_id, TASKS_DB[task_id])

        result = {}

        if req.task_type == "full_pipeline":
            # 全流程
            data = await crawl_market_data(req.params)
            analysis = await analyze_data_with_ai(data, req.task_type)
            report_path = await generate_ppt_report(analysis, "CoBM 市場分析報告")

            email_to = req.email_to or "scott365888@gmail.com"
            email_result = await send_email_report(report_path, email_to)

            result = {
                "crawl": data,
                "analysis": analysis,
                "report": report_path,
                "report_name": os.path.basename(report_path),
                "email": email_result
            }

        elif req.task_type == "report":
            report_path = await generate_ppt_report(req.params, req.params.get("title", "CoBM 報告"))
            result = {"report": report_path}

        elif req.task_type == "crawl":
            result = await crawl_market_data(req.params)

        elif req.task_type == "analyze":
            result = await analyze_data_with_ai(req.params.get("data", {}), req.task_type)

        TASKS_DB[task_id]["status"] = "completed"
        TASKS_DB[task_id]["completed_at"] = datetime.now().isoformat()
        TASKS_DB[task_id]["result"] = result
        save_task(task_id, TASKS_DB[task_id])

        # 回調
        if req.callback_url:
            try:
                async with httpx.AsyncClient(timeout=30.0) as client:
                    await client.post(req.callback_url, json=result)
            except:
                pass

    except Exception as e:
        TASKS_DB[task_id]["status"] = "failed"
        TASKS_DB[task_id]["error"] = str(e)
        save_task(task_id, TASKS_DB[task_id])

# ============== API 端點 ==============

@app.post("/api/v1/tasks", response_model=Dict)
async def create_task(req: TaskRequest, background_tasks: BackgroundTasks, api_key: str = Depends(verify_api_key)):
    """創建新任務"""
    task_id = str(uuid.uuid4())[:8]

    TASKS_DB[task_id] = {
        "task_id": task_id,
        "type": req.task_type,
        "params": req.params,
        "status": "pending",
        "created_at": datetime.now().isoformat(),
        "callback_url": req.callback_url,
        "email_to": req.email_to or "scott365888@gmail.com"
    }
    save_task(task_id, TASKS_DB[task_id])

    background_tasks.add_task(execute_task, task_id, req)

    return {
        "task_id": task_id,
        "status": "pending",
        "message": f"任務已創建，ID: {task_id}",
        "check_url": f"/api/v1/tasks/{task_id}"
    }

@app.get("/api/v1/tasks/{task_id}")
async def get_task(task_id: str, api_key: str = Depends(verify_api_key)):
    """查詢任務狀態"""
    if task_id not in TASKS_DB:
        raise HTTPException(status_code=404, detail="任務不存在")

    task = TASKS_DB[task_id]
    return {
        "task_id": task_id,
        "status": task["status"],
        "created_at": task["created_at"],
        "completed_at": task.get("completed_at"),
        "result": task.get("result"),
        "error": task.get("error")
    }

@app.get("/api/v1/tasks/{task_id}/download")
async def download_report(task_id: str, api_key: str = Depends(verify_api_key)):
    """下載報告"""
    if task_id not in TASKS_DB:
        raise HTTPException(status_code=404, detail="任務不存在")

    task = TASKS_DB[task_id]
    if task["status"] != "completed":
        raise HTTPException(status_code=400, detail="任務未完成")

    result = task.get("result", {})
    report_path = result.get("report")

    if not report_path or not os.path.exists(report_path):
        raise HTTPException(status_code=404, detail="報告文件不存在")

    return FileResponse(
        report_path,
        media_type='application/octet-stream',
        filename=os.path.basename(report_path)
    )

@app.get("/api/v1/tasks")
async def list_tasks(api_key: str = Depends(verify_api_key)):
    """列出所有任務"""
    tasks = []
    for tid, task in TASKS_DB.items():
        tasks.append({
            "task_id": tid,
            "type": task.get("type"),
            "status": task.get("status"),
            "created_at": task.get("created_at")
        })
    return {"tasks": tasks, "total": len(tasks)}

# 健康檢查
@app.get("/health")
async def health():
    return {
        "status": "healthy",
        "time": datetime.now().isoformat(),
        "version": "1.0"
    }

@app.get("/")
async def root():
    return {
        "name": "CoBM Cloud Agent API",
        "version": "1.0",
        "docs": "/docs",
        "endpoints": {
            "create_task": "POST /api/v1/tasks",
            "check_status": "GET /api/v1/tasks/{task_id}",
            "download": "GET /api/v1/tasks/{task_id}/download",
            "list_tasks": "GET /api/v1/tasks",
            "health": "GET /health"
        },
        "example": {
            "task_type": "full_pipeline",
            "params": {},
            "email_to": "scott365888@gmail.com"
        }
    }

# Lambda 適配器
handler = Mangum(app)

# ============== 本地運行 ==============

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
