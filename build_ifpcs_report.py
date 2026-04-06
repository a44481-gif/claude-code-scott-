# -*- coding: utf-8 -*-
"""
CCS集成母排 × 全球储能市场
智能融合电力连接系统（IFPCS）综合商业计划书
包含：全赛道痛点体系 + 全球市场经济数据 + 竞争分析
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from copy import deepcopy
import os

SRC_CCS = r'd:/claude mini max 2.7/storage_market_output/CCS_Global_Storage_2025_2030.pptx'
OUTPUT = r'd:/claude mini max 2.7/storage_market_output/CCS_Global_IFPCS_Final_Report.pptx'

# ─────────────────────────────────────────────
# 核心工具函數
# ─────────────────────────────────────────────
def find_cSld(tree):
    for elem in tree.iter():
        if elem.tag == qn('p:cSld'):
            return elem
    return None

def deep_copy_slide(target_prs, source_prs, slide_idx):
    """深度複製幻燈片（保留圖片和格式）"""
    src_slide = source_prs.slides[slide_idx]
    src_tree = src_slide.part._element
    blank = target_prs.slide_layouts[6]
    new_slide = target_prs.slides.add_slide(blank)
    new_tree = new_slide.part._element
    cSld_src = find_cSld(src_tree)
    cSld_new = find_cSld(new_tree)
    if cSld_src is not None and cSld_new is not None:
        bg_src = cSld_src.find(qn('p:bg'))
        bg_new = cSld_new.find(qn('p:bg'))
        if bg_new is not None:
            cSld_new.remove(bg_new)
        if bg_src is not None:
            cSld_new.insert(0, deepcopy(bg_src))
        spTree_src = cSld_src.find(qn('p:spTree'))
        spTree_new = cSld_new.find(qn('p:spTree'))
        if spTree_new is not None:
            for child in list(spTree_new):
                if child.tag not in (qn('p:nvGrpSpPr'), qn('p:grpSpPr')):
                    spTree_new.remove(child)
        if spTree_src is not None:
            for child in spTree_src:
                if child.tag not in (qn('p:nvGrpSpPr'), qn('p:grpSpPr')):
                    spTree_new.append(deepcopy(child))
    return new_slide

def add_divider(prs, part_num, title1, title2="", bg=(0x0A, 0x29, 0x5E)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*bg)
    def t(text, size, bold, color_rgb, y, h):
        tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.3), h)
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text
        p.font.size = Pt(size); p.font.bold = bold
        p.font.color.rgb = RGBColor(*color_rgb); p.alignment = PP_ALIGN.CENTER
    t(f"PART {part_num}", 20, True, (0xFF,0xD7,0x00), Inches(1.8), Inches(0.8))
    t(title1, 36, True, (255,255,255), Inches(2.6), Inches(1.2))
    if title2: t(title2, 28, True, (255,255,255), Inches(3.8), Inches(1.0))
    return slide

def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x0A, 0x29, 0x5E)
    def t(text, size, bold, color_rgb, y, h):
        tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.3), h)
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text
        p.font.size = Pt(size); p.font.bold = bold
        p.font.color.rgb = RGBColor(*color_rgb); p.alignment = PP_ALIGN.CENTER
    t("CCS集成母排 × 全球储能市场", 38, True, (255,255,255), Inches(1.2), Inches(1.5))
    t("智能融合电力连接系统（IFPCS）", 30, True, (0xFF,0xD7,0x00), Inches(2.7), Inches(1.0))
    t("综合商业计划书 2026-2030", 22, False, (0xCC,0xD9,0xEC), Inches(3.8), Inches(0.8))
    t("全赛道痛点分析 · 全球市场经济数据 · 竞争格局对标", 16, False, (0x99,0xAA,0xCC), Inches(4.7), Inches(0.6))
    t("2026年4月", 14, False, (0x77,0x88,0xAA), Inches(5.5), Inches(0.5))
    return slide

def add_table_slide(prs, title, headers, rows, col_widths=None, title_bg=(0x0A,0x29,0x5E)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xF0,0xF4,0xF8)
    # Title bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*title_bg); bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(12.7), Inches(0.65))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)
    # Table
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [Inches(12.3 / n_cols)] * n_cols
    else:
        col_widths = [Inches(w) for w in col_widths]
    row_h = min(Inches(0.42), Inches(5.5 / (len(rows)+1)))
    tbl = slide.shapes.add_table(len(rows)+1, n_cols, Inches(0.3), Inches(1.05),
                                  sum(col_widths), row_h * (len(rows)+1)).table
    for j, h in enumerate(headers):
        c = tbl.cell(0, j); c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = RGBColor(*title_bg)
        p2 = c.text_frame.paragraphs[0]; p2.font.size = Pt(10); p2.font.bold = True
        p2.font.color.rgb = RGBColor(255,255,255); p2.alignment = PP_ALIGN.CENTER
        tbl.columns[j].width = col_widths[j]
    for i, row in enumerate(rows):
        bg = RGBColor(255,255,255) if i % 2 == 0 else RGBColor(0xE3,0xF0,0xFF)
        for j, val in enumerate(row):
            c = tbl.cell(i+1, j); c.text = str(val)
            c.fill.solid(); c.fill.fore_color.rgb = bg
            p3 = c.text_frame.paragraphs[0]; p3.font.size = Pt(9)
            p3.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
    return slide

def add_kpi_slide(prs, title, kpis, title_bg=(0x0A,0x29,0x5E)):
    """kpis = list of (number, unit, label, sub)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xF0,0xF4,0xF8)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*title_bg); bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(12.7), Inches(0.65))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)
    n = len(kpis)
    bw = 12.3 / n
    for i, (num, unit, label, sub) in enumerate(kpis):
        x = 0.3 + i * bw
        box = slide.shapes.add_shape(1, Inches(x), Inches(1.15), Inches(bw-0.12), Inches(2.8))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(255,255,255); box.line.color.rgb = RGBColor(200,200,200)
        ntb = slide.shapes.add_textbox(Inches(x), Inches(1.35), Inches(bw-0.12), Inches(0.9))
        np_ = ntb.text_frame.paragraphs[0]; np_.text = num
        np_.font.size = Pt(38); np_.font.bold = True
        np_.font.color.rgb = RGBColor(*title_bg); np_.alignment = PP_ALIGN.CENTER
        utb = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(bw-0.12), Inches(0.4))
        up = utb.text_frame.paragraphs[0]; up.text = unit
        up.font.size = Pt(14); up.font.bold = True
        up.font.color.rgb = RGBColor(0xFF,0xD7,0x00); up.alignment = PP_ALIGN.CENTER
        ltb = slide.shapes.add_textbox(Inches(x+0.05), Inches(2.65), Inches(bw-0.18), Inches(0.5))
        lp = ltb.text_frame.paragraphs[0]; lp.text = label
        lp.font.size = Pt(10); lp.font.bold = True
        lp.font.color.rgb = RGBColor(0x33,0x33,0x33); lp.alignment = PP_ALIGN.CENTER
        if sub:
            stb = slide.shapes.add_textbox(Inches(x+0.05), Inches(3.1), Inches(bw-0.18), Inches(0.4))
            sp_ = stb.text_frame.paragraphs[0]; sp_.text = sub
            sp_.font.size = Pt(8); sp_.font.color.rgb = RGBColor(0x88,0x88,0x88); sp_.alignment = PP_ALIGN.CENTER
    return slide

def add_three_col_slide(prs, title, cols, title_bg=(0x0A,0x29,0x5E)):
    """cols = list of (header, [bullet list])"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xF0,0xF4,0xF8)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*title_bg); bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(12.7), Inches(0.65))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)
    n = len(cols)
    bw = 12.3 / n
    for i, (hdr, bullets) in enumerate(cols):
        x = 0.3 + i * bw
        # Header
        hbox = slide.shapes.add_shape(1, Inches(x), Inches(1.05), Inches(bw-0.1), Inches(0.5))
        hbox.fill.solid(); hbox.fill.fore_color.rgb = RGBColor(0x1A,0x5C,0x9E); hbox.line.fill.background()
        htb = slide.shapes.add_textbox(Inches(x+0.05), Inches(1.1), Inches(bw-0.15), Inches(0.4))
        hp = htb.text_frame.paragraphs[0]; hp.text = hdr
        hp.font.size = Pt(12); hp.font.bold = True; hp.font.color.rgb = RGBColor(255,255,255); hp.alignment = PP_ALIGN.CENTER
        # Content box
        cbox = slide.shapes.add_shape(1, Inches(x), Inches(1.6), Inches(bw-0.1), Inches(5.3))
        cbox.fill.solid(); cbox.fill.fore_color.rgb = RGBColor(255,255,255); cbox.line.color.rgb = RGBColor(200,200,200)
        ctb = slide.shapes.add_textbox(Inches(x+0.1), Inches(1.7), Inches(bw-0.25), Inches(5.1))
        ctf = ctb.text_frame; ctf.word_wrap = True
        for j, bullet in enumerate(bullets):
            p2 = ctf.paragraphs[0] if j == 0 else ctf.add_paragraph()
            p2.text = bullet; p2.font.size = Pt(9.5)
            p2.space_before = Pt(3)
    return slide

def add_text_slide(prs, title, lines, title_bg=(0x0A,0x29,0x5E)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xF0,0xF4,0xF8)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*title_bg); bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(12.7), Inches(0.65))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)
    cbox = slide.shapes.add_shape(1, Inches(0.3), Inches(1.05), Inches(12.73), Inches(6.2))
    cbox.fill.solid(); cbox.fill.fore_color.rgb = RGBColor(255,255,255); cbox.line.color.rgb = RGBColor(200,200,200)
    ctb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.9))
    ctf = ctb.text_frame; ctf.word_wrap = True
    for j, line in enumerate(lines):
        p2 = ctf.paragraphs[0] if j == 0 else ctf.add_paragraph()
        p2.text = line; p2.font.size = Pt(11)
        p2.space_before = Pt(4)
    return slide

def add_thanks_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x0A, 0x29, 0x5E)
    def t(text, size, bold, color_rgb, y, h):
        tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.3), h)
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text
        p.font.size = Pt(size); p.font.bold = bold
        p.font.color.rgb = RGBColor(*color_rgb); p.alignment = PP_ALIGN.CENTER
    t("谢谢聆听", 52, True, (255,255,255), Inches(2.0), Inches(1.5))
    t("THANK YOU FOR YOUR TIME", 22, False, (0xFF,0xD7,0x00), Inches(3.7), Inches(0.8))
    t("智能融合电力连接系统（IFPCS）综合商业计划书  |  2026年4月", 14, False, (0x77,0x88,0xAA), Inches(4.8), Inches(0.5))
    return slide

# ─────────────────────────────────────────────
# 加載源文件
# ─────────────────────────────────────────────
print("載入CCS商業方案PPT...")
prs_ccs = Presentation(SRC_CCS)
print(f"CCS PPT: {len(prs_ccs.slides)} 張")

# 建立新簡報
prs = Presentation()
prs.slide_width = prs_ccs.slide_width
prs.slide_height = prs_ccs.slide_height

# ─────────────────────────────────────────────
# 1. 封面
# ─────────────────────────────────────────────
print("[1/8] 加入封面...")
add_cover(prs)

# ─────────────────────────────────────────────
# 2. PART 0 - 全賽道痛點體系（3頁）
# ─────────────────────────────────────────────
print("[2/8] 建立全賽道痛點體系...")

add_divider(prs, "0", "全赛道用电市场痛点体系", "从单一产品缺陷到全产业链系统性瓶颈")

# 頁面1：5層痛點金字塔
add_table_slide(prs,
    "传统电力连接方案的5层系统性痛点",
    ["痛点层级", "全赛道通用核心问题", "量化影响", "IFPCS解决方案"],
    [
        ["技术性能层", "大电流温升/振动松脱/盐雾失效/电磁干扰",
         "80%电池热失控源于连接点；-40℃性能衰减60%",
         "一体化铜排均流+热压密封+电磁屏蔽层"],
        ["生产制造层", "人工依赖高/良率低/换型慢/交付长",
         "人工布线良率85%；单车型换型需15-30天",
         "模块化预装+全自动化产线；12天样品交付"],
        ["现场应用层", "安装错误/故障难查/单点故障连锁",
         "现场安装错误率12%；故障排查平均48小时",
         "即插即用+BMS智能诊断；故障风险降90%"],
        ["全生命周期层", "研发贵/装配贵/维护贵/回收难",
         "传统线束占整车成本8-12%；售后成本是制造成本3倍",
         "标准化平台；研发周期3-6个月；可拆解环保设计"],
        ["产业生态层", "供应链波动/认证周期长/迭代滞后",
         "铜价波动10%影响成本5-8%；海外认证6-12个月",
         "国产化替代+预认证通用模块+三代技术储备"],
    ],
    col_widths=[1.3, 3.2, 3.2, 4.6]
)

# 頁面2：三大黃金賽道差異化致命痛點
add_three_col_slide(prs,
    "三大核心赛道的差异化致命痛点",
    [
        ("新能源汽车（主战场）",
         ["800V高压平台：局部温升超标、电弧击穿风险",
          "CTC/CTB技术：线束维修不可达、整车报废率上升",
          "4680大圆柱：多连接点均流不均、一致性差",
          "",
          "量化数据：",
          "800V车型传统线束故障率是400V的2.3倍",
          "CTP技术对连接可靠性要求提升60%",
          "单车线束维修成本平均超3000元"]),
        ("工商业储能（增长极）",
         ["液冷储能系统：防水密封要求高、线缆腐蚀老化快",
          "集装箱式部署：空间极致压缩、布线难度指数级上升",
          "电网调频：高频充放电导致连接点疲劳断裂",
          "",
          "量化数据：",
          "储能系统连接故障占总停机时间40%",
          "集装箱储能布线人工成本占比35%",
          "年维护成本超初始成本的15%"]),
        ("AI数据中心（蓝海）",
         ["超高功率密度：>30kW/柜无法适配传统线束",
          "高压直流供电：拉弧风险高、安全防护难度大",
          "7×24运行：单点故障导致单柜损失超10万元/小时",
          "",
          "量化数据：",
          "数据中心断电每分钟损失超5万美元",
          "BBU备电系统年市场规模超150亿元",
          "AI集群单次断电损失可达千万美元级"]),
    ]
)

# 頁面3：七大高潛力賽道專屬痛點
add_table_slide(prs,
    "七大高潜力赛道的差异化致命痛点与技术壁垒",
    ["赛道", "核心极端工况", "专属致命痛点", "IFPCS核心价值"],
    [
        ["电动重卡/工程机械", "重载振动/高低温/粉尘", "接头松动断裂；停机成本>5000元/小时",
         "抗震一体化+防腐蚀密封+免维护设计"],
        ["船舶/海上风电", "盐雾腐蚀/潮湿/摇摆", "绝缘老化快；海上运维成本是陆地10倍",
         "船级社认证材料+全密封结构+长寿命设计"],
        ["轨道交通", "高频振动/强电磁干扰", "信号采集失真；寿命需匹配30年",
         "抗干扰屏蔽层+高疲劳强度材料+全生命周期监测"],
        ["户用储能/便携式电源", "低温/频繁移动/消费级安全", "低温性能衰减；跌落损坏；消费者起火焦虑",
         "低温适配+抗冲击结构+多重安全保护"],
        ["无人机/工业机器人", "轻量化/快速充电/高动态", "重量限制续航；关节弯折疲劳",
         "超薄柔性设计+集成散热+高弯折寿命线路"],
        ["航空航天", "极端高低温/真空/高可靠性", "重量克级要求；单点故障即灾难性后果",
         "极致轻量化+冗余设计+航天级材料"],
        ["光伏逆变器/充电桩", "户外暴晒/雨水侵蚀/高功率", "老化快；防护等级低；功率密度不足",
         "户外专用材料+IP67防护+高功率密度集成"],
    ],
    col_widths=[1.8, 2.0, 3.2, 5.3]
)

# ─────────────────────────────────────────────
# 3. 複製原始CCS報告（deep_copy_slide）
# ─────────────────────────────────────────────
print("[3/8] 複製原始CCS商業方案...")
for i in range(len(prs_ccs.slides)):
    deep_copy_slide(prs, prs_ccs, i)
    print(f"  CCS {i+1}/{len(prs_ccs.slides)}", end="\r")
print()

# ─────────────────────────────────────────────
# 4. PART 3 - 全球市場經濟數據（4頁）
# ─────────────────────────────────────────────
print("[4/8] 建立全球市場經濟數據體系...")

add_divider(prs, "3", "全球市场经济数据补全", "区域格局 · 盈利结构 · 供应链 · 竞争对标", bg=(0x1A,0x5C,0x9E))

# 全球區域市場規模
add_kpi_slide(prs,
    "全球市场格局：先深耕中国 → 再突破欧美 → 后覆盖新兴",
    [
        ("350", "亿元", "2028年全球IFPCS市场总规模", "中国占比57%"),
        ("200", "亿元", "2028年中国IFPCS市场规模", "全球最大单一市场"),
        ("600", "亿元", "2030年全球IFPCS预测规模", "CAGR约28%"),
        ("45%+", "毛利率", "IFPCS综合毛利率目标", "高端定制贡献70%利润"),
    ]
)

add_table_slide(prs,
    "全球区域市场规模拆分与增长预测",
    ["区域", "2028年市场规模", "全球占比", "CAGR(26-30)", "核心驱动因素"],
    [
        ["中国", "200亿元", "57%", "32%", "新能源汽车+工商业储能+AI数据中心三重爆发"],
        ["欧洲", "70亿元", "20%", "24%", "欧盟碳中和政策+商用车电动化强制落地"],
        ["北美", "42亿元", "12%", "28%", "AI算力基建爆发+皮卡电动化+户用储能"],
        ["亚太其他", "28亿元", "8%", "26%", "东南亚产能转移+日韩储能出口需求"],
        ["其他新兴", "10亿元", "3%", "31%", "拉美/中东/非洲新能源起步"],
    ],
    col_widths=[1.2, 1.8, 1.3, 1.5, 6.5]
)

# 分產品價格帶與毛利率
add_table_slide(prs,
    "盈利结构：高端与定制级市场贡献70%以上利润",
    ["产品等级", "应用场景", "单价区间(元/套)", "行业平均毛利率", "IFPCS目标毛利率", "市场策略"],
    [
        ["入门级", "低端乘用车/户用储能", "200-500", "15-20%", "20-25%", "规模化走量，建立市场份额"],
        ["中端级", "中高端乘用车/工商业储能", "500-1500", "25-35%", "35-45%", "核心主力，贡献主要营收"],
        ["高端级", "800V平台/AI数据中心/电动重卡", "1500-5000", "35-50%", "50-60%", "重点突破，构建技术壁垒"],
        ["定制级", "航空航天/轨道交通/船舶", "5000+", "50-70%", "60-75%", "标杆项目，打造品牌溢价"],
    ],
    col_widths=[1.1, 2.4, 1.4, 1.7, 1.7, 4.0]
)

# 供應鏈上下游與成本控制
add_three_col_slide(prs,
    "供应链壁垒：垂直整合+精益生产构建成本护城河",
    [
        ("成本结构对比",
         ["传统线束成本结构：",
          "  铜材 45% | 人工 25% | 绝缘 15% | 其他 15%",
          "",
          "IFPCS成本结构：",
          "  铜材 32% | 人工 10% | 绝缘 18% | PCB 20%",
          "",
          "IFPCS核心成本优势：",
          "  一体化设计减少原材料用量30%",
          "  全自动化生产降低人工成本60%",
          "  规模化集中采购降原材料成本10-15%",
          "",
          "风险应对：",
          "  与江铜/铜陵有色签3年锁价协议",
          "  推进铝代铜技术研发",
          "  建立2-3家备选供应商"]),
        ("供应链全景",
         ["上游原材料：",
          "  铜/铝（占IFPCS成本32%）",
          "  绝缘材料（国产替代率已达70%）",
          "  PCB/连接器/氟塑料",
          "",
          "核心工艺内制：",
          "  铜排精密成型+热压密封工艺",
          "  自动化组装+在线测试",
          "  质量追溯系统（MES集成）",
          "",
          "下游直供客户：",
          "  电芯厂（模组→PACK）",
          "  整车厂（CTP/CTC集成）",
          "  储能厂（集装箱系统）",
          "  数据中心（BBU备电）"]),
        ("价值链分析",
         ["上游议价能力：",
          "  铜材集中采购+长协锁价",
          "  绝缘材料国产化，成本可控",
          "",
          "中游制造优势：",
          "  全制程自动化，良率99.5%+",
          "  12天样品交付，响应速度领先",
          "",
          "下游客户粘性：",
          "  嵌入式合作，共同研发",
          "  全生命周期服务",
          "  换型成本低，锁定客户"]),
    ]
)

# 競爭對手對標分析
add_table_slide(prs,
    "竞争格局：外资份额持续下滑，本土技术型企业崛起",
    ["竞争梯队", "代表企业", "核心优势", "核心劣势", "IFPCS差异化优势"],
    [
        ["第一梯队(≥14%)", "安波福/莱尼(外资)", "全球品牌/客户资源/规模效应",
         "响应慢/定制化差/价格高30%+",
         "样品交付快3倍(12天vs36天)/成本低30%/快速适配迭代"],
        ["第二梯队(5-10%)", "汇川/中航光电/巴斯巴", "本土客户资源/交付速度",
         "产品线单一/高端市场渗透不足",
         "全赛道覆盖10大赛道/IFPCS集成技术领先一代"],
        ["第三梯队(<5%)", "众多中小厂商", "价格极低",
         "无研发能力/质量不稳定/产能不足",
         "全制程自动化/良率99.5%+/核心专利壁垒"],
    ],
    col_widths=[1.4, 1.8, 2.2, 2.4, 5.5]
)

# ─────────────────────────────────────────────
# 5. 客戶價值測算（可選增值頁）
# ─────────────────────────────────────────────
print("[5/8] 建立客戶價值測算...")

add_three_col_slide(prs,
    "客户价值：全生命周期成本降低30%以上",
    [
        ("新能源汽车（800V车型·年销10万辆）",
         ["研发成本：1200万元/车型 → 720万元（节约480万）",
          "制造成本：800元/车 → 600元（节约2000万/年）",
          "维护成本：200元/车 → 80元（节约1200万/年）",
          "",
          "单车全生命周期节约：520元",
          "年销10万辆节约总额：3680万元",
          "",
          "额外收益：",
          "  电池包能量密度提升5%",
          "  续航增加20km",
          "  整车报废率下降15%"]),
        ("工商业储能（1GWh/年·头部储能厂）",
         ["传统方案：120元/kWh，总成本1.2亿元/年",
          "IFPCS方案：85元/kWh，总成本8500万元/年",
          "",
          "年节约成本：3500万元",
          "安装效率提升：50%",
          "",
          "额外收益：",
          "  集装箱空间节约15%",
          "  调试周期缩短30%",
          "  系统可靠性提升，质保成本下降"]),
        ("AI数据中心（1000柜·大型IDC）",
         ["传统方案故障损失：",
          "  单点故障平均停机4小时",
          "  损失超400万元/次",
          "",
          "IFPCS方案：",
          "  故障风险降低90%",
          "  年减少损失超3000万元",
          "",
          "额外收益：",
          "  备电响应时间<10ms",
          "  BBU系统体积减少30%",
          "  全天候智能监控"])],
)

# ─────────────────────────────────────────────
# 6. 總結頁
# ─────────────────────────────────────────────
print("[6/8] 加入總結頁...")
add_thanks_slide(prs)

# ─────────────────────────────────────────────
# 保存
# ─────────────────────────────────────────────
prs.save(OUTPUT)
size = os.path.getsize(OUTPUT)
total = len(prs.slides)
print(f"\n=== 完成 ===")
print(f"文件：{OUTPUT}")
print(f"大小：{size/1024/1024:.1f} MB")
print(f"總頁數：{total} 張")
print(f"  封面 1張")
print(f"  Part 0 痛點體系 3張")
print(f"  CCS商業方案 {len(prs_ccs.slides)}張")
print(f"  Part 3 市場數據 5張（含分隔頁）")
print(f"  客戶價值測算 1張")
print(f"  總結頁 1張")
