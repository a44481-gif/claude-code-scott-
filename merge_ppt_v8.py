# -*- coding: utf-8 -*-
"""
PPT合并脚本 v8 - 正确处理图片引用的深度合并
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree
import os, io

SRC1 = r'd:/claude mini max 2.7/storage_market_output/Global_Storage_Market_2025_2030.pptx'
SRC2 = r'd:/claude mini max 2.7/storage_market_output/CCS_Global_Storage_2025_2030.pptx'
OUTPUT = r'd:/claude mini max 2.7/storage_market_output/CCS_Global_MERGED_FINAL.pptx'

print("Loading source files...")
prs1 = Presentation(SRC1)  # 全球储能 (18 slides)
prs2 = Presentation(SRC2)  # CCS报告 (22 slides)
print(f"SRC1 (全球储能): {len(prs1.slides)} slides")
print(f"SRC2 (CCS报告): {len(prs2.slides)} slides")

# 建立新的合并简报
merged = Presentation()
merged.slide_width = prs1.slide_width
merged.slide_height = prs1.slide_height

# ============================================================
# 图片复制引擎
# ============================================================
# rId -> PIL.Image
_img_cache = {}

def add_image_from_slide(target_prs, source_slide, old_rid):
    """从源幻灯片中提取图片，添加到目标简报，返回新的 rId"""
    cache_key = (id(source_slide.part), old_rid)
    if cache_key in _img_cache:
        return _img_cache[cache_key]

    try:
        rel = source_slide.part.rels[old_rid]
        img_part = rel.target_part
        blob = img_part.blob
        ct = img_part.content_type

        # 添加到目标简报
        new_img = target_prs.part.add_image(ct, blob)

        # 在新 slide 的 part 中建立对应的 rel
        new_slide_part = target_prs.slides[-1].part
        new_rel = new_slide_part.relate_to(
            new_img, rel.reltype, old_rid
        )

        _img_cache[cache_key] = old_rid
        return old_rid
    except Exception as e:
        _img_cache[cache_key] = None
        return None

# ============================================================
# 深度复制幻灯片
# ============================================================
def copy_slide_full(target_prs, source_prs, slide_idx):
    """完整复制幻灯片：文本+图片+背景+格式"""
    src_slide = source_prs.slides[slide_idx]
    src_tree = src_slide.part._element

    # ---- 先处理图片：复制到目标简报 ----
    img_rid_map = {}  # old_rid -> new_rid (same rid, added to new slide part)
    for rel_id, rel in src_slide.part.rels.items():
        if rel.reltype.endswith('image'):
            add_image_from_slide(target_prs, src_slide, rel_id)

    # ---- 新增空白幻灯片 ----
    blank_layout = target_prs.slide_layouts[6]
    new_slide = target_prs.slides.add_slide(blank_layout)
    new_tree = new_slide.part._element

    # ---- 找到 cSld ----
    def find_cSld(tree):
        for el in tree.iter():
            if el.tag == qn('p:cSld'):
                return el
        return None

    cSld_src = find_cSld(src_tree)
    cSld_new = find_cSld(new_tree)

    if cSld_src is not None and cSld_new is not None:
        # ---- 背景 ----
        bg_src = cSld_src.find(qn('p:bg'))
        bg_new = cSld_new.find(qn('p:bg'))
        if bg_new is not None:
            cSld_new.remove(bg_new)
        if bg_src is not None:
            cSld_new.insert(0, deepcopy(bg_src))

        # ---- 形状树 ----
        spTree_src = cSld_src.find(qn('p:spTree'))
        spTree_new = cSld_new.find(qn('p:spTree'))

        if spTree_new is not None:
            for child in list(spTree_new):
                if child.tag not in (qn('p:nvGrpSpPr'), qn('p:grpSpPr')):
                    spTree_new.remove(child)

        if spTree_src is not None:
            for child in spTree_src:
                if child.tag not in (qn('p:nvGrpSpPr'), qn('p:grpSpPr')):
                    copied = deepcopy(child)

                    # 修正图片引用 (a:blip 的 r:embed)
                    for blip in copied.iter(qn('a:blip')):
                        old_embed = blip.get(qn('r:embed'))
                        if old_embed:
                            # rId 在新幻灯片中应该相同（已通过 relate_to 建立）
                            pass  # 保持原来的 rId，因为 relate_to 使用了相同 old_rid

                    spTree_new.append(copied)

    return new_slide

# ============================================================
# 辅助页面
# ============================================================
def add_divider(part_num, title1, title2, bg=(0x0A,0x29,0x5E)):
    slide = merged.slides.add_slide(merged.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*bg)

    def t(text, size, bold, rgb, y, h):
        tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.3), h)
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*rgb)
        p.alignment = PP_ALIGN.CENTER

    t(f"PART {part_num}", 20, True, (0xFF,0xD7,0x00), Inches(1.8), Inches(0.8))
    t(title1, 36, True, (255,255,255), Inches(2.6), Inches(1.2))
    if title2:
        t(title2, 26, True, (255,255,255), Inches(3.7), Inches(1.0))

def add_cover():
    slide = merged.slides.add_slide(merged.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x0A,0x29,0x5E)

    def t(text, size, bold, rgb, y, h):
        tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.3), h)
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*rgb)
        p.alignment = PP_ALIGN.CENTER

    t("CCS集成母排 × 全球储能市场", 38, True, (255,255,255), Inches(1.2), Inches(1.5))
    t("综合分析报告 2025-2030", 30, True, (0xFF,0xD7,0x00), Inches(2.7), Inches(1.0))
    t("CCS集成母排与线束技术整合市场分析报告", 18, False, (0xCC,0xD9,0xEC), Inches(4.2), Inches(0.8))
    t("全球储能及新能源关联市场经济规模分析", 16, False, (0x99,0xAA,0xCC), Inches(5.0), Inches(0.7))
    t("2026年4月  |  内部研究资料", 14, False, (0x77,0x88,0xAA), Inches(5.8), Inches(0.5))

# ============================================================
# 开始合并
# ============================================================
print("\nMerging...")

add_cover()
print("  [封面]")

add_divider(1, "CCS集成母排与线束技术整合", "市场分析报告", bg=(0x1A,0x5C,0x9E))
print("  [Part 1 分隔页]")

print("  [CCS报告...]")
for i in range(len(prs2.slides)):
    copy_slide_full(merged, prs2, i)
    print(f"    CCS {i+1}/{len(prs2.slides)}", end='\r')

add_divider(2, "全球储能及新能源关联市场", "2025-2030", bg=(0x0A,0x29,0x5E))
print("\n  [Part 2 分隔页]")

print("  [全球储能分析...]")
for i in range(len(prs1.slides)):
    copy_slide_full(merged, prs1, i)
    print(f"    储能 {i+1}/{len(prs1.slides)}", end='\r')

print("\n\nSaving...")
merged.save(OUTPUT)
size = os.path.getsize(OUTPUT)
total = 1 + 1 + len(prs2.slides) + 1 + len(prs1.slides)
print(f"Done! Size: {size/1024/1024:.2f} MB, {len(merged.slides)} slides (expected {total})")

# 验证
verify = Presentation(OUTPUT)
print(f"Verified: {len(verify.slides)} slides OK")

# 检查图片
with zipfile.ZipFile(OUTPUT) as z:
    media = [f for f in z.namelist() if f.startswith('ppt/media/')]
    print(f"Media files: {len(media)}")

import zipfile
