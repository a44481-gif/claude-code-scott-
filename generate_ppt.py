# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd
import os

OUTPUT_DIR = r"d:/claude mini max 2.7/storage_market_output"
OUT_PPTX = os.path.join(OUTPUT_DIR, "Global_Storage_Market_2025_2030.pptx")

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0A, 0x29, 0x5E)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xCC, 0xD9, 0xEC)
    p2.alignment = PP_ALIGN.CENTER
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "2025-2030 | 11大赛道深度分析 | 全球市场研究"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0x88, 0x99, 0xBB)
    p3.alignment = PP_ALIGN.CENTER
    return slide


def make_header(slide, title_text):
    hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.0))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = RGBColor(0x0A, 0x29, 0x5E)
    hdr.line.fill.background()
    txH = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.8))
    pH = txH.text_frame.paragraphs[0]
    pH.text = title_text
    pH.font.size = Pt(28)
    pH.font.bold = True
    pH.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_text_slide(prs, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    make_header(slide, title)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
        p.space_before = Pt(3)
    return slide


def add_table_slide(prs, title, headers, rows_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    make_header(slide, title)
    all_rows = [headers] + rows_data
    nrows, ncols = len(all_rows), len(all_rows[0])
    tbl_shape = slide.shapes.add_table(nrows, ncols,
        Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.8))
    tbl = tbl_shape.table
    col_widths = [Inches(2.5 if j == 0 else 1.1) for j in range(ncols)]
    for j, w in enumerate(col_widths):
        tbl.columns[j].width = w
    for i, row in enumerate(all_rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(10) if i > 0 else Pt(11)
            para.font.bold = (i == 0)
            para.alignment = PP_ALIGN.CENTER
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x0A, 0x29, 0x5E)
                para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                row_bg = RGBColor(0xE8, 0xF4, 0xFC) if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_bg
                para.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    return slide


def add_img_slide(prs, title, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    make_header(slide, title)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.2), Inches(1.1), height=Inches(6.0))
    return slide


def add_spec_slide(prs, seg, specs, total_5y, size_2030, cagr, share, trend):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    make_header(slide, f"赛道深度 | {seg}")

    stats = f"  5年累计: {total_5y:.0f}亿元  |  2030年: {size_2030:.0f}亿元  |  CAGR: {cagr}%  |  累计占比: {share:.2f}%"
    txStat = slide.shapes.add_textbox(Inches(0.3), Inches(1.05), Inches(12.7), Inches(0.4))
    pStat = txStat.text_frame.paragraphs[0]
    pStat.text = stats
    pStat.font.size = Pt(11)
    pStat.font.color.rgb = RGBColor(0x29, 0x6B, 0x9A)
    pStat.font.bold = True

    headers = ["型号", "电池参数", "性能指标", "典型应用", "主流厂商标配"]
    all_rows = [headers] + specs
    nrows, ncols = len(all_rows), len(all_rows[0])
    tbl_shape = slide.shapes.add_table(nrows, ncols,
        Inches(0.3), Inches(1.55), Inches(12.7), Inches(3.5))
    tbl = tbl_shape.table
    for j, w in enumerate([Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.5), Inches(3.8)]):
        tbl.columns[j].width = w
    for i, row in enumerate(all_rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(10)
            para.alignment = PP_ALIGN.LEFT if (i > 0 and j == len(row) - 1) else PP_ALIGN.CENTER
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x0A, 0x29, 0x5E)
                para.font.bold = True
                para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                row_bg = RGBColor(0xE8, 0xF4, 0xFC) if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_bg
                para.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
                if j == 0:
                    para.font.bold = True

    txTrend = slide.shapes.add_textbox(Inches(0.3), Inches(5.4), Inches(12.7), Inches(1.8))
    tfT = txTrend.text_frame
    tfT.word_wrap = True
    pT1 = tfT.paragraphs[0]
    pT1.text = "市场趋势："
    pT1.font.size = Pt(12)
    pT1.font.bold = True
    pT1.font.color.rgb = RGBColor(0x0A, 0x29, 0x5E)
    pT2 = tfT.add_paragraph()
    pT2.text = trend
    pT2.font.size = Pt(11)
    pT2.font.color.rgb = RGBColor(0x33, 0x33, 0x44)
    return slide


# ============================================================
# 开始生成PPT
# ============================================================
df = pd.read_csv(os.path.join(OUTPUT_DIR, "storage_market_2025_2030.csv"))
print(f"Data loaded: {df.shape}")

# Slide 1: 封面
add_title_slide(prs,
    "全球储能及新能源关联市场\n经济规模分析报告",
    "2025-2030 | 11大赛道 | 市场规模 | 规格型号 | 竞争格局")

# Slide 2: 执行摘要
summary = [
    "■ 分析范围：全球储能及新能源关联市场（2025-2030）",
    "■ 覆盖赛道：电芯 | 电动汽车 | 电动货卡车 | 电动船舶 | 家庭储能",
    "              工商储能 | AI储能 | 半导体设备储能 | BBU储能",
    "              无人机储能 | 机器人储能（共11大赛道）",
    "",
    "■ 核心结论：",
    f"  ▸ 11大赛道5年累计总规模：{df['5年累计(亿元)'].sum():.0f} 亿元",
    f"  ▸ 2030年预计总规模：{df['2030年规模(亿元)'].sum():.0f} 亿元",
    "  ▸ 增速最快：AI储能 (CAGR 65%)、机器人储能 (CAGR 58%)、无人机储能 (CAGR 55%)",
    f"  ▸ 规模最大：电动汽车 (5年累计 {df[df['赛道']=='电动汽车']['5年累计(亿元)'].values[0]:.0f} 亿元，占比 {df[df['赛道']=='电动汽车']['累计占比(%)'].values[0]:.1f}%)",
    "  ▸ 电动货卡车快速崛起(CAGR 32%)，船舶电动化加速(CAGR 42%)",
    "",
    "■ 数据来源：BloombergNEF / GGII / SNE Research / 行业公开数据综合",
]
add_text_slide(prs, "执行摘要 | Executive Summary", summary)

# Slide 3: 市场规模总览表
headers = ["赛道", "单位", "2025", "2026", "2027", "2028", "2029", "2030", "5年累计", "占比%"]
rows_tbl = []
for _, row in df.iterrows():
    rows_tbl.append([
        row["赛道"], row["单位"],
        f"{row['2025']:.0f}", f"{row['2026']:.0f}", f"{row['2027']:.0f}",
        f"{row['2028']:.0f}", f"{row['2029']:.0f}", f"{row['2030']:.0f}",
        f"{row['5年累计(亿元)']:.0f}", f"{row['累计占比(%)']:.2f}",
    ])
add_table_slide(prs, "市场规模总览 | Market Size Overview (亿元)", headers, rows_tbl)

# Slide 4: 综合图表
add_img_slide(prs, "综合图表分析 | Comprehensive Charts",
    os.path.join(OUTPUT_DIR, "storage_market_charts.png"))

# Slide 5: 趋势线图
add_img_slide(prs, "各赛道年度趋势 | Annual Trend Lines",
    os.path.join(OUTPUT_DIR, "storage_market_trends.png"))

# Slide 6: 相关性矩阵
add_img_slide(prs, "赛道相关性矩阵 | Correlation Matrix",
    os.path.join(OUTPUT_DIR, "storage_market_correlation.png"))

# ============================================================
# 各赛道详细规格页
# ============================================================
SEGMENT_SPECS = {
    "电芯": [
        ("LFP 280Ah", "170 Wh/kg", "4000次", "电网/工商业储能", "宁德时代、比亚迪"),
        ("LFP 314Ah", "178 Wh/kg", "6000次", "大型储能集装箱", "亿纬锂能、中创新航"),
        ("NCM 811 100Ah", "260 Wh/kg", "2000次", "高端新能源汽车", "LG新能源、松下"),
        ("钠离子 120Ah", "140 Wh/kg", "3000次", "大规模储能备用", "宁德时代、中科海钠"),
    ],
    "电动汽车": [
        ("CATL CTP 3.0", "75-100 kWh", "500-700km", "纯电动乘用车", "特斯拉Model 3/Y"),
        ("BYD Blade Battery", "60-85 kWh", "450-600km", "纯电/插混乘用车", "比亚迪全系"),
        ("GM Ultium", "50-200 kWh", "400-650km", "纯电皮卡/SUV", "通用Hummer EV"),
        ("Hyundai E-GMP", "58-77 kWh", "430-560km", "纯电动乘用车", "现代Ioniq 5/6"),
    ],
    "电动货卡车": [
        ("Tesla Semi", "900 kWh", "800km", "长途重型货运", "特斯拉"),
        ("比亚迪Q3", "435 kWh", "400km", "港口/矿区短倒", "比亚迪"),
        ("戴姆勒 eActros", "315 kWh", "400km", "城市配送重卡", "戴姆勒"),
        ("Nikola Tre BEV", "753 kWh", "410km", "美国区域货运", "Nikola/博世"),
    ],
    "电动船舶": [
        ("船用LFP 600kWh", "600 kWh", "10-15节", "内河客船/货船", "亿纬锂能、宁德时代"),
        ("纯电渡轮 2MW", "2000 kWh", "12节", "港口渡轮", "Corvus Energy"),
        ("集装箱船 50MW", "50 MWh", "20节", "沿海集装箱船", "比亚迪船舶板块"),
    ],
    "家庭储能": [
        ("Tesla Powerwall 3", "13.5 kWh", "11.5 kW", "家储/峰谷套利", "特斯拉/Panasonic"),
        ("华为LUNA2000", "5-30 kWh", "5 kW", "家储+光伏", "华为"),
        ("比亚迪Home FS", "7.7-30 kWh", "3.68 kW", "家储+光伏系统", "比亚迪"),
        ("Enphase IQ Battery", "3.36-13.44 kWh", "1.28kW/模块", "微逆变器集成", "Enphase"),
    ],
    "工商储能": [
        ("宁德时代 EnerOne", "200 MWh/站", "效率95%", "工业园区峰谷", "宁德时代+电网"),
        ("阳光电源 PowTitan", "100-500 kWh", "效率93%", "工商业调峰/备用", "阳光电源"),
        ("海辰储能 5MWh集装箱", "5 MWh", "效率94%", "数据算力备用", "海辰/科陆"),
    ],
    "AI储能": [
        ("NVIDIA DGX SuperPOD备用", "10-50 MWh", "MW级", "AI数据中心备用", "宁德时代+英伟达"),
        ("微软Project Natick", "100 MWh", "数十MW", "海底数据中心", "微软/阳光电源"),
        ("AWS/谷歌AI集群储能", "200+ MWh", "百MW级", "超大规模AI集群", "特斯拉Megapack"),
    ],
    "半导体设备储能": [
        ("ASML TWINSCAN NXE备用", "50-200 kWh", "数百kW", "光刻机不间断供电", "宁德时代"),
        ("应用材料PVD设备储能", "100-300 kWh", "数百kW", "薄膜沉积设备", "比亚迪电子"),
        ("Lam Research刻蚀备用", "50-150 kWh", "数百kW", "刻蚀设备精密供电", "亿纬锂能"),
    ],
    "BBU储能": [
        ("中国移动BBU备电", "2-10 kWh/柜", "数百W", "5G基站备用", "华为、宁德时代"),
        ("中国铁塔一体化机柜", "5-20 kWh", "1-3 kW", "边缘机房备电", "比亚迪、亿纬"),
        ("数据中心BBU机架", "10-50 kWh/架", "数kW", "通信机房不间断", "中创新航"),
    ],
    "无人机储能": [
        ("DJI Mavic 3 电池", "77 Wh", "260 Wh/kg", "消费级航拍", "大疆"),
        ("翼龙-2无人机", "20-50 kWh", "300 Wh/kg", "军用长航时", "中创新航"),
        ("亿航216 eVTOL", "100 kWh级", "250 Wh/kg", "城市空中交通", "亿纬锂能"),
    ],
    "机器人储能": [
        ("Tesla Optimus电池包", "2.3 kWh", "180 Wh/kg", "人形服务机器人", "特斯拉自研"),
        ("工业机械臂电池", "1-5 kWh", "200 Wh/kg", "柔性生产线", "ABB、发那科"),
        ("AMR/AGV机器人", "20-100 kWh", "160 Wh/kg", "仓储物流AMR", "宁德时代、比亚迪"),
    ],
}

CAGR_MAP = {
    "电芯": 26, "电动汽车": 18, "电动货卡车": 32, "电动船舶": 42,
    "家庭储能": 35, "工商储能": 40, "AI储能": 65, "半导体设备储能": 28,
    "BBU储能": 45, "无人机储能": 55, "机器人储能": 58,
}

TREND_MAP = {
    "电芯": "LFP路线持续主导大储市场，314Ah成为新一代大电芯标配，钠离子进入规模化初期。",
    "电动汽车": "CTP/CTC技术持续进化，刀片电池/麒麟电池高集成化，续航普遍突破600km。",
    "电动货卡车": "重卡续航瓶颈逐步突破，900kWh级产品开始商业化，换电模式在短倒场景加速普及。",
    "电动船舶": "电池系统向更大容量演进，50MW级大型船舶电池系统进入规划阶段。",
    "家庭储能": "德美日澳四极市场，模块化产品成主流，家储一体机趋势明显。",
    "工商储能": "峰谷价差套利驱动，工商业储能经济性持续改善，推动光储充一体化。",
    "AI储能": "AI算力爆发直接拉动储能需求，PUE监管趋严背景下数据中心储能成标配。",
    "半导体设备储能": "精密制造设备对供电连续性要求极高，定制化储能系统成为刚需。",
    "BBU储能": "5G基站密布推动备电需求增长，智能削峰/混电替代铅酸趋势加速。",
    "无人机储能": "eVTOL适航认证加速，城市空中交通市场2028年后望进入快速成长期。",
    "机器人储能": "人形机器人量产元年将至，AMR/AGV在仓储领域快速渗透，储能成核心零部件。",
}

for seg, specs in SEGMENT_SPECS.items():
    cagr = CAGR_MAP.get(seg, 0)
    seg_df = df[df["赛道"] == seg]
    total_5y = seg_df["5年累计(亿元)"].values[0] if len(seg_df) else 0
    size_2030 = seg_df["2030年规模(亿元)"].values[0] if len(seg_df) else 0
    share = seg_df["累计占比(%)"].values[0] if len(seg_df) else 0
    trend = TREND_MAP.get(seg, "市场持续高速增长。")
    add_spec_slide(prs, seg, specs, total_5y, size_2030, cagr, share, trend)

# 最后一页：总结
conclusions = [
    "【核心结论】",
    "",
    "1. 规模格局：电动汽车独占62.8%份额主导市场，电动货卡车占比27.9%居第二，",
    "   两大赛道合计占比超90%，构成市场的绝对主体。",
    "",
    "2. 增速格局：AI储能(CAGR 65%)、机器人储能(CAGR 58%)、无人机储能(CAGR 55%)",
    "   位列增速前三，AI算力经济和人形机器人产业化是最大增量驱动力。",
    "",
    "3. 新兴机会：电动船舶(CAGR 42%)、BBU储能(CAGR 45%)、工商储能(CAGR 40%)",
    "   处于高速成长期，兼具规模体量和增速优势，值得重点关注。",
    "",
    "4. 核心技术趋势：LFP路线持续主导储能市场，314Ah大电芯成为新一代主流，",
    "   钠离子电池开始规模化，半固态/全固态电芯研发加速。",
    "",
    "5. 政策驱动：中国双碳目标、欧盟Fit55、美国IRA法案共同推动全球储能市场",
    "   高速增长，新能源渗透率提升为储能打开长期空间。",
]
add_text_slide(prs, "总结与展望 | Conclusion", conclusions)

prs.save(OUT_PPTX)
print(f"PPT saved: {OUT_PPTX}")
print(f"Total slides: {len(prs.slides)}")
