# -*- coding: utf-8 -*-
"""
PPT合并脚本 v9 - 最终版：Base64内嵌图片，完整保留所有内容
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from copy import deepcopy
import zipfile, re, base64, os

SRC1 = r'd:/claude mini max 2.7/storage_market_output/Global_Storage_Market_2025_2030.pptx'
SRC2 = r'd:/claude mini max 2.7/storage_market_output/CCS_Global_Storage_2025_2030.pptx'
OUTPUT = r'd:/claude mini max 2.7/storage_market_output/CCS_Global_MERGED_FINAL.pptx'

print("Loading source files...")
prs1 = Presentation(SRC1)   # 全球储能 18 slides
prs2 = Presentation(SRC2)   # CCS报告  22 slides
print(f"SRC1 (全球储能): {len(prs1.slides)} slides")
print(f"SRC2 (CCS报告): {len(prs2.slides)} slides")

# ============================================================
# 1. 用 python-pptx 的 API 提取所有图片 (最可靠)
# ============================================================
def extract_all_images(source_prs):
    """用 python-pptx API 提取所有图片，返回 {(slide_partname, rId): (blob, mime, target_ref)}"""
    result = {}
    for slide in source_prs.slides:
        partname = slide.part.partname          # e.g. /ppt/slides/slide1.xml
        for rid, rel in slide.part.rels.items():
            if rel.reltype.endswith('image'):
                try:
                    img_part = rel.target_part
                    blob = img_part.blob
                    ct = img_part.content_type
                    result[(partname, rid)] = (blob, ct, rel.target_ref)
                except Exception:
                    pass
    return result

print("Extracting images via python-pptx API...")
img1 = extract_all_images(prs1)
img2 = extract_all_images(prs2)
print(f"Images from 全球储能: {len(img1)}")
print(f"Images from CCS报告: {len(img2)}")

# ============================================================
# 2. 构建合并简报
# ============================================================
merged = Presentation()
merged.slide_width = prs1.slide_width
merged.slide_height = prs1.slide_height

# ============================================================
# 3. 核心：复制幻灯片（含Base64图片）
# ============================================================
def find_cSld(tree):
    for el in tree.iter():
        if el.tag == qn('p:cSld'):
            return el
    return None

def resolve_target_path(rel_path):
    """解析 ../media/xxx 相对路径 -> ppt/media/xxx"""
    parts = rel_path.split('/')
    resolved = []
    for p in parts:
        if p == '..':
            if resolved: resolved.pop()
        elif p and p != '.':
            resolved.append(p)
    return '/'.join(resolved)

def copy_slide_with_base64_images(target_prs, source_prs, slide_idx, img_map):
    src_slide = source_prs.slides[slide_idx]
    src_tree = src_slide.part._element
    src_partname = src_slide.part.partname

    blank_layout = target_prs.slide_layouts[6]
    new_slide = target_prs.slides.add_slide(blank_layout)
    new_tree = new_slide.part._element

    cSld_src = find_cSld(src_tree)
    cSld_new = find_cSld(new_tree)

    if cSld_src is not None and cSld_new is not None:
        # 背景
        bg_src = cSld_src.find(qn('p:bg'))
        bg_new = cSld_new.find(qn('p:bg'))
        if bg_new is not None:
            cSld_new.remove(bg_new)
        if bg_src is not None:
            cSld_new.insert(0, deepcopy(bg_src))

        # 形状树
        spTree_src = cSld_src.find(qn('p:spTree'))
        spTree_new = cSld_new.find(qn('p:spTree'))

        if spTree_new is not None:
            for child in list(spTree_new):
                if child.tag not in (qn('p:nvGrpSpPr'), qn('p:grpSpPr')):
                    spTree_new.remove(child)

        if spTree_src is not None:
            for child in spTree_src:
                if child.tag == qn('p:nvGrpSpPr') or child.tag == qn('p:grpSpPr'):
                    spTree_new.append(deepcopy(child))
                else:
                    copied = deepcopy(child)

                    # 处理图片: 将 r:embed blip 替换为 base64 内嵌
                    for blip in copied.iter(qn('a:blip')):
                        embed = blip.get(qn('r:embed'))
                        if embed:
                            key = (src_partname, embed)
                            if key in img_map:
                                blob, mime, _ = img_map[key]
                                b64 = base64.b64encode(blob).decode('ascii')
                                data_uri = f'data:{mime};base64,{b64}'
                                # 清除 r:embed，添加 xlink:href (base64内嵌)
                                blip.attrib.clear()
                                blip.set('{http://www.w3.org/1999/xlink}href', data_uri)
                                # 添加 cstate="none"
                                parent = blip.getparent() if hasattr(blip, 'getparent') else None
                                if parent is not None:
                                    cstate = parent.find(qn('a:cstate'))
                                    if cstate is None:
                                        cstate = parent.makeelement(qn('a:cstate'), {})
                                        parent.append(cstate)
                                    cstate.set('val', 'none')
                                # 处理拉伸
                                stretch = parent.find(qn('a:stretch')) if parent is not None else None
                                if stretch is None and parent is not None:
                                    stretch = parent.makeelement(qn('a:stretch'), {})
                                    parent.append(stretch)
                                fillrect = stretch.makeelement(qn('a:fillRect')) if stretch is not None else None

                    spTree_new.append(copied)

    return new_slide

# ============================================================
# 4. 辅助页面
# ============================================================
def add_divider(part_num, line1, line2="", bg=(0x0A,0x29,0x5E)):
    slide = merged.slides.add_slide(merged.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*bg)
    def t(text, size, bold, rgb, y, h):
        tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.3), h)
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*rgb)
        p.alignment = PP_ALIGN.CENTER
    t(f"PART {part_num}", 20, True, (0xFF,0xD7,0x00), Inches(1.8), Inches(0.8))
    t(line1, 36, True, (255,255,255), Inches(2.6), Inches(1.2))
    if line2:
        t(line2, 26, True, (255,255,255), Inches(3.7), Inches(1.0))

def add_cover():
    slide = merged.slides.add_slide(merged.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x0A,0x29,0x5E)
    def t(text, size, bold, rgb, y, h):
        tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.3), h)
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*rgb)
        p.alignment = PP_ALIGN.CENTER
    t("CCS集成母排 × 全球储能市场", 38, True, (255,255,255), Inches(1.2), Inches(1.5))
    t("综合分析报告 2025-2030", 30, True, (0xFF,0xD7,0x00), Inches(2.7), Inches(1.0))
    t("CCS集成母排与线束技术整合市场分析报告", 18, False, (0xCC,0xD9,0xEC), Inches(4.2), Inches(0.8))
    t("全球储能及新能源关联市场经济规模分析", 16, False, (0x99,0xAA,0xCC), Inches(5.0), Inches(0.7))
    t("2026年4月  |  内部研究资料", 14, False, (0x77,0x88,0xAA), Inches(5.8), Inches(0.5))

# ============================================================
# 5. 合并流程
# ============================================================
print("\nMerging...")

add_cover()
print("  [封面]")

add_divider(1, "CCS集成母排与线束技术整合", "市场分析报告", bg=(0x1A,0x5C,0x9E))
print("  [Part 1 分隔页]")

print("  [CCS报告...]")
for i in range(len(prs2.slides)):
    copy_slide_with_base64_images(merged, prs2, i, img2)
    print(f"    CCS {i+1}/{len(prs2.slides)}", end='\r')

add_divider(2, "全球储能及新能源关联市场", "2025-2030")
print("\n  [Part 2 分隔页]")

print("  [全球储能分析...]")
for i in range(len(prs1.slides)):
    copy_slide_with_base64_images(merged, prs1, i, img1)
    print(f"    储能 {i+1}/{len(prs1.slides)}", end='\r')

print("\n\nSaving...")
merged.save(OUTPUT)
size = os.path.getsize(OUTPUT)
total = 1 + 1 + len(prs2.slides) + 1 + len(prs1.slides)
print(f"Done! Size: {size/1024/1024:.2f} MB, {len(merged.slides)} slides (expected {total})")

verify = Presentation(OUTPUT)
print(f"Verified: {len(verify.slides)} slides OK")

# 检查base64内容
with zipfile.ZipFile(OUTPUT) as z:
    media = [f for f in z.namelist() if f.startswith('ppt/media/')]
    print(f"Media files: {len(media)}")
    # Check for base64 in slides
    slide_data = z.read('ppt/slides/slide24.xml').decode('utf-8', errors='ignore')
    has_base64 = 'data:image' in slide_data
    print(f"Slide 24 has base64 images: {has_base64}")
