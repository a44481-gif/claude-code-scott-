#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

BG_DARK = RGBColor(13, 17, 23)
BG_CARD = RGBColor(22, 27, 34)
BORDER_COLOR = RGBColor(48, 54, 61)
PRIMARY = RGBColor(30, 111, 239)
CYAN = RGBColor(0, 212, 255)
GREEN = RGBColor(63, 185, 80)
ORANGE = RGBColor(240, 136, 62)
WHITE = RGBColor(240, 246, 252)
GRAY = RGBColor(139, 148, 158)
LIGHT_GRAY = RGBColor(180, 185, 195)
PURPLE = RGBColor(155, 89, 182)

def sf(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def sl(shape, color=None, width=Pt(1)):
    shape.line.width = width
    if color:
        shape.line.color.rgb = color
    else:
        shape.line.fill.background()

def tb(slide, l, t, w, h, text, sz=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, fn="Microsoft YaHei"):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = fn
    p.alignment = align
    return txb

def rect(slide, l, t, w, h, fill=None, lc=None, lw=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill: sf(s, fill)
    if lc: sl(s, lc, lw)
    else: s.line.fill.background()
    return s

def bg(slide):
    rect(slide, 0, 0, Inches(13.33), Inches(7.5), fill=BG_DARK)

def sn(slide, n):
    tb(slide, Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3), f"{n}/25", sz=9, color=GRAY, align=PP_ALIGN.RIGHT)

def ft(slide):
    tb(slide, Inches(0.3), Inches(7.1), Inches(5), Inches(0.3), "PD Adapter Market Analysis | 2025-2030", sz=8, color=GRAY)

def hdr(slide, title, sub=None):
    rect(slide, 0, 0, Inches(13.33), Inches(0.06), fill=PRIMARY)
    tb(slide, Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.55), title, sz=28, bold=True, color=WHITE)
    if sub:
        tb(slide, Inches(0.4), Inches(0.68), Inches(12.5), Inches(0.35), sub, sz=12, color=GRAY)

# ==============================================================================
# SLIDE 1: Cover
# ==============================================================================
def s01(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, Inches(13.33), Inches(0.08), fill=PRIMARY)
    for i in range(3):
        rect(s, Inches(0.5+i*4.2), Inches(2.1), Inches(3.5), Inches(0.015), fill=PRIMARY)
    tb(s, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.1),
       "PD 大電流適配器集成技術應用市場分析", sz=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.55),
       "Global PD High-Current Adapter Integration Technology Market Analysis 2025-2030",
       sz=16, color=GRAY, align=PP_ALIGN.CENTER)
    rect(s, Inches(2.3), Inches(4.4), Inches(8.7), Inches(1.0), fill=BG_CARD, lc=BORDER_COLOR)
    badges = [("$142億 → $287億", PRIMARY), ("CAGR 15.2%", CYAN), ("2025-2030", GREEN)]
    for i, (txt, col) in enumerate(badges):
        bx = Inches(2.5+i*2.9)
        rect(s, bx, Inches(4.55), Inches(2.7), Inches(0.72), fill=col)
        tb(s, bx, Inches(4.67), Inches(2.7), Inches(0.48), txt, sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, Inches(4.5), Inches(6.55), Inches(4.3), Inches(0.38),
       "市場研究報告旗艦版 | 2025年4月", sz=11, color=GRAY, align=PP_ALIGN.CENTER)
    sn(s, 1)
    return s

# ==============================================================================
# SLIDE 2: Disclaimer
# ==============================================================================
def s02(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    hdr(s, "免責聲明與研究方法論")
    items = [
        ("數據來源", "本報告數據來源包括：行業協會公開報告、上市公司年報及季報、主要OEM/ODM廠商財務數據、海關進出口統計、第三方調研機構數據庫、以及一級供應商訪談記錄。"),
        ("估算說明", "市場規模採用多種方法交叉驗證：額值法（出貨量×均價）、供應鏈調研法、以及終端應用拆解法的綜合平均。區域市場數據經過海關和產業調研雙重校正。"),
        ("研究方法", "本報告採用PEST宏觀分析、波特五力模型、價值鏈分析、及GaN/Si功率半導體滲透率S曲線模型相結合的綜合研究框架。"),
        ("核心假設", "研究假設包括：全球GDP2025-2030年均增長4.5%、消費電子景氣度溫和復甦、中國廠商持續價格競爭常態化、GaN成本年均下降12-15%等關鍵前提。"),
    ]
    cols = [PRIMARY, CYAN, GREEN, ORANGE]
    for i, (title, body) in enumerate(items):
        top = Inches(1.15)+i*Inches(1.52)
        rect(s, Inches(0.4), top, Inches(12.5), Inches(1.35), fill=BG_CARD, lc=BORDER_COLOR)
        rect(s, Inches(0.4), top, Inches(0.12), Inches(1.35), fill=cols[i])
        tb(s, Inches(0.65), top+Inches(0.08), Inches(12.0), Inches(0.32), title, sz=14, bold=True, color=WHITE)
        tb(s, Inches(0.65), top+Inches(0.42), Inches(12.0), Inches(0.8), body, sz=10.5, color=LIGHT_GRAY)
    ft(s); sn(s, 2); return s

# ==============================================================================
# SLIDE 3: Executive Summary
# ==============================================================================
def s03(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    hdr(s, "核心發現摘要")
    metrics = [
        ("2025 市場規模", "$142億", PRIMARY),
        ("2030 市場規模", "$287億", CYAN),
        ("複合年增長率", "15.2%", GREEN),
        ("2030 集成滲透率", "92%", ORANGE),
    ]
    cw = Inches(2.95)
    for i, (lbl, val, col) in enumerate(metrics):
        lft = Inches(0.4)+i*(cw+Inches(0.1))
        rect(s, lft, Inches(0.82), cw, Inches(1.72), fill=BG_CARD, lc=BORDER_COLOR)
        rect(s, lft, Inches(0.82), cw, Inches(0.42), fill=col)
        tb(s, lft, Inches(0.88), cw, Inches(0.32), lbl, sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(s, lft, Inches(1.32), cw, Inches(1.05), val, sz=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, Inches(0.4), Inches(2.72), Inches(12.5), Inches(0.38), "三大核心投資機會", sz=15, bold=True, color=WHITE)
    opps = [
        ("1", "GaN單芯片全集成方案", "Navitas、PI、南芯等龍頭搶佔的下一代PD適配器技術高地，預計2027年進入爆發期"),
        ("2", "中國PD SoC芯片出海", "南芯、智融等廠商的140W以上SoC方案性價比全球領先，有望複製韋爾、博裝微的出海路徑"),
        ("3", "汽車200-240W市場", "800V EV快充和AI車載設備拉動汽車PD需求，2026-2028年CAGR達38%"),
    ]
    for i, (num, title, desc) in enumerate(opps):
        top = Inches(3.2)+i*Inches(1.35)
        rect(s, Inches(0.4), top, Inches(12.5), Inches(1.15), fill=BG_CARD, lc=BORDER_COLOR)
        rect(s, Inches(0.52), top+Inches(0.22), Inches(0.58), Inches(0.58), fill=PRIMARY)
        tb(s, Inches(0.52), top+Inches(0.3), Inches(0.58), Inches(0.45), num, sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(s, Inches(1.25), top+Inches(0.08), Inches(11.5), Inches(0.35), title, sz=13, bold=True, color=WHITE)
        tb(s, Inches(1.25), top+Inches(0.48), Inches(11.5), Inches(0.55), desc, sz=10.5, color=LIGHT_GRAY)
    ft(s); sn(s, 3); return s

# ==============================================================================
# SLIDE 4: Market Overview
# ==============================================================================
def s04(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    hdr(s, "市場全景 2025→2030")
    data = [("2025",14.2),("2026",17.1),("2027",20.6),("2028",24.0),("2029",26.8),("2030",28.7)]
    max_v = 30.0
    n = len(data)
    bw = (Inches(11.0)) / n
    g = Inches(0.04)
    bh_max = Inches(3.0)
    for i, (lbl, val) in enumerate(data):
        bh = (val/max_v)*bh_max
        bl = Inches(0.6)+i*bw
        rect(s, bl, Inches(3.5)-bh, bw-g, bh, fill=PRIMARY)
        tb(s, bl, Inches(3.55)-bh, bw-g, Inches(0.25), f"${val}B", sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(s, bl, Inches(3.5), bw-g, Inches(0.25), lbl, sz=10, color=GRAY, align=PP_ALIGN.CENTER)
    tb(s, Inches(0.4), Inches(3.95), Inches(12.5), Inches(0.35), "五大增長驅動力", sz=15, bold=True, color=WHITE)
    drivers = [
        ("AI PC/筆電", "140-240W PD快充需求爆發，2025-2028 CAGR達12%", PRIMARY),
        ("旗艦手機", "100-240W快充全面普及，摺疊手機拉動需求", CYAN),
        ("電動汽車", "800V平台車載PD適配器標配，200W+需求涌現", GREEN),
        ("GaN成本下降", "GaN單芯片集成方案成本接近離散方案，滲透率加速", ORANGE),
        ("快充協議統一", "USB PD3.1 EPR 240W標準全球採用率提升", PURPLE),
    ]
    for i, (title, desc, col) in enumerate(drivers):
        top = Inches(4.4)+i*Inches(0.58)
        rect(s, Inches(0.4), top, Inches(12.5), Inches(0.52), fill=BG_CARD, lc=BORDER_COLOR)
        rect(s, Inches(0.4), top, Inches(0.1), Inches(0.52), fill=col)
        tb(s, Inches(0.62), top+Inches(0.04), Inches(2.0), Inches(0.28), title, sz=11, bold=True, color=WHITE)
        tb(s, Inches(2.7), top+Inches(0.04), Inches(10.0), Inches(0.42), desc, sz=10, color=LIGHT_GRAY)
    ft(s); sn(s, 4); return s

# ==============================================================================
# SLIDE 5: Global Market Size
# ==============================================================================
def s05(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    hdr(s, "全球市場規模 2025-2030", "Unit: USD Billions | CAGR 15.2%")
    data = [("2025",14.2),("2026",17.1),("2027",20.6),("2028",24.0),("2029",26.8),("2030",28.7)]
    max_v = 30.0
    n = len(data)
    bw = (Inches(11.0)) / n
    g = Inches(0.05)
    bh_max = Inches(3.2)
    for i, (lbl, val) in enumerate(data):
        bh = (val/max_v)*bh_max
        bl = Inches(0.7)+i*bw
        rect(s, bl, Inches(3.55)-bh, bw-g, bh, fill=PRIMARY)
        tb(s, bl, Inches(3.6)-bh, bw-g, Inches(0.25), f"${val}B", sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(s, bl, Inches(3.5), bw-g, Inches(0.25), lbl, sz=10, color=GRAY, align=PP_ALIGN.CENTER)
    rect(s, Inches(0.4), Inches(3.9), Inches(12.5), Inches(0.05), fill=CYAN)
    milestones = [
        ("2026", "突破$170億", CYAN),
        ("2028", "出貨量>7.5億台", GREEN),
        ("2030", "集成滲透率達92%", ORANGE),
    ]
    for i, (yr, ms, col) in enumerate(milestones):
        lft = Inches(0.4)+i*Inches(4.2)
        rect(s, lft, Inches(4.05), Inches(3.8), Inches(0.55), fill=BG_CARD, lc=BORDER_COLOR)
        rect(s, lft, Inches(4.05), Inches(0.1), Inches(0.55), fill=col)
        tb(s, lft+Inches(0.15), Inches(4.1), Inches(1.0), Inches(0.25), yr, sz=12, bold=True, color=col)
        tb(s, lft+Inches(1.1), Inches(4.12), Inches(2.5), Inches(0.35), ms, sz=11, color=WHITE)
    rect(s, Inches(9.5), Inches(0.95), Inches(3.3), Inches(0.72), fill=BG_CARD, lc=CYAN, lw=Pt(2))
    tb(s, Inches(9.6), Inches(1.0), Inches(3.1), Inches(0.28), "6年複合增長率", sz=11, color=GRAY)
    tb(s, Inches(9.6), Inches(1.25), Inches(3.1), Inches(0.35), "CAGR = 15.2%", sz=18, bold=True, color=CYAN)
    ft(s); sn(s, 5); return s

# ==============================================================================
# SLIDE 6: Power Structure
# ==============================================================================
def s06(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    hdr(s, "功率結構演變 2025 vs 2030", "PD Adapter Market Power Segment Shift")
    tb(s, Inches(0.4), Inches(1.1), Inches(5.5), Inches(0.4), "2025 市場份額", sz=16, bold=True, color=WHITE)
    p2025 = [("65W", 52, PRIMARY),("100W", 28, CYAN),("140W", 12, GREEN),("240W", 8, ORANGE)]
    y = Inches(1.55)
    for lbl, pct, col in p2025:
        rect(s, Inches(0.5), y, Inches(0.25), Inches(0.3), fill=col)
        tb(s, Inches(0.82), y-Inches(0.02), Inches(3.5), Inches(0.3), f"{lbl}: {pct}%", sz=13, color=WHITE)
        y += Inches(0.42)
    total_w = Inches(5.0)
    x = Inches(0.5)
    for lbl, pct, col in p2025:
        w = total_w * pct / 100
        rect(s, x, Inches(3.4), w, Inches(0.5), fill=col)
        if pct >= 10:
            tb(s, x+Inches(0.05), Inches(3.48), w-Inches(0.1), Inches(0.35), f"{pct}%", sz=10, bold=True, color=WHITE)
        x += w
    tb(s, Inches(7.0), Inches(1.1), Inches(5.5), Inches(0.4), "2030 市場份額", sz=16, bold=True, color=WHITE)
    p2030 = [("65W", 35, PRIMARY),("100W", 30, CYAN),("140W", 20, GREEN),("240W", 15, ORANGE)]
    y = Inches(1.55)
    for lbl, pct, col in p2030:
        rect(s, Inches(7.1), y, Inches(0.25), Inches(0.3), fill=col)
        tb(s, Inches(7.42), y-Inches(0.02), Inches(3.5), Inches(0.3), f"{lbl}: {pct}%", sz=13, color=WHITE)
        y += Inches(0.42)
    x = Inches(7.1)
    for lbl, pct, col in p2030:
        w = total_w * pct / 100
        rect(s, x, Inches(3.4), w, Inches(0.5), fill=col)
        if pct >= 10:
            tb(s, x+Inches(0.05), Inches(3.48), w-Inches(0.1), Inches(0.35), f"{pct}%", sz=10, bold=True, color=WHITE)
        x += w
    tb(s, Inches(0.4), Inches(4.15), Inches(12.5), Inches(0.35), "各功率段CAGR對比", sz=15, bold=True, color=WHITE)
    cagrs = [("65W","8.5%",Inches(0.4)),("100W","15.2%",Inches(3.2)),("140W","22.0%",Inches(6.0)),("240W","28.5%",Inches(8.8))]
    for lbl, cagr, lft in cagrs:
        rect(s, lft, Inches(4.6), Inches(2.6), Inches(1.3), fill=BG_CARD, lc=BORDER_COLOR)
        rect(s, lft, Inches(4.6), Inches(2.6), Inches(0.38), fill=PRIMARY)
        tb(s, lft, Inches(4.65), Inches(2.6), Inches(0.3), lbl, sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(s, lft, Inches(5.1), Inches(2.6), Inches(0.65), cagr, sz=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    tb(s, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.5),
       "洞察：240W功率段增長最快（CAGR 28.5%），主要由AI PC遊戲本和EV車載適配器拉動；65W份額收縮但總量仍居首位",
       sz=11, color=LIGHT_GRAY)
    ft(s); sn(s, 6); return s
