#!/usr/bin/env python3
"""
Generate PPTX Business Presentation for PD High-Current Adapter Market Analysis
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BLUE = RGBColor(0x1a, 0x1a, 0x2e)
MED_BLUE = RGBColor(0x0f, 0x34, 0x60)
ACCENT_RED = RGBColor(0xe9, 0x45, 0x60)
ACCENT_ORANGE = RGBColor(0xf5, 0xa6, 0x23)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf0, 0xf2, 0xf5)
GRAY = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x0f, 0x9b, 0x6e)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BLUE)

    for i in range(4):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0), Inches(5.5 + i * 0.12),
                                      prs.slide_width, Inches(0.12))
        rect.fill.solid()
        rect.fill.fore_color.rgb = MED_BLUE if i < 3 else ACCENT_RED
        rect.line.fill.background()

    txBox0 = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
    tf = txBox0.text_frame
    p = tf.paragraphs[0]
    p.text = "MARKET ANALYSIS REPORT  |  2026年4月"
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(2.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in ["全球PD大電流適配器", "集成技術市場分析", "2025-2030"]:
        color = ACCENT_RED if line == "2025-2030" else WHITE
        if tf.paragraphs[0].text:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(44) if line != "2025-2030" else Pt(38)
        p.font.bold = True
        p.font.color.rgb = color

    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(9), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "芯片級集成  ·  GaN/SiC合封  ·  系統級SoC  ·  商業價值落地"
    p2.font.size = Pt(14)
    p2.font.color.rgb = LIGHT_GRAY

    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.4))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "報告版本：v1.0  |  機密文件  |  僅供內部參考"
    p3.font.size = Pt(10)
    p3.font.color.rgb = GRAY

    return slide

def add_kpi_slide(prs, title, kpis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GRAY)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    header.fill.solid(); header.fill.fore_color.rgb = DARK_BLUE; header.line.fill.background()

    txTitle = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.7))
    tf = txTitle.text_frame
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = WHITE

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), prs.slide_width, Inches(0.06))
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT_RED; accent.line.fill.background()

    n = len(kpis)
    card_w = Inches(2.8); card_h = Inches(2.0)
    gap = Inches(0.35)
    total_w = n * card_w + (n - 1) * gap
    start_x = (prs.slide_width - total_w) / 2
    start_y = Inches(1.5)
    colors = [ACCENT_RED, MED_BLUE, ACCENT_ORANGE, DARK_BLUE, GREEN, RGBColor(0x6b,0x4c,0x9a)]

    for i, (label, value, sub) in enumerate(kpis):
        cx = start_x + i * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, start_y, card_w, card_h)
        card.fill.solid(); card.fill.fore_color.rgb = colors[i % len(colors)]; card.line.fill.background()

        tv = slide.shapes.add_textbox(cx, start_y + Inches(0.25), card_w, Inches(0.8))
        pv = tv.text_frame.paragraphs[0]; pv.text = value; pv.font.size = Pt(30); pv.font.bold = True; pv.font.color.rgb = WHITE; pv.alignment = PP_ALIGN.CENTER

        tl = slide.shapes.add_textbox(cx, start_y + Inches(1.1), card_w, Inches(0.5))
        pl = tl.text_frame.paragraphs[0]; pl.text = label; pl.font.size = Pt(12); pl.font.color.rgb = WHITE; pl.alignment = PP_ALIGN.CENTER

        if sub:
            ts = slide.shapes.add_textbox(cx, start_y + Inches(1.5), card_w, Inches(0.4))
            ps = ts.text_frame.paragraphs[0]; ps.text = sub; ps.font.size = Pt(9); ps.font.color.rgb = RGBColor(0xdd,0xdd,0xdd); ps.alignment = PP_ALIGN.CENTER

    return slide

def add_table_slide(prs, title, headers, rows, footer=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GRAY)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    header.fill.solid(); header.fill.fore_color.rgb = DARK_BLUE; header.line.fill.background()

    txTitle = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.7))
    tf = txTitle.text_frame
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = WHITE

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), prs.slide_width, Inches(0.06))
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT_RED; accent.line.fill.background()

    cols = len(headers)
    rcount = len(rows) + 1
    table_h = min(5.0, 0.38 * rcount + 0.1)
    tbl = slide.shapes.add_table(rcount, cols, Inches(0.4), Inches(1.35), Inches(9.2), Inches(table_h)).table

    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE
        cp = cell.text_frame.paragraphs[0]; cp.font.bold = True; cp.font.size = Pt(10); cp.font.color.rgb = WHITE; cp.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else RGBColor(0xf5,0xf7,0xfa)
            cp = cell.text_frame.paragraphs[0]; cp.font.size = Pt(9); cp.font.color.rgb = DARK_BLUE; cp.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    if footer:
        txF = slide.shapes.add_textbox(Inches(0.4), Inches(6.8), Inches(9), Inches(0.4))
        pf = txF.text_frame.paragraphs[0]; pf.text = footer; pf.font.size = Pt(9); pf.font.color.rgb = GRAY; pf.alignment = PP_ALIGN.RIGHT

    return slide

def add_content_slide(prs, title, bullets, footer=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GRAY)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    header.fill.solid(); header.fill.fore_color.rgb = DARK_BLUE; header.line.fill.background()

    txTitle = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.7))
    tf = txTitle.text_frame
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = WHITE

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), prs.slide_width, Inches(0.06))
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT_RED; accent.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5))
    tf2 = txBox.text_frame; tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = "  " + bullet
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(10)

    if footer:
        txF = slide.shapes.add_textbox(Inches(0.4), Inches(6.8), Inches(9), Inches(0.4))
        pf = txF.text_frame.paragraphs[0]; pf.text = footer; pf.font.size = Pt(9); pf.font.color.rgb = GRAY; pf.alignment = PP_ALIGN.RIGHT

    return slide

def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, MED_BLUE)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.5), prs.slide_width, Inches(0.08))
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT_RED; accent.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame; tf.word_wrap = True
    lines = title.split("\n")
    for j, line in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

    return slide


def main():
    print("Generating PD Adapter Market Analysis PPT...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: Cover
    add_cover_slide(prs)

    # SLIDE 2: KPIs
    add_kpi_slide(prs, "執行摘要 — 核心指標", [
        ("2025市場規模", "$142億", "USD Billion"),
        ("2030市場規模", "$287億", "CAGR 15.2%"),
        ("累計市場", "$1,314億", "2025-2030"),
        ("GaN滲透率", "90%", "2030年"),
        ("最快增速", "240W段", "CAGR 32.4%"),
    ])

    # SLIDE 3: Executive Summary
    add_content_slide(prs, "執行摘要 — 核心結論", [
        "GaN集成技術是最大增長驅動力：滲透率 45%(2025) → 92%(2030)，CAGR 15.2%",
        "140W+功率段增速最快 (CAGR 28.7%)，AI PC/遊戲本為核心增長引擎",
        "中國廠商快速崛起：南芯/智融/東科已具備挑戰國際大廠能力",
        "車載/新能源為最高毛利細分市場 (CAGR 38%)，AEC-Q100認證形成天然壁壘",
        "芯片廠商拿走價值鏈50%+毛利，為整條價值鏈最肥環節",
        "投資建議：超配GaN集成芯片廠商 + 140W以上功率段 + 中國SoC廠商IPO前布局",
    ])

    # SLIDE 4: 市場邊界定義
    add_table_slide(prs, "市場邊界定義 — 功率區間", [
        "功率段", "標準協議", "電流", "典型應用場景"
    ], [
        ["65W", "PD3.0 PPS", "3A / 5A", "主流筆電 / 智能手機"],
        ["100W", "PD3.0 EPR", "5A", "旗艦筆電"],
        ["140W", "PD3.1 EPR", "5A", "AI PC / 遊戲本"],
        ["240W", "PD3.1 EPR", "7A", "工作站 / 車載 / 工業"],
    ])

    # SLIDE 5: 市場規模 KPIs
    add_kpi_slide(prs, "全球市場規模 (2025-2030)", [
        ("2025年", "$14.2B", "485M 出貨量"),
        ("2030年", "$28.7B", "978M 出貨量"),
        ("CAGR", "15.2%", "6年複合增長"),
        ("累計總量", "$1,314億", "USD"),
    ])

    # SLIDE 6: 市場規模詳情
    add_table_slide(prs, "全球市場規模詳情 (2025-2030)", [
        "年份", "市場規模", "出貨量", "集成滲透率", "同比增速"
    ], [
        ["2025", "$14.2B", "485M件", "45%", "-"],
        ["2026", "$17.1B", "572M件", "58%", "+20.4%"],
        ["2027", "$20.6B", "672M件", "72%", "+20.5%"],
        ["2028", "$24.0B", "778M件", "82%", "+16.5%"],
        ["2029", "$26.8B", "878M件", "88%", "+11.7%"],
        ["2030", "$28.7B", "978M件", "92%", "+7.1%"],
    ])

    # SLIDE 7: 功率結構
    add_table_slide(prs, "功率結構演變 (2025 vs 2030)", [
        "功率段", "2025份額", "2030份額", "2025市場", "2030市場", "CAGR"
    ], [
        ["65W", "52%", "35%", "$7.38B", "$10.05B", "8.2%"],
        ["100W", "28%", "30%", "$3.98B", "$8.61B", "16.5%"],
        ["140W", "12%", "20%", "$1.70B", "$5.74B", "28.7%"],
        ["240W", "8%", "15%", "$1.14B", "$4.31B", "32.4%"],
    ])

    # SLIDE 8: 區域市場
    add_table_slide(prs, "區域市場格局 (2025 vs 2030)", [
        "地區", "2025份額", "2030份額", "CAGR", "關鍵驅動因素"
    ], [
        ["中國", "48%", "45%", "13.8%", "UFCS強制認證"],
        ["北美", "20%", "19%", "14.2%", "AI設備普及"],
        ["歐洲", "17%", "20%", "18.5%", "USB-C強制令執法"],
        ["東南亞", "8%", "10%", "22.0%", "製造業轉移、收入提升"],
        ["日韓", "5%", "4%", "12.5%", "成熟市場"],
    ])

    # SLIDE 9: 技術滲透率
    add_table_slide(prs, "集成技術滲透率路線圖 (2025-2030)", [
        "年份", "GaN集成", "SiC採用", "單芯片方案", "數字控制", "AI自適應"
    ], [
        ["2025", "45%", "8%", "25%", "40%", "5%"],
        ["2026", "55%", "12%", "35%", "50%", "12%"],
        ["2027", "65%", "18%", "48%", "62%", "25%"],
        ["2028", "75%", "25%", "60%", "75%", "40%"],
        ["2029", "82%", "32%", "72%", "85%", "55%"],
        ["2030", "90%", "40%", "85%", "95%", "70%"],
    ])

    # SLIDE 10: 技術方案對比
    add_table_slide(prs, "芯片級集成方案對比", [
        "架構", "代表廠商", "效率", "功率密度", "BOM指數", "成熟度"
    ], [
        ["GaN集成 PowiGaN", "Power Integrations", "96%", "25 W/in\u00b3", "0.75", "成熟"],
        ["GaN HEMT+控制器", "Navitas", "95%", "18 W/in\u00b3", "0.85", "成長中"],
        ["PFC+LLC單芯片", "PI, ON Semi", "95%", "20 W/in\u00b3", "0.70", "成長中"],
        ["SiC MOSFET", "Wolfspeed", "97%", "22 W/in\u00b3", "1.20", "早期"],
        ["協議SoC+GaN", "南芯, Richtek", "94%", "19 W/in\u00b3", "0.72", "成長中"],
    ])

    # SLIDE 11: 關鍵芯片
    add_content_slide(prs, "關鍵芯片方案深度分析", [
        "Power Integrations InnoSwitch4 (PowiGaN)：功率最高220W, 效率>95%, 旗艦筆電標配，代表客戶Dell/HP/ASUS ROG",
        "Navitas GaNSense：1000+ GaN專利, 首創感測保護, 手機/筆電高端市場，代表客戶Samsung/小米/OPPO",
        "東科 DK8607 (ACF單芯片)：最高120W, 性價比最強, 中國消費市場，代表客戶Anker/倍思/品勝",
        "南芯 SC9712 (協議+DC-DC一體化)：最高140W, PD3.1+UFCS雙協議, 多口充电器, 已達國際領先水平",
        "TI UCC29950 (全數字控制)：最高240W, 最智能控制, 車載/醫療/工作站",
    ])

    # SLIDE 12: 技術路線圖
    add_table_slide(prs, "技術路線圖 (2025-2030)", [
        "年份", "集成度", "開關頻率", "功率密度", "協議標準"
    ], [
        ["2025", "PFC+LLC 雙芯片", "200-500 kHz", "20-25 W/in\u00b3", "PD3.1 + UFCS興起"],
        ["2026", "GaN合封主流化", "500-1000 kHz", "25-32 W/in\u00b3", "PD3.1 + UFCS強制"],
        ["2027", "單芯片GaN SoC", "1000-2000 kHz", "32-40 W/in\u00b3", "多協議融合"],
        ["2028", "全數字單芯片", "2000-3000 kHz", "40-50 W/in\u00b3", "AI感知功率分配"],
        ["2029", "GaN+SiC混合", "3000-5000 kHz", "50-60 W/in\u00b3", "預測性功率管理"],
        ["2030", "單片GaN SoC", "5000-10000 kHz", "60-80 W/in\u00b3", "傳感器+控制一體化"],
    ])

    # SLIDE 13: Section - 競爭格局
    add_section_slide(prs, "第三章\n競爭格局")

    # SLIDE 14: 國際廠商
    add_table_slide(prs, "國際廠商競爭格局", [
        "廠商", "市場份額", "核心技術", "代表產品", "目標市場"
    ], [
        ["Power Integrations", "22%", "GaN/SiC合封", "InnoSwitch4", "旗艦/工業"],
        ["Infineon", "20%", "GaN+SiC雙平台", "CoolMOS", "工業/汽車"],
        ["Navitas", "18%", "GaN Power IC", "GaNFast", "高端充电器"],
        ["TI", "15%", "數字電源", "UCD系列", "廣泛市場"],
        ["ON Semi", "6%", "MOSFET", "Nexperia", "成本敏感"],
    ])

    # SLIDE 15: 中國廠商
    add_table_slide(prs, "中國廠商快速崛起", [
        "廠商", "市場份額", "核心技術", "代表產品", "目標市場"
    ], [
        ["南芯半導體", "8%", "協議+GaN SoC", "SC9712", "手機/筆電"],
        ["智融科技", "6%", "多口協議IC", "SW35xx", "充電寶/充电器"],
        ["東科半導體", "6%", "GaN控制,ACF", "DK8607", "消費適配器"],
        ["英集芯", "5%", "協議IC", "IP2188", "消費適配器"],
        ["杰華特", "4%", "數字LLC/PFC", "JW72xx", "工業"],
    ])

    # SLIDE 16: 終端應用
    add_table_slide(prs, "終端應用需求分析", [
        "應用場景", "功率(W)", "CAGR", "關鍵需求", "代表廠商"
    ], [
        ["AI PC/筆電", "140-240", "28.7%", "超薄,多口,96%+效率", "Apple,Dell,HP"],
        ["遊戲設備", "120-320", "25.0%", "高穩定性,低噪聲", "Sony,Microsoft,Valve"],
        ["旗艦手機", "100-150", "8.0%", "迷你化,UFCS+PD", "Samsung,小米,OPPO"],
        ["汽車IVI", "45-120", "40.0%", "AEC-Q,寬溫", "Tesla,比亞迪,NIO"],
        ["戶外儲能", "100-240", "30.0%", "便攜,大容量", "Jackery,EcoFlow"],
    ])

    # SLIDE 17: BOM成本
    add_table_slide(prs, "BOM成本與毛利分析", [
        "功率段", "離散BOM", "集成BOM", "節省比例", "零售價", "毛利"
    ], [
        ["25W-45W", "$5.20", "$4.10", "21.2%", "$12", "52%"],
        ["65W-100W", "$11.50", "$8.80", "23.5%", "$28", "56%"],
        ["100W-140W", "$18.20", "$14.50", "20.3%", "$48", "58%"],
        ["140W-240W", "$28.50", "$22.00", "22.8%", "$85", "62%"],
        ["240W以上", "$48.00", "$38.00", "20.8%", "$160", "68%"],
    ])

    # SLIDE 18: 挑戰與風險
    add_content_slide(prs, "挑戰與風險評估", [
        "GaN良率與成本：投資8寸GaN晶圓是關鍵制約，滲透速度放緩風險",
        "中國廠商價格戰：全行業毛利侵蝕，需聚焦差異化（車規/工業/特殊協議）",
        "地緣政治風險：中美科技脫鉤影響供應鏈，優先配置中國本土廠商",
        "專利壁壘：Navitas 1000+專利, PI 800+專利，形成強技術護城河",
        "熱管理瓶頸：100W+設計關鍵挑戰，先進封裝是核心解決方案",
        "認證週期：車規/醫療18-24個月認證，形成天然時間壁壘保護毛利",
    ])

    # SLIDE 19: 投資建議
    add_content_slide(prs, "投資建議", [
        "Tier 1 — 超配：Navitas Semiconductor (GaN集成最強, AI PC受益者, PS 8-12x)",
        "Tier 1 — 超配：南芯半導體 (中國SoC領導者, UFCS+PD3.1雙協議, Pre-IPO)",
        "Tier 2 — 配置：Power Integrations (穩定盈利, 旗艦定位, PE 25-30x)",
        "Tier 2 — 配置：智融科技 (多口市場領導者, Pre-IPO)",
        "Tier 3 — 長期：汽車PD芯片 (高壁壘, 高回報, AEC-Q100認證護城河)",
        "功率段選擇：優先140W+高增速段 (CAGR 28.7%)，避免65W紅海價格戰",
    ])

    # SLIDE 20: 商業建議
    add_content_slide(prs, "商業建議總結", [
        "1. 現在布局GaN集成芯片：技術窗口2025-2027，錯過再等5年",
        "2. 聚焦140W+功率段：CAGR 28.7%藍海市場，車載/工業，毛利率達62%+",
        "3. 中國市場戰略性優先：UFCS強制推行，本土芯片廠商爆發窗口",
        "4. 切入路徑：消費電子積累規模 → 多口桌面充电器提升毛利 → 汽車電子建立壁壘",
        "5. 風險對策：差異化避開價格戰、交叉授權繞過專利、中國廠商贴近供應鏈",
    ])

    # SLIDE 21: 法規標準
    add_table_slide(prs, "附錄 — 法規標準參考", [
        "標準", "應用場景", "版本", "關鍵要求"
    ], [
        ["USB PD 3.1 EPR", "通用充電", "240W", "全球主流標準"],
        ["UFCS 1.0", "中國市場", "240W", "強制認證"],
        ["IEC 60601", "醫療設備", "3rd Edition", "安全標準,2xMOPP"],
        ["ISO 26262", "汽車電子", "ASIL B/C", "功能安全"],
        ["AEC-Q100", "汽車電子", "Grade 1", "可靠性標準"],
        ["EU USB-C", "歐洲強制令", "2024/2026", "執法時間線"],
    ])

    # SLIDE 22: 術語表
    add_content_slide(prs, "附錄 — 專業術語表", [
        "GaN (氮化鎵)：寬禁帶半導體，功率應用下一代材料，效率比矽高50%，滲透率從45%升至90%",
        "ACF (有源鉗位反激)：開關電源拓撲結構，實現軟開關，東科DK8607為中國ACF最強代表",
        "SoC (系統級芯片)：高集成半導體方案，PFC+控制器+GaN一體化，南芯SC9712已達國際領先",
        "EPR (擴展功率範圍)：USB PD3.1高功率規範，支持240W，AI PC標配接口",
        "UFCS (融合快充)：中國統一快充標準，強制認證，國產芯片最大爆發機會窗口",
        "BOM (物料清單)：集成方案比離散方案節省BOM成本20-23%，毛利提升10-15個百分點",
    ])

    # SLIDE 23: Thank you
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BLUE)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.6), prs.slide_width, Inches(0.08))
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT_RED; accent.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "謝謝"; p.font.size = Pt(56); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph(); p2.text = "PD High-Current Adapter Market Analysis 2025-2030"; p2.font.size = Pt(16); p2.font.color.rgb = LIGHT_GRAY; p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph(); p3.text = "\n本報告僅供內部參考  |  數據截止2026年4月  |  投資決策需自行判斷"; p3.font.size = Pt(10); p3.font.color.rgb = GRAY; p3.alignment = PP_ALIGN.CENTER

    # Save
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(script_dir, "PD_Adapter_Market_PPT_2025_2030.pptx")
    prs.save(output_path)
    size = os.path.getsize(output_path)
    print(f"\nPPTX saved: {output_path}")
    print(f"File size: {size/1024:.1f} KB")
    print(f"Total slides: {len(prs.slides)}")
    return output_path


if __name__ == "__main__":
    main()
