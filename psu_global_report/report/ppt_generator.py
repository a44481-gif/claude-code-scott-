"""
PPT 報告生成器
使用 python-pptx 生成全球 PSU 市場每日分析報告
"""

import logging
import sys
from datetime import datetime, date
from collections import defaultdict
from pathlib import Path
from typing import Optional

# 確保 Windows 控制台支援 UTF-8 emoji
if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# 顏色主題
THEME = {
    "primary": RGBColor(0x1A, 0x73, 0xE8),      # Google Blue
    "secondary": RGBColor(0x0D, 0x47, 0xA1),    # Deep Blue
    "accent": RGBColor(0x00, 0x96, 0x88),       # Teal
    "warning": RGBColor(0xFF, 0x6D, 0x00),       # Orange
    "success": RGBColor(0x00, 0xC8, 0x53),       # Green
    "text": RGBColor(0x2C, 0x3E, 0x50),
    "light_text": RGBColor(0x7F, 0x8C, 0x8D),
    "bg_light": RGBColor(0xF5, 0xF6, 0xFA),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}


class PPTReporter:
    """PPT 報告生成器"""

    def __init__(self, config: dict):
        self.config = config
        self.today_str = date.today().isoformat()
        self.products = []
        self.stats = {}
        self.analysis_text = ""

    def set_data(self, products: list[dict], analysis_text: str):
        self.products = products
        self.analysis_text = analysis_text
        self.stats = self._compute_stats()

    def _compute_stats(self) -> dict:
        stats = {
            "total": len(self.products),
            "by_region": defaultdict(int),
            "by_brand": defaultdict(int),
            "by_platform": defaultdict(int),
            "by_wattage": defaultdict(int),
            "price_ranges": [],
            "products_with_price": 0,
            "total_price_usd": 0.0,
            "avg_price": 0.0,
        }
        for p in self.products:
            stats["by_region"][p.get("region", "Unknown")] += 1
            stats["by_brand"][p.get("brand", "Unknown")] += 1
            stats["by_platform"][p.get("platform", "Unknown")] += 1

            import re
            wattage = p.get("wattage") or ""
            if wattage:
                wm = re.search(r'(\d+)', wattage)
                if wm:
                    w = int(wm.group(1))
                    bucket = f"{w}W"
                    stats["by_wattage"][bucket] += 1

            try:
                pv = float(p.get("price") or 0)
                stats["price_ranges"].append(pv)
                stats["total_price_usd"] += pv
                stats["products_with_price"] += 1
            except (ValueError, TypeError):
                pass

        if stats["products_with_price"] > 0:
            stats["avg_price"] = stats["total_price_usd"] / stats["products_with_price"]

        return stats

    def generate(self, output_path: Optional[str] = None) -> Optional[str]:
        """生成 PPT 報告"""
        if not HAS_PPTX:
            logger.warning("python-pptx 未安裝，無法生成 PPT 報告")
            return None

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # 1. 封面
        self._add_cover(prs)
        # 2. 目錄
        self._add_toc(prs)
        # 3. 全球概覽
        self._add_global_overview(prs)
        # 4. 地區分析
        self._add_regional_analysis(prs)
        # 5. 品牌分析
        self._add_brand_analysis(prs)
        # 6. 平台對比
        self._add_platform_comparison(prs)
        # 7. 瓦數分析
        self._add_wattage_analysis(prs)
        # 8. AI 洞察
        self._add_ai_insights(prs)
        # 9. 總結
        self._add_summary(prs)

        if output_path is None:
            output_path = str(
                Path(self.config.get("output", {}).get("report_dir", "reports"))
                / f"全球PSU報告_{self.today_str}.pptx"
            )

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        logger.info(f"✅ PPT 報告已生成: {output_path}")
        return output_path

    # ─── Slide 1: 封面 ────────────────────────────────────────────

    def _add_cover(self, prs):
        slide = prs.slide_layouts[6]  # 空白
        slide = prs.slides.add_slide(slide)

        # 背景
        self._fill_slide_bg(slide, THEME["secondary"])

        # 主標題
        self._add_text_box(
            slide, "全球電商平台電源供應器（PSU）", 0.5, 1.8, 12.3, 1.2,
            font_size=44, bold=True, color=THEME["white"],
            align=PP_ALIGN.CENTER
        )
        self._add_text_box(
            slide, "每日市場分析報告", 0.5, 3.0, 12.3, 0.8,
            font_size=32, bold=False, color=THEME["accent"],
            align=PP_ALIGN.CENTER
        )

        # 日期和平台數
        self._add_text_box(
            slide,
            f"📅 {self.today_str}  |  🌐 覆蓋 35+ 全球電商平台  |  "
            f"📊 {self.stats.get('total', 0)} 個商品數據",
            0.5, 4.2, 12.3, 0.6,
            font_size=18, color=THEME["white"], align=PP_ALIGN.CENTER
        )

        # 生成時間
        self._add_text_box(
            slide,
            f"報告生成時間：{datetime.now().strftime('%Y-%m-%d %H:%M')}",
            0.5, 6.5, 12.3, 0.5,
            font_size=12, color=RGBColor(0xBB, 0xDE, 0xFB),
            align=PP_ALIGN.CENTER
        )

    # ─── Slide 2: 目錄 ───────────────────────────────────────────

    def _add_toc(self, prs):
        slide = prs.slide_layouts[6].slide if hasattr(prs.slide_layouts[6], 'slide') else prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self._fill_slide_bg(slide, THEME["bg_light"])

        self._add_text_box(
            slide, "📋 報告目錄", 0.5, 0.3, 12.3, 0.8,
            font_size=32, bold=True, color=THEME["secondary"],
            align=PP_ALIGN.LEFT
        )

        toc_items = [
            ("1", "🌍 全球數據概覽", "收集商品總數、覆蓋地區與平台"),
            ("2", "🗺️ 地區市場分析", "北美/歐洲/俄羅斯/亞洲/南美/中東"),
            ("3", "🏷️ 品牌競爭分析", "華碩/技嘉/微星/海盗船/海韻等"),
            ("4", "🛒 平台對比", "各平台價格水平與產品豐富度"),
            ("5", "⚡ 瓦數與認證分析", "主流瓦數段與 80+ 認證分佈"),
            ("6", "🤖 AI 市場洞察", "MiniMax AI 分析關鍵結論"),
            ("7", "📌 總結與建議", "市場洞察與策略建議"),
        ]

        for i, (num, title, desc) in enumerate(toc_items):
            y = 1.3 + i * 0.82
            self._add_text_box(
                slide, num, 0.5, y, 0.6, 0.6,
                font_size=22, bold=True, color=THEME["white"],
                align=PP_ALIGN.CENTER,
                bg_color=THEME["primary"]
            )
            self._add_text_box(
                slide, title, 1.2, y, 6, 0.5,
                font_size=18, bold=True, color=THEME["text"]
            )
            self._add_text_box(
                slide, desc, 1.2, y + 0.38, 10, 0.4,
                font_size=12, color=THEME["light_text"]
            )

    # ─── Slide 3: 全球概覽 ───────────────────────────────────────

    def _add_global_overview(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._fill_slide_bg(slide, THEME["white"])

        self._add_slide_header(slide, "🌍 全球數據概覽", self.today_str)

        # KPI 卡片
        kpis = [
            ("📦", "總商品數", str(self.stats.get("total", 0)), THEME["primary"]),
            ("🌐", "覆蓋地區", str(len(self.stats.get("by_region", {}))), THEME["accent"]),
            ("🛒", "覆蓋平台", str(len(self.stats.get("by_platform", {}))), THEME["warning"]),
            ("💰", "均價(USD)", f"${self.stats.get('avg_price', 0):.1f}", THEME["success"]),
        ]
        for i, (icon, label, value, color) in enumerate(kpis):
            x = 0.4 + i * 3.2
            self._add_kpi_card(slide, x, 1.1, 2.9, 1.5, icon, label, value, color)

        # 地區分布（文字）
        self._add_text_box(
            slide, "📍 各地區分佈", 0.5, 2.8, 12, 0.5,
            font_size=18, bold=True, color=THEME["secondary"]
        )

        region_data = sorted(
            self.stats.get("by_region", {}).items(),
            key=lambda x: x[1], reverse=True
        )
        if region_data:
            bar_total = self.stats.get("total", 1)
            y = 3.3
            for region, count in region_data[:8]:
                pct = count / bar_total * 100
                self._add_bar_row(slide, 0.5, y, pct, f"{region} ({count}個, {pct:.1f}%)", THEME["primary"])
                y += 0.45
        else:
            self._add_text_box(
                slide, "（今日未收集到數據）", 0.5, 3.3, 12, 0.5,
                font_size=14, color=THEME["light_text"]
            )

    # ─── Slide 4: 地區市場分析 ──────────────────────────────────

    def _add_regional_analysis(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._fill_slide_bg(slide, THEME["white"])
        self._add_slide_header(slide, "🗺️ 地區市場分析", self.today_str)

        regions = [
            ("🌎 北美", "Amazon US, Newegg, Best Buy", "北美消費者偏好高瓦數（750W+），"
             "80+ Gold/Platinum 為主流，性價比意識強"),
            ("🇪🇺 歐洲", "Amazon DE/UK, MediaMarkt, Saturn", "環保意識高，80+ Gold 以上為標配，"
             "德系品牌（be quiet!）強勢，價格偏高"),
            ("🇷🇺 俄羅斯", "Wildberries, Ozon, Yandex Market", "价格敏感度高，本地品牌，性價比優先，"
             "PSU 瓦數偏低（450-650W）"),
            ("🇨🇳 中國", "京東, 天貓, 淘寶", "全球最大 PSU 市场，品牌競爭激烈，"
             "價格區間最廣（¥100-¥2000+）"),
            ("🇹🇼 台灣", "PChome, Momo", "PChome 24h 為主战场，品牌齊全，"
             "價格與中國接近，用戶注重CP值"),
            ("🇯🇵 日本", "Amazon JP, Rakuten", "對品質要求高，全模組產品受歡迎，"
             "海韻 Seasonic 口碑佳，價格區間中高"),
            ("🇰🇷 韓國", "Gmarket, Coupang", "Gmarket 為主，需韓文關鍵字優化，"
             "本地化品牌偏好強，Corsair 強勢"),
            ("🇮🇳 印度", "Flipkart, Amazon.in", "價格為最敏感因素，基礎款暢銷，"
             "Flipkart 低價策略明顯，INR 3000-5000 為主流"),
        ]

        cols = 2
        for i, (region, platforms, insight) in enumerate(regions):
            col = i % cols
            row = i // cols
            x = 0.3 + col * 6.5
            y = 1.1 + row * 1.55
            self._add_region_card(slide, x, y, 6.2, 1.45, region, platforms, insight)

    # ─── Slide 5: 品牌分析 ──────────────────────────────────────

    def _add_brand_analysis(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._fill_slide_bg(slide, THEME["white"])
        self._add_slide_header(slide, "🏷️ 品牌競爭分析", self.today_str)

        brand_data = sorted(
            self.stats.get("by_brand", {}).items(),
            key=lambda x: x[1], reverse=True
        )

        if brand_data:
            bar_total = sum(c for _, c in brand_data)
            y = 1.3
            brand_colors = {
                "華碩（ASUS）": RGBColor(0x00, 0x5B, 0xB0),
                "技嘉（GIGABYTE）": RGBColor(0xE1, 0x26, 0x26),
                "微星（MSI）": RGBColor(0xCC, 0x00, 0x00),
                "海盜船（Corsair）": RGBColor(0xFF, 0x6D, 0x00),
                "海韻（Seasonic）": RGBColor(0x00, 0x96, 0x88),
                "安鈦克（Antec）": RGBColor(0x7F, 0x8C, 0x8D),
                "酷冷至尊（Cooler Master）": RGBColor(0x00, 0xA1, 0xE1),
                "BQT": RGBColor(0xFF, 0xC1, 0x07),
                "九州風神（DeepCool）": RGBColor(0x4A, 0x90, 0xE2),
                "聯力（Lian Li）": RGBColor(0x2C, 0x3E, 0x50),
            }

            for brand, count in brand_data[:10]:
                pct = count / bar_total * 100 if bar_total > 0 else 0
                color = brand_colors.get(brand, THEME["primary"])
                label = f"{brand} ({count}個, {pct:.1f}%)"
                self._add_bar_row(slide, 0.5, y, pct, label, color)
                y += 0.45
        else:
            self._add_text_box(
                slide, "（今日未收集到品牌數據）", 0.5, 1.3, 12, 0.5,
                font_size=14, color=THEME["light_text"]
            )

    # ─── Slide 6: 平台對比 ──────────────────────────────────────

    def _add_platform_comparison(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._fill_slide_bg(slide, THEME["white"])
        self._add_slide_header(slide, "🛒 平台對比分析", self.today_str)

        platform_data = sorted(
            self.stats.get("by_platform", {}).items(),
            key=lambda x: x[1], reverse=True
        )

        if platform_data:
            # 左側：柱狀圖文字版
            self._add_text_box(
                slide, "各平台商品數量", 0.5, 1.1, 6, 0.5,
                font_size=16, bold=True, color=THEME["secondary"]
            )
            y = 1.6
            for platform, count in platform_data[:12]:
                pct = count / self.stats.get("total", 1) * 100
                self._add_bar_row(slide, 0.5, y, pct, f"{platform} ({count})", THEME["accent"])
                y += 0.42

            # 右側：平台特點說明
            self._add_text_box(
                slide, "主要平台特點", 6.8, 1.1, 6, 0.5,
                font_size=16, bold=True, color=THEME["secondary"]
            )
            insights = [
                ("京東/天貓（中國）", "覆蓋最廣，價格區間最大，品牌最齊全"),
                ("Amazon（全球）", "各站點覆盖广，價格透明，评论数据丰富"),
                ("PChome/Momo（台灣）", "發貨快速，品牌齊全，中文界面友好"),
                ("Newegg（北美）", "專業 DIY 用戶首選，規格齊全"),
                ("Mercado Libre（南美）", "覆蓋多國，物流分散，價格含運費"),
                ("Noon（中東）", "新品平台為主，國際品牌多"),
            ]
            y = 1.6
            for platform, insight in insights:
                self._add_text_box(
                    slide, f"• {platform}", 6.8, y, 6, 0.35,
                    font_size=13, bold=True, color=THEME["text"]
                )
                self._add_text_box(
                    slide, f"  {insight}", 6.8, y + 0.32, 6, 0.35,
                    font_size=11, color=THEME["light_text"]
                )
                y += 0.78
        else:
            self._add_text_box(
                slide, "（今日未收集到數據）", 0.5, 1.3, 12, 0.5,
                font_size=14, color=THEME["light_text"]
            )

    # ─── Slide 7: 瓦數分析 ──────────────────────────────────────

    def _add_wattage_analysis(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._fill_slide_bg(slide, THEME["white"])
        self._add_slide_header(slide, "⚡ 瓦數與認證分析", self.today_str)

        wattage_data = sorted(
            self.stats.get("by_wattage", {}).items(),
            key=lambda x: x[1], reverse=True
        )

        self._add_text_box(
            slide, "熱門瓦數分佈", 0.5, 1.1, 6, 0.5,
            font_size=16, bold=True, color=THEME["secondary"]
        )

        if wattage_data:
            y = 1.6
            for w, count in wattage_data[:10]:
                pct = count / sum(c for _, c in wattage_data) * 100
                self._add_bar_row(slide, 0.5, y, pct, f"{w} ({count}個)", THEME["warning"])
                y += 0.42
        else:
            self._add_text_box(
                slide, "（今日未收集到瓦數數據）", 0.5, 1.6, 6, 0.5,
                font_size=14, color=THEME["light_text"]
            )

        # 右側：瓦數趨勢說明
        self._add_text_box(
            slide, "全球瓦數趨勢觀察", 6.8, 1.1, 6, 0.5,
            font_size=16, bold=True, color=THEME["secondary"]
        )
        wattage_insights = [
            ("850W-1000W", "全球主流高端，RTX 4090/4080 帶動需求"),
            ("750W", "北美/歐洲主流，性價比最高的旗艦級門票"),
            ("650W", "亞洲/新興市場主流，主流遊戲配置"),
            ("550W", "基礎入門，全球通用，價格競爭激烈"),
            ("1000W+", "歐美專業用戶/工作站，份額小但增長快"),
            ("80+ Gold", "全球最普及標準，性價比最優"),
            ("80+ Platinum+", "歐美高端，商用水冷/旗艦配置"),
        ]
        y = 1.6
        for spec, desc in wattage_insights:
            self._add_text_box(
                slide, f"▸ {spec}", 6.8, y, 6, 0.35,
                font_size=13, bold=True, color=THEME["text"]
            )
            self._add_text_box(
                slide, f"  {desc}", 6.8, y + 0.3, 6, 0.35,
                font_size=11, color=THEME["light_text"]
            )
            y += 0.72

    # ─── Slide 8: AI 洞察 ────────────────────────────────────────

    def _add_ai_insights(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._fill_slide_bg(slide, THEME["bg_light"])
        self._add_slide_header(slide, "🤖 AI 市場洞察", self.today_str)

        if self.analysis_text:
            # 解析 AI 分析文字，取關鍵段落
            lines = self.analysis_text.split("\n")
            insight_lines = []
            for line in lines:
                line = line.strip()
                if line and (line.startswith("##") or line.startswith("**") or "建議" in line or "洞察" in line or "觀察" in line):
                    insight_lines.append(line.replace("##", "").replace("**", ""))
                if len(insight_lines) >= 20:
                    break

            if insight_lines:
                y = 1.1
                for line in insight_lines[:15]:
                    if line.startswith("#"):
                        line = line.lstrip("#").strip()
                        self._add_text_box(
                            slide, line, 0.5, y, 12, 0.5,
                            font_size=16, bold=True, color=THEME["secondary"]
                        )
                    else:
                        self._add_text_box(
                            slide, line[:150], 0.5, y, 12, 0.4,
                            font_size=12, color=THEME["text"]
                        )
                    y += 0.42
            else:
                self._add_text_box(
                    slide, self.analysis_text[:2000], 0.5, 1.1, 12, 6,
                    font_size=11, color=THEME["text"]
                )
        else:
            self._add_text_box(
                slide,
                "⚠️ 今日分析文字為空，請檢查數據收集情況",
                0.5, 1.1, 12, 5, font_size=14, color=THEME["light_text"]
            )

    # ─── Slide 9: 總結 ───────────────────────────────────────────

    def _add_summary(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._fill_slide_bg(slide, THEME["secondary"])

        self._add_text_box(
            slide, "📌 總結與建議", 0.5, 0.4, 12.3, 0.8,
            font_size=32, bold=True, color=THEME["white"],
            align=PP_ALIGN.CENTER
        )

        summary_items = [
            ("全球覆蓋",
             f"今日收集涵蓋 {len(self.stats.get('by_region', {}))} 個地區、"
             f"{len(self.stats.get('by_platform', {}))} 個平台的 "
             f"{self.stats.get('total', 0)} 個 PSU 商品數據"),
            ("市場格局",
             "北美/歐洲市場以 850W 80+ Gold 為主流，亞洲市場瓦數多樣，"
             "俄羅斯/南美偏向基礎款，價格差異顯著"),
            ("品牌態勢",
             "華碩/微星/技嘉覆蓋最廣，海盜船在歐美強勢，海韻在日本口碑佳，"
             "DeepCool/Lian Li 在亞洲市場滲透率提升"),
            ("區域機會",
             "台灣 PChome 24h 為核心入口，建議強化日文/韓文關鍵字優化，"
             "印度市場需價格下探至 INR 3000 以内"),
            ("行動建議",
             "1) 持續監控京東/天貓價格波動  "
             "2) 針對台灣/日本市場推出中高瓦數產品  "
             "3) 佈局印度低價位段"),
        ]

        y = 1.5
        for title, content in summary_items:
            self._add_text_box(
                slide, f"▶ {title}", 0.5, y, 12, 0.5,
                font_size=16, bold=True, color=THEME["accent"]
            )
            self._add_text_box(
                slide, content, 0.7, y + 0.45, 11.8, 0.5,
                font_size=12, color=THEME["white"]
            )
            y += 1.1

        self._add_text_box(
            slide,
            f"報告生成時間：{datetime.now().strftime('%Y-%m-%d %H:%M')}  |  "
            f"Powered by Claude Code AI",
            0.5, 7.0, 12.3, 0.4,
            font_size=11, color=RGBColor(0xBB, 0xDE, 0xFB),
            align=PP_ALIGN.CENTER
        )

    # ─── 通用工具函數 ────────────────────────────────────────────

    def _fill_slide_bg(self, slide, color: RGBColor):
        from pptx.util import Inches
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.33), Inches(7.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    def _add_text_box(
        self, slide, text: str, x: float, y: float, w: float, h: float,
        font_size: int = 12, bold: bool = False, color=None,
        align=PP_ALIGN.LEFT, bg_color: RGBColor | None = None
    ):
        from pptx.util import Inches, Pt

        if bg_color is not None:
            # Add background rectangle behind text
            bg = slide.shapes.add_shape(
                1, Inches(x), Inches(y), Inches(w), Inches(h)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = bg_color
            bg.line.fill.background()

        if color is None:
            color = THEME["text"]
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        return txBox

    def _add_slide_header(self, slide, title: str, subtitle: str = ""):
        # 頂部藍色條
        header = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.33), Inches(0.9)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = THEME["primary"]
        header.line.fill.background()

        self._add_text_box(
            slide, title, 0.3, 0.1, 9, 0.7,
            font_size=26, bold=True, color=THEME["white"]
        )
        self._add_text_box(
            slide, subtitle, 9.5, 0.25, 3.5, 0.4,
            font_size=12, color=RGBColor(0xBB, 0xDE, 0xFB),
            align=PP_ALIGN.RIGHT
        )

    def _add_kpi_card(
        self, slide, x: float, y: float, w: float, h: float,
        icon: str, label: str, value: str, color: RGBColor
    ):
        card = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()

        self._add_text_box(
            slide, icon, x + 0.1, y + 0.1, w - 0.2, 0.4,
            font_size=24, color=THEME["white"], align=PP_ALIGN.CENTER
        )
        self._add_text_box(
            slide, value, x + 0.1, y + 0.45, w - 0.2, 0.6,
            font_size=22, bold=True, color=THEME["white"],
            align=PP_ALIGN.CENTER
        )
        self._add_text_box(
            slide, label, x + 0.1, y + 1.0, w - 0.2, 0.35,
            font_size=11, color=THEME["white"],
            align=PP_ALIGN.CENTER
        )

    def _add_region_card(
        self, slide, x: float, y: float, w: float, h: float,
        region: str, platforms: str, insight: str
    ):
        card = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = THEME["bg_light"]
        card.line.color.rgb = THEME["primary"]
        card.line.width = Pt(1)

        self._add_text_box(
            slide, region, x + 0.1, y + 0.08, w - 0.2, 0.35,
            font_size=13, bold=True, color=THEME["primary"]
        )
        self._add_text_box(
            slide, platforms, x + 0.1, y + 0.42, w - 0.2, 0.3,
            font_size=10, color=THEME["light_text"]
        )
        self._add_text_box(
            slide, insight, x + 0.1, y + 0.72, w - 0.2, 0.65,
            font_size=10, color=THEME["text"]
        )

    def _add_bar_row(
        self, slide, x: float, y: float, pct: float, label: str, color: RGBColor
    ):
        max_width = 11.0
        bar_width = max(pct / 100 * max_width, 0.1)
        bar_height = 0.35

        # 標籤
        self._add_text_box(
            slide, label, x, y, 4.5, 0.32,
            font_size=10, color=THEME["text"]
        )
        # 背景條
        bg = slide.shapes.add_shape(
            1, Inches(x + 4.5), Inches(y), Inches(max_width), Inches(bar_height)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xE8, 0xEA, 0xED)
        bg.line.fill.background()

        # 數值條
        bar = slide.shapes.add_shape(
            1, Inches(x + 4.5), Inches(y), Inches(bar_width), Inches(bar_height)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
