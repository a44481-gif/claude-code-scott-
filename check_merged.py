from pptx import Presentation
prs = Presentation('d:/claude mini max 2.7/storage_market_output/CCS_Global_MERGED_v6.pptx')
print(f"Total slides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    texts = []
    pic_count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t[:40])
        if shape.shape_type == 13:
            pic_count += 1
    print(f"Slide {i+1}: {texts[:2] if texts else '(blank)'} | pics={pic_count}")
