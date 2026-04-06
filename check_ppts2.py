from pptx import Presentation
prs1 = Presentation('d:/claude mini max 2.7/storage_market_output/Global_Storage_Market_2025_2030.pptx')
prs2 = Presentation('d:/claude mini max 2.7/storage_market_output/CCS_全球储能市场综合报告_2025-2030_FINAL.pptx')
print('src1 slides:', len(prs1.slides))
print('src2 slides:', len(prs2.slides))
for i, slide in enumerate(prs2.slides):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t[:30])
    print(f'src2 slide {i+1}: {texts[:3]}')
