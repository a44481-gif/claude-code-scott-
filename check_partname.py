from pptx import Presentation

src2 = r'd:/claude mini max 2.7/storage_market_output/CCS_Global_Storage_2025_2030.pptx'
prs2 = Presentation(src2)
print(f"Total slides: {len(prs2.slides)}")
for i in range(min(3, len(prs2.slides))):
    slide = prs2.slides[i]
    print(f"Slide {i} partname: {slide.part.partname}")
    for rid, rel in slide.part.rels.items():
        print(f"  {rid}: {rel.reltype[-30:]} -> {rel.target_ref}")
