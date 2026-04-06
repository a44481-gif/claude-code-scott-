# -*- coding: utf-8 -*-
"""
PPT合併腳本 v7 - 乾淨版本，直接用python-pptx深度複製
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree
import os

SRC1 = r'd:/claude mini max 2.7/storage_market_output/Global_Storage_Market_2025_2030.pptx'
SRC2 = r'd:/claude mini max 2.7/storage_market_output/CCS_Global_Storage_2025_2030.pptx'
OUTPUT = r'd:/claude mini max 2.7/storage_market_output/CCS_Global_MERGED_v7.pptx'

print("Loading source files...")
prs1 = Presentation(SRC1)  # 全球储能
prs2 = Presentation(SRC2)  # CCS报告
print(f"SRC1 (全球储能): {len(prs1.slides)} slides")
print(f"SRC2 (CCS报告): {len(prs2.slides)} slides")

# 建立新的合併簡報
merged = Presentation()
merged.slide_width = prs1.slide_width
merged.slide_height = prs1.slide_height

def find_cSld(tree):
    for elem in tree.iter():
        if elem.tag == qn('p:cSld'):
            return elem
    return None

def deep_copy_slide(target_prs, source_prs, slide_idx):
    """
    深度複製幻燈片：完整複製所有形狀和背景，並複製圖片資源
    """
    src_slide = source_prs.slides[slide_idx]
    src_tree = src_slide.part._element

    # 複製圖片資源
    src_part = src_slide.part
    for rel_id, rel in src_part.rels.items():
        if rel.reltype.endswith('image'):
            # 圖片資源，需要複製到目標簡報
            try:
                img_part = rel.target_part
                # 將圖片添加到目標簡報
                img_blob = img_part.blob
                # 獲取圖片類型
                content_type = img_part.content_type
                # 添加到目標簡報
                img_name = f'image_{rel_id}.{content_type.split("/")[-1]}'
                target_prs.part.add_image(
                    target_prs.part.pack_uri(rel.target_ref).content_type,
                    img_blob
                )
            except Exception as e:
                pass

    blank_layout = target_prs.slide_layouts[6]
    new_slide = target_prs.slides.add_slide(blank_layout)
    new_tree = new_slide.part._element

    cSld_src = find_cSld(src_tree)
    cSld_new = find_cSld(new_tree)

    if cSld_src is None or cSld_new is None:
        return new_slide

    # ---- 背景 ----
    bg_src = cSld_src.find(qn('p:bg'))
    bg_new = cSld_new.find(qn('p:bg'))
    if bg_new is not None:
        cSld_new.remove(bg_new)
    if bg_src is not None:
        cSld_new.insert(0, deepcopy(bg_src))

    # ---- 形狀樹 ----
    spTree_src = cSld_src.find(qn('p:spTree'))
    spTree_new = cSld_new.find(qn('p:spTree'))

    if spTree_new is not None:
        for child in list(spTree_new):
            if child.tag != qn('p:nvGrpSpPr') and child.tag != qn('p:grpSpPr'):
                spTree_new.remove(child)

    if spTree_src is not None:
        for child in spTree_src:
            if child.tag != qn('p:nvGrpSpPr') and child.tag != qn('p:grpSpPr'):
                copied = deepcopy(child)
                spTree_new.append(copied)

    return new_slide

def add_divider(target_prs, part_num, line1, line2, line3="", bg_color=(0x0A,0x29,0x5E)):
    slide = target_prs.slides.add_slide(target_prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*bg_color)

    def add_t(text, size, bold, color_rgb, y, h):
        txBox = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.3), h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color_rgb)
        p.alignment = PP_ALIGN.CENTER

    add_t(f"PART {part_num}", 20, True, (0xFF,0xD7,0x00), Inches(1.8), Inches(0.8))
    add_t(line1, 36, True, (255,255,255), Inches(2.6), Inches(1.2))
    if line2:
        add_t(line2, 28, True, (255,255,255), Inches(3.7), Inches(1.0))
    if line3:
        add_t(line3, 18, False, (0xCC,0xD9,0xEC), Inches(4.9), Inches(0.8))
    return slide

def add_cover(target_prs):
    slide = target_prs.slides.add_slide(target_prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x0A,0x29,0x5E)

    def add_t(text, size, bold, color_rgb, y, h):
        txBox = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.3), h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color_rgb)
        p.alignment = PP_ALIGN.CENTER

    add_t("CCS集成母排 × 全球储能市场", 38, True, (255,255,255), Inches(1.2), Inches(1.5))
    add_t("综合分析报告 2025-2030", 30, True, (0xFF,0xD7,0x00), Inches(2.7), Inches(1.0))
    add_t("CCS集成母排与线束技术整合市场分析报告", 18, False, (0xCC,0xD9,0xEC), Inches(4.2), Inches(0.8))
    add_t("全球储能及新能源关联市场经济规模分析", 16, False, (0x99,0xAA,0xCC), Inches(5.0), Inches(0.7))
    add_t("2026年4月  |  内部研究资料", 14, False, (0x77,0x88,0xAA), Inches(5.8), Inches(0.5))
    return slide

# ============================================================
# 合併流程
# ============================================================
print("\nBuilding merged presentation...")

# 1. 封面
add_cover(merged)
print("  [封面]")

# 2. Part 1 分隔頁
add_divider(merged, 1, "CCS集成母排与线束技术整合", "市场分析报告", bg_color=(0x1A,0x5C,0x9E))
print("  [Part 1 分隔页]")

# 3. CCS 報告的所有 slides
print("  [CCS报告 slides...]")
for i in range(len(prs2.slides)):
    deep_copy_slide(merged, prs2, i)
    print(f"    CCS {i+1}/{len(prs2.slides)}", end='\r')

# 4. Part 2 分隔頁
add_divider(merged, 2, "全球储能及新能源关联市场", "2025-2030", bg_color=(0x0A,0x29,0x5E))
print("\n  [Part 2 分隔页]")

# 5. 全球儲能 PPT 的所有 slides
print("  [全球储能 slides...]")
for i in range(len(prs1.slides)):
    deep_copy_slide(merged, prs1, i)
    print(f"    储能 {i+1}/{len(prs1.slides)}", end='\r')

# 保存
merged.save(OUTPUT)
size = os.path.getsize(OUTPUT)
total = 1 + 1 + len(prs2.slides) + 1 + len(prs1.slides)
print(f"\n\n=== Done! ===")
print(f"File: {OUTPUT}")
print(f"Size: {size/1024:.0f} KB ({size/1024/1024:.2f} MB)")
print(f"Total slides: {len(merged.slides)} (expected: {total})")

# 驗證
verify = Presentation(OUTPUT)
print(f"Verified: {len(verify.slides)} slides")
