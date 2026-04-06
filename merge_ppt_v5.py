# -*- coding: utf-8 -*-
"""
PPT合併腳本 v5 - 最可靠的 _element 深層複製方式
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree
import os, shutil, re

SRC1 = r'd:/claude mini max 2.7/storage_market_output/Global_Storage_Market_2025_2030.pptx'
SRC2 = r'd:/claude mini max 2.7/storage_market_output/CCS_Global_Storage_2025_2030.pptx'
OUTPUT = r'd:/claude mini max 2.7/storage_market_output/CCS_Global_MERGED_v5.pptx'

prs1 = Presentation(SRC1)
prs2 = Presentation(SRC2)
print(f"SRC1 (全球储能): {len(prs1.slides)} slides")
print(f"SRC2 (CCS报告): {len(prs2.slides)} slides")

# 建立新的空白簡報
merged = Presentation()
merged.slide_width = prs1.slide_width
merged.slide_height = prs1.slide_height

# 複製 prs2 的所有 slides
def add_slide_copy(target_prs, source_prs, slide_idx):
    """使用 _element 深層複製複製幻燈片"""
    src_slide = source_prs.slides[slide_idx]
    src_tree = src_slide.part._element

    # 新增一張空白幻燈片
    blank_layout = target_prs.slide_layouts[6]
    new_slide = target_prs.slides.add_slide(blank_layout)
    new_tree = new_slide.part._element

    # 找到 cSld 元素
    def find_cSld(tree):
        for elem in tree.iter():
            if elem.tag == qn('p:cSld'):
                return elem
        return None

    cSld_src = find_cSld(src_tree)
    cSld_new = find_cSld(new_tree)

    if cSld_src is not None and cSld_new is not None:
        # 處理背景
        bg_src = cSld_src.find(qn('p:bg'))
        bg_new = cSld_new.find(qn('p:bg'))
        if bg_new is not None:
            cSld_new.remove(bg_new)
        if bg_src is not None:
            cSld_new.insert(0, deepcopy(bg_src))

        # 處理 spTree
        spTree_src = cSld_src.find(qn('p:spTree'))
        spTree_new = cSld_new.find(qn('p:spTree'))

        if spTree_new is not None:
            # 移除 nvGrpSpPr 和 grpSpPr 以外的所有元素
            for child in list(spTree_new):
                if child.tag != qn('p:nvGrpSpPr') and child.tag != qn('p:grpSpPr'):
                    spTree_new.remove(child)

        if spTree_src is not None:
            # 深層複製所有形狀（跳過 nvGrpSpPr 和 grpSpPr）
            for child in spTree_src:
                if child.tag != qn('p:nvGrpSpPr') and child.tag != qn('p:grpSpPr'):
                    spTree_new.append(deepcopy(child))

    return new_slide

# 加入綜合封面
print("加入综合封面...")
cover = merged.slides.add_slide(merged.slide_layouts[6])
cover.background.fill.solid()
cover.background.fill.fore_color.rgb = RGBColor(0x0A, 0x29, 0x5E)

def add_text(slide, text, size, bold, color, x, y, w, h, align=PP_ALIGN.CENTER):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align

add_text(cover, "CCS集成母排 × 全球储能市场", 38, True, RGBColor(255,255,255),
         Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.5))
add_text(cover, "综合分析报告 2025-2030", 30, True, RGBColor(0xFF, 0xD7, 0x00),
         Inches(0.5), Inches(2.7), Inches(12.3), Inches(1.0))
add_text(cover, "CCS集成母排与线束技术整合市场分析报告", 18, False, RGBColor(0xCC,0xD9,0xEC),
         Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8))
add_text(cover, "全球储能及新能源关联市场经济规模分析", 16, False, RGBColor(0x99,0xAA,0xCC),
         Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.7))
add_text(cover, "2026年4月  |  内部研究资料", 14, False, RGBColor(0x77,0x88,0xAA),
         Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5))

# Part 1 分隔頁
print("加入 Part 1 分隔页...")
div1 = merged.slides.add_slide(merged.slide_layouts[6])
div1.background.fill.solid()
div1.background.fill.fore_color.rgb = RGBColor(0x1A, 0x5C, 0x9E)
add_text(div1, "PART 1", 20, True, RGBColor(0xFF, 0xD7, 0x00),
         Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.8))
add_text(div1, "CCS集成母排与线束技术整合", 36, True, RGBColor(255,255,255),
         Inches(0.5), Inches(2.6), Inches(12.3), Inches(1.2))
add_text(div1, "市场分析报告", 28, True, RGBColor(255,255,255),
         Inches(0.5), Inches(3.7), Inches(12.3), Inches(1.0))

# 複製 CCS PPT 的所有 slides
print("复制 CCS 报告...")
for i in range(len(prs2.slides)):
    add_slide_copy(merged, prs2, i)
    print(f"  CCS {i+1}/{len(prs2.slides)}", end='\r')

# Part 2 分隔頁
print("\n加入 Part 2 分隔页...")
div2 = merged.slides.add_slide(merged.slide_layouts[6])
div2.background.fill.solid()
div2.background.fill.fore_color.rgb = RGBColor(0x0A, 0x29, 0x5E)
add_text(div2, "PART 2", 20, True, RGBColor(0xFF, 0xD7, 0x00),
         Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.8))
add_text(div2, "全球储能及新能源关联市场", 36, True, RGBColor(255,255,255),
         Inches(0.5), Inches(2.6), Inches(12.3), Inches(1.2))
add_text(div2, "2025-2030", 28, True, RGBColor(0xFF, 0xD7, 0x00),
         Inches(0.5), Inches(3.7), Inches(12.3), Inches(1.0))

# 複製全球儲能 PPT 的所有 slides
print("复制全球储能分析...")
for i in range(len(prs1.slides)):
    add_slide_copy(merged, prs1, i)
    print(f"  储能 {i+1}/{len(prs1.slides)}", end='\r')

# 保存
merged.save(OUTPUT)
size = os.path.getsize(OUTPUT)
print(f"\n\n保存完成!")
print(f"文件: {OUTPUT}")
print(f"大小: {size/1024:.0f} KB ({size/1024/1024:.2f} MB)")
print(f"总幻灯片: {len(merged.slides)} 张")

# 驗證
verify = Presentation(OUTPUT)
print(f"驗證: {len(verify.slides)} slides OK")
print(f"  - 封面 1张 + CCS报告 {len(prs2.slides)}张 + 分隔页2张 + 储能 {len(prs1.slides)}张 = {1+len(prs2.slides)+2+len(prs1.slides)}张")
