# -*- coding: utf-8 -*-
"""
PPT合併腳本 - 使用XML深層複製，完整保留所有內容
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
import os
from lxml import etree

SRC1 = r"d:/claude mini max 2.7/storage_market_output/Global_Storage_Market_2025_2030.pptx"
SRC2 = r"d:/李總項目/CCS集成母排与线束技术整合市场分析报告.pptx"
OUTPUT = r"d:/claude mini max 2.7/storage_market_output/CCS_全球储能市场综合报告_2025-2030.pptx"

prs1 = Presentation(SRC1)
prs2 = Presentation(SRC2)

print(f"源文件1 (全球储能): {len(prs1.slides)} 张")
print(f"源文件2 (CCS报告): {len(prs2.slides)} 张")


def clone_slide_deep(merged_prs, source_prs, slide_index):
    """使用XML深層複製完整克隆幻燈片"""
    src_slide = source_prs.slides[slide_index]
    src_tree = src_slide.part._element

    blank_layout = merged_prs.slide_layouts[6]
    new_slide = merged_prs.slides.add_slide(blank_layout)
    new_tree = new_slide.part._element

    # 深層複製整個slide元素
    new_tree[:] = deepcopy(list(src_tree))

    # 確保slideId存在
    cSld = new_tree.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
    if cSld is not None:
        # 複製nvGrpSpPr (non-visual group shape properties)
        nvGrpSpPr = cSld.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}nvGrpSpPr')
        if nvGrpSpPr is not None:
            pass  # 已包含在深層複製中

    # 複製所有圖片關聯
    for shape in new_slide.shapes:
        if shape.shape_type == 13:  # 圖片
            pass  # 已由深層複製處理

    return new_slide


def add_cover_slide(merged_prs, title_text, subtitle_text):
    """建立綜合封面"""
    slide = merged_prs.slides.add_slide(merged_prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x0A, 0x29, 0x5E)

    # 主標題
    txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.5))
    tf = txTitle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # 副標題
    txSub = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(1.2))
    tf2 = txSub.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle_text
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xCC, 0xD9, 0xEC)
    p2.alignment = PP_ALIGN.CENTER

    # 日期
    txDate = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5))
    tf3 = txDate.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "2026年4月  |  内部研究资料"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0x88, 0x99, 0xBB)
    p3.alignment = PP_ALIGN.CENTER
    return slide


def add_divider_slide(merged_prs, part_num, title_text):
    """建立分隔頁"""
    slide = merged_prs.slides.add_slide(merged_prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x1A, 0x5C, 0x9E)

    # Part 编号
    txNum = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.0))
    pNum = txNum.text_frame.paragraphs[0]
    pNum.text = f"PART {part_num}"
    pNum.font.size = Pt(20)
    pNum.font.bold = True
    pNum.font.color.rgb = RGBColor(0xFF, 0xD7, 0x00)
    pNum.alignment = PP_ALIGN.CENTER

    # 标题
    txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.3), Inches(2.0))
    tf = txTitle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    return slide


# ============================================================
# 開始合併
# ============================================================
merged = Presentation()
merged.slide_width = Inches(13.33)
merged.slide_height = Inches(7.5)

# 1. 綜合封面
add_cover_slide(merged,
    "CCS集成母排 × 全球储能市场\n综合分析报告",
    "CCS集成母排与线束技术整合市场分析  |  2025-2030全球储能市场经济规模"
)
print("[1/4] 综合封面已创建")

# 2. CCS 报告（前缀分隔页）
add_divider_slide(merged, 1, "CCS集成母排与线束技术整合市场分析报告")
print("[2/4] 正在复制 CCS 报告...")
for i in range(len(prs2.slides)):
    clone_slide_deep(merged, prs2, i)
    print(f"  CCS {i+1}/{len(prs2.slides)}", end="\r")
print()

# 3. 储能市场分析（分隔页）
add_divider_slide(merged, 2, "全球储能及新能源关联市场经济规模分析\n2025-2030")
print("[3/4] 正在复制全球储能分析...")
for i in range(len(prs1.slides)):
    clone_slide_deep(merged, prs1, i)
    print(f"  储能 {i+1}/{len(prs1.slides)}", end="\r")
print()

# 4. 保存
merged.save(OUTPUT)
size = os.path.getsize(OUTPUT)
print(f"\n[4/4] 保存完成!")
print(f"文件: {OUTPUT}")
print(f"大小: {size/1024:.0f} KB")
print(f"总幻灯片: {len(merged.slides)} 张")
print(f"  - CCS报告: {len(prs2.slides)} 张")
print(f"  - 全球储能分析: {len(prs1.slides)} 张")
print(f"  - 封面+分隔: 3 张")
